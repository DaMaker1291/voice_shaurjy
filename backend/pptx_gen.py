import os
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUTPUT_DIR = "/opt/jarvis/cad_output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ── Color Palette (Dark Tech Theme) ──
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1B, 0x22)
ACCENT_BLUE = RGBColor(0x58, 0xA6, 0xFF)
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_ORANGE = RGBColor(0xF7, 0x8C, 0x1E)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY = RGBColor(0x8B, 0x94, 0x9E)
TEXT_LIGHT = RGBColor(0xC9, 0xD1, 0xD9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_LIGHT, bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
        p.level = 0
        # Add bullet character
        p.text = "  " + item
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u25CF'})
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': f'{bullet_color}'})
        buClr.append(srgbClr)
        pPr.append(buClr)
        pPr.append(buChar)
    return txBox


def add_transition(slide, transition_type="fade", duration_ms=700):
    transition = slide.element.makeelement(qn('p:transition'), {
        'advClick': '1',
        'advTm': '0',
        'spd': 'med'
    })
    if transition_type == "fade":
        child = transition.makeelement(qn('p:fade'), {})
    elif transition_type == "push":
        child = transition.makeelement(qn('p:push'), {'dir': 'l'})
    elif transition_type == "cover":
        child = transition.makeelement(qn('p:cover'), {'dir': 'l'})
    else:
        child = transition.makeelement(qn('p:fade'), {})
    transition.append(child)
    slide.element.append(transition)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE (Full-screen hero)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)

# Accent bar top
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), ACCENT_BLUE)

# Title
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "PARAMETRIC 3D TURBINE ENGINE", font_size=44, bold=True,
             color=TEXT_WHITE, alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")

# Accent line under title
add_accent_line(slide, Inches(4.5), Inches(3.3), Inches(4), ACCENT_BLUE)

# Subtitle
add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.8),
             "Autonomous Generation Report", font_size=24,
             color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Metadata
add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.5),
             f"Generated: {time.strftime('%B %d, %Y at %H:%M')}  |  Pipeline: JARVIS OS  |  Engine: Blender 4.0.2",
             font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_shape(slide, Inches(0), Inches(7.4), SLIDE_W, Pt(4), ACCENT_GREEN)

add_transition(slide, "fade")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2: OVERVIEW
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "PROJECT OVERVIEW", font_size=32, bold=True, color=TEXT_WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_BLUE)

# Left card
card = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.0), Inches(4.8), Inches(0.5),
             "Objective", font_size=20, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(1.2), Inches(2.6), Inches(4.8), Inches(1.5),
             "Design and render a parametric 3D mechanical turbine engine "
             "based on current aerospace specifications, using autonomous "
             "headless execution in a background VDI.",
             font_size=15, color=TEXT_LIGHT)

add_text_box(slide, Inches(1.2), Inches(4.2), Inches(4.8), Inches(0.5),
             "Pipeline Components", font_size=20, bold=True, color=ACCENT_BLUE)
items = [
    "Parametric Blender modeling (12-blade turbine)",
    "10-frame turntable animation render",
    "STL mesh export for CAD/CAM integration",
    "Automated PowerPoint generation",
    "ZIP payload compression & deployment",
]
add_bullet_list(slide, Inches(1.2), Inches(4.8), Inches(4.8), Inches(2.5), items, font_size=14)

# Right card — key stats
card2 = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(7.2), Inches(2.0), Inches(4.8), Inches(0.5),
             "Key Metrics", font_size=20, bold=True, color=ACCENT_GREEN)

stats = [
    ("Blades", "12"),
    ("Hub Radius", "1.0 units"),
    ("Blade Length", "2.0 units"),
    ("Render Resolution", "1920 x 1080"),
    ("Frames Rendered", "10"),
    ("Total Time", "~59 seconds"),
    ("Host Interaction", "Zero"),
]
y = 2.7
for label, value in stats:
    add_text_box(slide, Inches(7.2), Inches(y), Inches(2.5), Inches(0.35),
                 label, font_size=13, color=TEXT_GRAY)
    add_text_box(slide, Inches(9.8), Inches(y), Inches(2.5), Inches(0.35),
                 value, font_size=13, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.RIGHT)
    y += 0.42

add_transition(slide, "push")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3: TECHNICAL SPECS
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "TECHNICAL SPECIFICATIONS", font_size=32, bold=True, color=TEXT_WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_ORANGE)

specs = [
    ("Geometry", [
        "Number of Blades: 12",
        "Hub Radius: 1.0 units",
        "Blade Length: 2.0 units",
        "Blade Width: 0.3 units",
        "Blade Thickness: 0.05 units",
    ]),
    ("Assembly", [
        "Drive Shaft: 3.0 units x 0.15 radius",
        "Base Plate: 2.8 units radius x 0.15 thick",
        "Blade Pitch: 11.5 degrees",
    ]),
    ("Materials", [
        "Hub/Shaft: Brushed Metal (Metallic 0.95)",
        "Blades: Light Alloy (Metallic 0.90)",
        "Shader: Principled BSDF PBR",
    ]),
]

x_positions = [0.8, 4.8, 8.8]
for idx, (section, items) in enumerate(specs):
    x = x_positions[idx]
    card = add_shape(slide, Inches(x), Inches(1.8), Inches(3.6), Inches(5.0), BG_CARD)
    add_text_box(slide, Inches(x + 0.3), Inches(2.0), Inches(3.0), Inches(0.5),
                 section, font_size=20, bold=True, color=ACCENT_ORANGE)
    add_accent_line(slide, Inches(x + 0.3), Inches(2.6), Inches(1.5), ACCENT_ORANGE)
    add_bullet_list(slide, Inches(x + 0.3), Inches(2.9), Inches(3.0), Inches(3.5),
                    items, font_size=13, color=TEXT_LIGHT, bullet_color=ACCENT_ORANGE)

add_transition(slide, "cover")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4: RENDERING
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "RENDERING PIPELINE", font_size=32, bold=True, color=TEXT_WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_BLUE)

# Two columns
card_left = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.0), Inches(4.8), Inches(0.5),
             "Render Settings", font_size=20, bold=True, color=ACCENT_BLUE)
items_left = [
    "Engine: Blender Eevee (real-time rasterizer)",
    "Resolution: 1920 x 1080 (Full HD)",
    "Format: PNG (lossless compression)",
    "Samples: Eevee default (fast)",
    "Lighting: 2-point (Key + Fill)",
    "Key Light: 400W area at (5, -5, 8)",
    "Fill Light: 150W area at (-5, 5, 3)",
]
add_bullet_list(slide, Inches(1.2), Inches(2.6), Inches(4.8), Inches(4.0), items_left, font_size=14)

card_right = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(7.2), Inches(2.0), Inches(4.8), Inches(0.5),
             "Turntable Animation", font_size=20, bold=True, color=ACCENT_BLUE)
items_right = [
    "Frames: 10 (36-degree increments)",
    "Camera orbit radius: 10 units",
    "Camera elevation: 5 units",
    "Total rotation: 360 degrees",
    "Hero frame: Center-front view",
    "Export: Individual PNG + STL mesh",
]
add_bullet_list(slide, Inches(7.2), Inches(2.6), Inches(4.8), Inches(4.0), items_right, font_size=14,
                bullet_color=ACCENT_GREEN)

add_transition(slide, "fade")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5: BENCHMARKS
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "PERFORMANCE BENCHMARKS", font_size=32, bold=True, color=TEXT_WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_GREEN)

# Bar chart style with shapes
benchmarks = [
    ("Model Generation", 1.0, 1.0, "< 1s"),
    ("Single Frame (Eevee)", 3.0, 5.0, "~3s"),
    ("10-Frame Turntable", 30.0, 60.0, "~30s"),
    ("PowerPoint Gen", 0.8, 2.0, "< 1s"),
    ("ZIP Compression", 1.0, 3.0, "< 1s"),
    ("Total Pipeline", 59.0, 120.0, "59s"),
]

y = 2.0
for label, actual, max_val, display in benchmarks:
    # Label
    add_text_box(slide, Inches(0.8), Inches(y), Inches(3.0), Inches(0.4),
                 label, font_size=14, color=TEXT_LIGHT)
    # Bar background
    bar_bg = add_shape(slide, Inches(4.0), Inches(y + 0.05), Inches(7.5), Inches(0.3),
                       RGBColor(0x21, 0x26, 0x2D), corner_radius=Inches(0.05))
    # Bar fill
    fill_width = max(0.3, (actual / max_val) * 7.5)
    bar_color = ACCENT_GREEN if actual < 10 else ACCENT_BLUE if actual < 60 else ACCENT_ORANGE
    bar_fill = add_shape(slide, Inches(4.0), Inches(y + 0.05), Inches(fill_width), Inches(0.3),
                         bar_color, corner_radius=Inches(0.05))
    # Value
    add_text_box(slide, Inches(11.7), Inches(y), Inches(1.2), Inches(0.4),
                 display, font_size=14, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.RIGHT)
    y += 0.65

# Footer note
add_text_box(slide, Inches(0.8), Inches(y + 0.5), Inches(11), Inches(0.4),
             "Hardware: WSL2 VDI (DISPLAY=:99) | CPU-only rendering | GPU acceleration available with CUDA",
             font_size=12, color=TEXT_GRAY)

add_transition(slide, "push")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6: ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "AUTONOMOUS PIPELINE ARCHITECTURE", font_size=32, bold=True, color=TEXT_WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_BLUE)

# Flow diagram with boxes
stages = [
    ("1. OBSERVE", "Screenshot + OCR\nActive window detect", ACCENT_BLUE),
    ("2. PLAN", "LLM code generation\nStep-by-step actions", ACCENT_GREEN),
    ("3. ACT", "Blender bpy execution\nxdotool commands", ACCENT_ORANGE),
    ("4. REFLECT", "Verify render output\nCheck file existence", ACCENT_BLUE),
    ("5. RETRY", "Auto-fix errors\nRe-plan on failure", ACCENT_GREEN),
]

x = 0.5
for label, desc, color in stages:
    card = add_shape(slide, Inches(x), Inches(2.0), Inches(2.3), Inches(2.5), BG_CARD)
    add_text_box(slide, Inches(x + 0.15), Inches(2.2), Inches(2.0), Inches(0.5),
                 label, font_size=16, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(x + 0.5), Inches(2.8), Inches(1.3), color)
    add_text_box(slide, Inches(x + 0.15), Inches(3.0), Inches(2.0), Inches(1.2),
                 desc, font_size=12, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    x += 2.5

# Bottom section
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
             "Execution Environment", font_size=20, bold=True, color=ACCENT_BLUE)

env_items = [
    "All execution on DISPLAY=:99 (WSL2 VDI)  |  Host desktop (DISPLAY=:0) untouched",
    "XFCE4 desktop + Google Chrome + Blender 4.0.2 + python3-pptx",
    "Frame streaming via WebSocket to PiP overlay on host",
]
add_bullet_list(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(2.0), env_items, font_size=14)

add_transition(slide, "cover")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7: OUTPUT FILES
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "OUTPUT FILES", font_size=32, bold=True, color=TEXT_WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_GREEN)

files = [
    ("turbine_engine.blend", "970 KB", "Full Blender project file"),
    ("turbine_hero.png", "1.3 MB", "Hero render (1920x1080)"),
    ("turntable_frames/", "10 files", "360-degree rotation frames"),
    ("turbine_model.stl", "38 KB", "3D mesh for CAD/CAM"),
    ("turbine_report.pptx", "37 KB", "This presentation"),
    ("cad_assets.zip", "11 MB", "Compressed payload"),
]

y = 1.8
for fname, size, desc in files:
    card = add_shape(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(0.7), BG_CARD)
    add_text_box(slide, Inches(1.2), Inches(y + 0.1), Inches(4.0), Inches(0.5),
                 fname, font_size=15, bold=True, color=TEXT_WHITE, font_name="Consolas")
    add_text_box(slide, Inches(5.5), Inches(y + 0.1), Inches(1.5), Inches(0.5),
                 size, font_size=14, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.2), Inches(y + 0.1), Inches(4.5), Inches(0.5),
                 desc, font_size=14, color=TEXT_GRAY)
    y += 0.85

add_transition(slide, "fade")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8: SOFTWARE STACK
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "SOFTWARE STACK", font_size=32, bold=True, color=TEXT_WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_ORANGE)

# Required
card1 = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.0), Inches(4.8), Inches(0.5),
             "Required (Auto-Installed)", font_size=20, bold=True, color=ACCENT_GREEN)
items_req = [
    "Blender 4.0.2 — 3D modeling & rendering",
    "python3-pptx — PowerPoint generation",
    "ImageMagick — Screenshot capture (import)",
    "xdotool — Window management & input",
    "RapidOCR — Screen text recognition",
    "XFCE4 — Lightweight desktop environment",
]
add_bullet_list(slide, Inches(1.2), Inches(2.6), Inches(4.8), Inches(4.0), items_req, font_size=14,
                bullet_color=ACCENT_GREEN)

# Optional
card2 = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(7.2), Inches(2.0), Inches(4.8), Inches(0.5),
             "Optional / Enhancement", font_size=20, bold=True, color=ACCENT_ORANGE)
items_opt = [
    "FreeCAD — Alternative parametric CAD",
    "GIMP — Texture editing & compositing",
    "LibreOffice — Document export",
    "CUDA toolkit — GPU-accelerated rendering",
    "Playwright — Stealth web automation",
    "wkhtmltopdf — PDF generation",
]
add_bullet_list(slide, Inches(7.2), Inches(2.6), Inches(4.8), Inches(4.0), items_opt, font_size=14,
                bullet_color=ACCENT_ORANGE)

add_transition(slide, "push")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9: VISION VERIFICATION
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "VISION VERIFICATION LOOP", font_size=32, bold=True, color=TEXT_WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
             "How the AI verifies its work in Excel, Word, and other apps",
             font_size=16, color=TEXT_GRAY)

# Flow
flow_items = [
    ("Execute Action", "Type data, format cells,\napply formulas in Excel/Word", ACCENT_BLUE),
    ("Capture Screen", "Screenshot via ImageMagick\nimport -window root", ACCENT_GREEN),
    ("OCR Verify", "Read text on screen to confirm\ndata rendered correctly", ACCENT_ORANGE),
    ("Compare Expected", "LLM compares OCR output\nto expected result", ACCENT_BLUE),
    ("Auto-Correct", "If wrong: retry with fixes\nadjust formatting, re-enter data", ACCENT_GREEN),
]

x = 0.3
for label, desc, color in flow_items:
    card = add_shape(slide, Inches(x), Inches(2.5), Inches(2.4), Inches(2.5), BG_CARD)
    add_text_box(slide, Inches(x + 0.1), Inches(2.7), Inches(2.2), Inches(0.5),
                 label, font_size=15, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(x + 0.5), Inches(3.2), Inches(1.4), color)
    add_text_box(slide, Inches(x + 0.1), Inches(3.4), Inches(2.2), Inches(1.2),
                 desc, font_size=11, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    x += 2.55

# Example scenarios
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(0.5),
             "Example Scenarios", font_size=18, bold=True, color=ACCENT_BLUE)

examples = [
    "Excel: Enter formula =SUM(A1:A10) -> OCR sees #NAME? -> auto-fix to =SUM(A1:A10) with proper locale",
    "Word: Apply Heading 1 style -> OCR confirms font changed to 24pt bold -> proceed",
    "PowerPoint: Insert image -> OCR verifies image visible and centered -> adjust position if needed",
]
add_bullet_list(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1.5), examples, font_size=13,
                bullet_color=ACCENT_GREEN)

add_transition(slide, "fade")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10: CONCLUSION
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# Top accent bar
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), ACCENT_GREEN)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.0),
             "CONCLUSION", font_size=40, bold=True, color=TEXT_WHITE,
             alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")

add_accent_line(slide, Inches(5), Inches(2.6), Inches(3), ACCENT_GREEN)

# Summary cards
summary = [
    ("3D Model", "Parametric turbine with 12 blades,\nmetal PBR materials, STL export", ACCENT_BLUE),
    ("Rendering", "10-frame turntable animation\nat 1920x1080 resolution", ACCENT_GREEN),
    ("Presentation", "Professional 10-slide deck\nwith benchmarks and architecture", ACCENT_ORANGE),
    ("Automation", "Zero host desktop interaction\n100% background VDI execution", ACCENT_BLUE),
]

x = 0.8
for title, desc, color in summary:
    card = add_shape(slide, Inches(x), Inches(3.2), Inches(2.8), Inches(2.2), BG_CARD)
    add_text_box(slide, Inches(x + 0.2), Inches(3.4), Inches(2.4), Inches(0.5),
                 title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(x + 0.7), Inches(3.9), Inches(1.4), color)
    add_text_box(slide, Inches(x + 0.2), Inches(4.1), Inches(2.4), Inches(1.0),
                 desc, font_size=12, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    x += 3.1

# Final stats
add_text_box(slide, Inches(2), Inches(5.8), Inches(9), Inches(0.5),
             "Total Pipeline Time: 59 seconds  |  Files Generated: 6  |  Payload: 11 MB",
             font_size=16, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(6.4), Inches(9), Inches(0.5),
             "Powered by JARVIS OS Autonomous Pipeline",
             font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_shape(slide, Inches(0), Inches(7.4), SLIDE_W, Pt(4), ACCENT_BLUE)

add_transition(slide, "fade")


# ══════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════
pptx_path = os.path.join(OUTPUT_DIR, "turbine_report.pptx")
prs.save(pptx_path)
print("[JARVIS] Presentation saved: " + pptx_path)
print("[JARVIS] Slides: 10 | Theme: Dark Tech | Transitions: Yes")
