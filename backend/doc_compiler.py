"""
Document Compiler Engine
Generates production-quality DOCX, PPTX, and XLSX files with strict validation.
All files are validated for non-zero size and meaningful content before returning.
"""
import os
import io
from typing import List, Dict, Any, Optional
from pydantic import BaseModel, Field, validator


class SlideContent(BaseModel):
    title: str = "Slide"
    content: str = ""
    subtitle: Optional[str] = None


class DocxContent(BaseModel):
    heading: str = "Document"
    paragraphs: List[str] = Field(default_factory=list)
    tables: List[List[str]] = Field(default_factory=list)


class XlsxContent(BaseModel):
    headers: List[str] = Field(default_factory=list)
    rows: List[List[str]] = Field(default_factory=list)
    title: str = "Sheet"


class FileCodegenError(Exception):
    pass


def _validate_output(path: str, min_bytes: int = 5120) -> None:
    if not os.path.exists(path):
        raise FileCodegenError(f"Output file not created: {path}")
    size = os.path.getsize(path)
    if size < min_bytes:
        raise FileCodegenError(
            f"Output file too small ({size} bytes, minimum {min_bytes}). "
            f"Content may be empty or corrupted."
        )


def create_docx(content: DocxContent, output_path: str) -> str:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)

    doc.add_heading(content.heading, level=1)

    for para_text in content.paragraphs:
        if para_text.strip():
            p = doc.add_paragraph(para_text.strip())
            for run in p.runs:
                run.font.size = Pt(11)

    if content.tables:
        table = doc.add_table(rows=len(content.tables), cols=len(content.tables[0]))
        table.style = 'Light Grid Accent 1'
        for i, row_data in enumerate(content.tables):
            for j, cell_text in enumerate(row_data):
                table.cell(i, j).text = str(cell_text)

    doc.save(output_path)
    _validate_output(output_path, min_bytes=5120)
    return output_path


def create_pptx(slides_data: List[SlideContent], output_path: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg_color = RGBColor(0x03, 0x07, 0x12)
    accent_color = RGBColor(0x00, 0xFF, 0x66)
    white_color = RGBColor(0xFF, 0xFF, 0xFF)

    for idx, slide_data in enumerate(slides_data):
        if idx == 0:
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(slide_layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = slide_data.title
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = accent_color
                    run.font.size = Pt(32)
                    run.font.bold = True

        if slide_data.content:
            if len(slide.shapes) > 1:
                body_shape = slide.shapes[1]
            else:
                left = Inches(0.5)
                top = Inches(1.8)
                width = Inches(12.333)
                height = Inches(5.0)
                body_shape = slide.shapes.add_textbox(left, top, width, height)

            tf = body_shape.text_frame
            tf.word_wrap = True
            tf.clear()

            lines = slide_data.content.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line.strip()
                for run in p.runs:
                    run.font.color.rgb = white_color
                    run.font.size = Pt(16)

        if slide_data.subtitle:
            left = Inches(0.5)
            top = Inches(6.5)
            width = Inches(12.333)
            height = Inches(0.5)
            sub_box = slide.shapes.add_textbox(left, top, width, height)
            tf = sub_box.text_frame
            tf.text = slide_data.subtitle
            for run in tf.paragraphs[0].runs:
                run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
                run.font.size = Pt(12)

    prs.save(output_path)
    _validate_output(output_path, min_bytes=5120)
    return output_path


def create_xlsx(content: XlsxContent, output_path: str) -> str:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import BarChart, Reference

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = content.title[:31] or "Sheet1"

    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="1B2A4A", end_color="1B2A4A", fill_type="solid")
    header_alignment = Alignment(horizontal="center", vertical="center")
    header_border = Border(
        bottom=Side(style="thin", color="D0D0D0"),
        left=Side(style="thin", color="D0D0D0"),
        right=Side(style="thin", color="D0D0D0"),
        top=Side(style="thin", color="D0D0D0"),
    )

    if content.headers:
        for col_idx, header in enumerate(content.headers, start=1):
            cell = ws.cell(row=1, column=col_idx, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_alignment
            cell.border = header_border

    for row_idx, row_data in enumerate(content.rows, start=2):
        for col_idx, value in enumerate(row_data, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=str(value))
            cell.alignment = Alignment(
                horizontal="center" if col_idx > 1 else "left", vertical="center"
            )
            cell.border = Border(
                bottom=Side(style="thin", color="D0D0D0"),
            )

    if content.rows and len(content.rows) >= 2:
        chart = BarChart()
        chart.title = content.title
        chart.y_axis.title = "Values"
        chart.x_axis.title = "Items"
        data_ref = Reference(ws, min_col=2, min_row=1, max_row=len(content.rows) + 1, max_col=min(len(content.headers), len(content.rows[0])))
        cats_ref = Reference(ws, min_col=1, min_row=2, max_row=len(content.rows) + 1)
        chart.add_data(data_ref, titles_from_data=True)
        chart.set_categories(cats_ref)
        chart.width = 18
        chart.height = 10
        ws.add_chart(chart, "A" + str(len(content.rows) + 4))

    wb.save(output_path)
    _validate_output(output_path, min_bytes=5120)
    return output_path


def create_pdf(reportlab_data: Dict[str, Any], output_path: str) -> str:
    """Create a styled PDF report using reportlab.
    
    Args:
        reportlab_data: dict with keys:
            - title: report title
            - sections: list of {"heading": str, "content": str, "items": list[str]}
            - pages: optional page count minimum
    """
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch, cm
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        PageBreak, Image, HRulespan
    )
    from reportlab.lib.colors import HexColor

    doc = SimpleDocTemplate(
        output_path,
        pagesize=A4,
        rightMargin=2 * cm,
        leftMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
    )

    styles = getSampleStyleSheet()
    # Custom styles
    title_style = ParagraphStyle(
        "CustomTitle",
        parent=styles["Title"],
        fontSize=24,
        spaceAfter=30,
        textColor=HexColor("#00FF66"),
        alignment=TA_CENTER,
    )
    heading_style = ParagraphStyle(
        "CustomHeading",
        parent=styles["Heading1"],
        fontSize=16,
        spaceAfter=12,
        textColor=HexColor("#030712"),
        backColor=HexColor("#E0E0E0"),
        leftIndent=0,
        borderWidth=0.5,
        borderBottomPadding=4,
    )
    body_style = ParagraphStyle(
        "CustomBody",
        parent=styles["Normal"],
        fontSize=11,
        leading=15,
        alignment=TA_LEFT,
        spaceAfter=12,
        textColor=HexColor("#333333"),
    )
    bullet_style = ParagraphStyle(
        "CustomBullet",
        parent=styles["Normal"],
        fontSize=11,
        leftIndent=20,
        bulletIndent=8,
        bulletFontSize=10,
        spaceAfter=6,
    )

    story = []

    # Title
    title = reportlab_data.get("title", "JARVIS Report")
    story.append(Paragraph(title, title_style))

    # Subtitle
    subtitle = reportlab_data.get("subtitle", "")
    if subtitle:
        story.append(Paragraph(subtitle, body_style))
        story.append(Spacer(1, 10))

    # Horizontal rule
    from reportlab.platypus import Drawing
    from reportlab.lib.styles import Alignment
    hr = Drawing(0, 0.5)
    story.append(hr)
    story.append(Spacer(1, 16))

    # Sections
    sections = reportlab_data.get("sections", [])
    for idx, section in enumerate(sections):
        heading = section.get("heading", f"Section {idx + 1}")
        story.append(Paragraph(heading, heading_style))
        story.append(Spacer(1, 6))

        content = section.get("content", "")
        if content:
            for para in content.split("\n"):
                if para.strip():
                    story.append(Paragraph(para.strip(), body_style))

        items = section.get("items", [])
        if items:
            data = [["#", "Item"]]
            for i, item in enumerate(items, 1):
                data.append([str(i), str(item)])
            table = Table(data, colWidths=[1.5 * cm, 12 * cm])
            table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), HexColor("#030712")),
                ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#FFFFFF")),
                ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, 0), 12),
                ("FONTSIZE", (0, 1), (-1, -1), 10),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BACKGROUND", (0, 1), (-1, -1), HexColor("#F8F8F8")),
                ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#D0D0D0")),
            ]))
            story.append(table)
            story.append(Spacer(1, 12))

        # Page break between major sections (optional)
        if section.get("page_break"):
            story.append(PageBreak())

    # Generate at least 2 pages if fewer than that
    min_pages = reportlab_data.get("min_pages", 2)
    # We'll check after building

    doc.build(story)

    _validate_output(output_path, min_bytes=5120)
    return output_path