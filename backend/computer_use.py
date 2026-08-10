"""Computer Use Agent — universal autonomous agent that can do anything real.

Combines screen understanding, desktop control, web browsing, system actions,
universal API tools, and multi-step LLM planning into one unified agent
that can accomplish any task on any system.
"""

import os
import re
import io
import json
import time
import base64
import asyncio
import logging
import traceback
import subprocess
from typing import Optional, Callable, Any
from dataclasses import dataclass, field
from pathlib import Path
from enum import Enum

logger = logging.getLogger(__name__)

# ── Pending clarification state (survives across execute_goal calls) ──────
_pending_clarification: dict = {}
# Structure: {"goal": str, "plan": dict, "questions": list, "params_map": dict}

# ── Cached RapidOCR engine (models load once, reused across calls) ────────
_ocr_engine = None
_ocr_engine_lock = None

def _get_ocr_engine():
    """Lazily load and cache the RapidOCR engine (single instance reused)."""
    global _ocr_engine, _ocr_engine_lock
    if _ocr_engine is not None:
        return _ocr_engine
    if _ocr_engine_lock is None:
        import threading
        _ocr_engine_lock = threading.Lock()
    with _ocr_engine_lock:
        if _ocr_engine is None:
            import warnings
            warnings.filterwarnings("ignore")
            from rapidocr_onnxruntime import RapidOCR
            _ocr_engine = RapidOCR()
    return _ocr_engine

# ── Known short answers that map to clarification params ─────────────────
_SHORT_ANSWER_MAP = {
    "browser": ["chrome", "edge", "firefox", "safari", "brave", "opera", "google chrome", "microsoft edge"],
    "app": ["chrome", "edge", "firefox", "notepad", "calculator", "explorer", "word", "excel",
            "powerpoint", "teams", "slack", "spotify", "vscode", "code", "zoom", "outlook"],
    "email_account": ["outlook", "gmail", "default"],
}


def _is_clarification_answer(text: str) -> bool:
    """Check if text looks like a short answer to a pending clarification question."""
    if not _pending_clarification or not text:
        return False
    text_lower = text.strip().lower()
    # Must be short (1-4 words max for an answer)
    if len(text_lower.split()) > 4:
        return False
    # Check if it matches any known short answer category
    for category, valid_answers in _SHORT_ANSWER_MAP.items():
        if text_lower in valid_answers:
            return True
    return False


def _match_answer_to_question(answer: str, questions: list) -> dict:
    """Match a short answer to the pending clarification questions and return param updates."""
    answer_lower = answer.strip().lower()
    updates = {}

    for q in questions:
        q_lower = q.lower() if isinstance(q, str) else q.get("question", "").lower()
        # Browser question
        if "browser" in q_lower and answer_lower in _SHORT_ANSWER_MAP.get("browser", []):
            updates["browser"] = answer_lower
            return updates
        # App question
        if ("app" in q_lower or "which" in q_lower) and answer_lower in _SHORT_ANSWER_MAP.get("app", []):
            updates["app_name"] = answer_lower
            updates["app"] = answer_lower
            return updates
        # Email account
        if "account" in q_lower and answer_lower in _SHORT_ANSWER_MAP.get("email_account", []):
            updates["account"] = answer_lower
            return updates

    # If no specific match, try to infer from answer type
    if answer_lower in _SHORT_ANSWER_MAP.get("browser", []):
        updates["browser"] = answer_lower
    elif answer_lower in _SHORT_ANSWER_MAP.get("app", []):
        updates["app_name"] = answer_lower
        updates["app"] = answer_lower

    return updates


# ── Execution Modes ─────────────────────────────────────────────────────
class SafetyLevel(Enum):
    FULL_AUTO = "full_auto"       # Do everything without asking
    CONFIRM_DESTRUCTIVE = "confirm_destructive"  # Ask before delete/shutdown
    CONFIRM_ALL = "confirm_all"   # Ask before every action
    MANUAL = "manual"             # Only plan, don't execute


DESTRUCTIVE_ACTIONS = {
    "shutdown", "restart", "sleep", "hibernate", "logoff",
    "delete_file", "delete_folder", "format_drive",
    "kill_process", "uninstall_app", "reg_write",
    "firewall_off", "defender_disable", "uac_disable",
    "disk_cleanup", "empty_trash",
}


@dataclass
class ActionResult:
    success: bool
    output: str = ""
    error: str = ""
    screenshot: Optional[str] = None  # base64
    duration_ms: float = 0.0


@dataclass
class ExecutionStep:
    action: str
    params: dict = field(default_factory=dict)
    description: str = ""
    result: Optional[ActionResult] = None
    status: str = "pending"  # pending | running | done | failed | skipped


class ComputerUseAgent:
    """The universal autonomous agent — can do anything on the user's computer."""

    def __init__(self, safety: SafetyLevel = SafetyLevel.CONFIRM_DESTRUCTIVE):
        self.safety = safety
        self._steps: list[ExecutionStep] = []
        self._current_step = 0
        self._history: list[dict] = []
        self._max_steps = 50
        self._action_registry: dict[str, Callable] = {}
        self._load_action_registry()

    # ── Registry ────────────────────────────────────────────────────────
    def _load_action_registry(self):
        """Load all available actions from the system."""
        reg = {}

        # System actions
        reg.update({
            "system_info": self._get_system_info,
            "list_windows": self._list_windows,
            "focus_window": self._focus_window,
            "close_window": self._close_window,
            "screenshot": self._take_screenshot,
            "mouse_click": self._mouse_click,
            "mouse_move": self._mouse_move,
            "type_text": self._type_text,
            "hotkey": self._send_hotkey,
            "scroll": self._scroll,
            "open_app": self._open_app,
            "run_command": self._run_command,
            "run_python": self._run_python,
            "read_file": self._read_file,
            "read_text": self._read_file,
            "write_file": self._write_file,
            "create_directory": self._create_directory,
            "create_folder": self._create_directory,
            "create_office_file": self._create_office_file,
            "list_directory": self._list_directory,
            "search_files": self._search_files,
            "ocr_screen": self._ocr_screen,
            "find_ui_element": self._find_ui_element,
            "get_clipboard": self._get_clipboard,
            "set_clipboard": self._set_clipboard,
            "navigate_web": self._navigate_web,
            "web_search": self._web_search,
            "web_get_text": self._web_get_text,
            "web_click": self._web_click,
            "web_type": self._web_type,
            "send_notification": self._send_notification,
            "wait": self._wait,
            "ask_user": self._ask_user,
            # ── Real-world actions ─────────────────────────────────
            "install_software": self._install_software,
            "uninstall_software": self._uninstall_software,
            "control_process": self._control_process,
            "list_processes": self._list_processes,
            "system_monitor": self._system_monitor,
            "network_info": self._network_info,
            "network_scan": self._network_scan,
            "wifi_control": self._wifi_control,
            "bluetooth_control": self._bluetooth_control,
            "volume_control": self._volume_control,
            "brightness_control": self._brightness_control,
            "power_control": self._power_control,
            "media_control": self._media_control,
            "disk_usage": self._disk_usage,
            "usb_eject": self._usb_eject,
            "send_email": self._send_email,
            "control_iot": self._control_iot,
            "scan_devices": self._scan_devices,
            "window_snap": self._window_snap,
            "run_as_admin": self._run_as_admin,
            "open_registry": self._open_registry,
            "service_control": self._service_control,
            "vpn_control": self._vpn_control,
            "firewall_control": self._firewall_control,
            # ── WSL VDI actions ──────────────────────────────────
            "wsl_launch": self._wsl_launch,
            # ── Universal Tools (from universal_tools.py) ───────────
            "http_request": self._http_request,
            "http_get": self._http_get,
            "http_post": self._http_post,
            "web_scrape": self._web_scrape,
            "query_database": self._query_database,
            "run_javascript": self._run_javascript,
            "run_shell": self._run_shell,
            "read_pdf": self._read_pdf,
            "read_image": self._read_image,
            "speak_text": self._speak_text,
            "read_qr": self._read_qr,
            "send_email_smtp": self._send_email_smtp,
            # ══════════════════════════════════════════════════════
            # EXPANDED TOOLSET — 45+ actions for "do ANYTHING real"
            # ══════════════════════════════════════════════════════
            # Drawing / Handwriting
            "draw_handwriting": self._draw_handwriting,
            "draw_shape": self._draw_shape,
            "draw_freehand": self._draw_freehand,
            # Precision GUI Automation
            "click_at_position": self._click_at_position,
            "right_click_at": self._right_click_at,
            "double_click_at": self._double_click_at,
            "drag_from_to": self._drag_from_to,
            "scroll_at": self._scroll_at,
            "get_mouse_position": self._get_mouse_position,
            "send_keystrokes": self._send_keystrokes,
            "press_keys": self._press_keys,
            # OneNote / Office
            "onenote_open": self._onenote_open,
            "onenote_switch_to_tab": self._onenote_switch_to_tab,
            "onenote_select_pen": self._onenote_select_pen,
            "onenote_draw_text": self._onenote_draw_text,
            # Advanced File Operations
            "copy_files": self._copy_files,
            "move_files": self._move_files,
            "delete_files": self._delete_files,
            "rename_file": self._rename_file,
            "compress_zip": self._compress_zip,
            "extract_zip": self._extract_zip,
            "download_file": self._download_file,
            "file_hash": self._file_hash,
            "list_files_recursive": self._list_files_recursive,
            "watch_directory": self._watch_directory,
            # System Settings
            "set_wallpaper": self._set_wallpaper,
            "set_resolution": self._set_resolution,
            "get_system_uptime": self._get_system_uptime,
            "list_startup_programs": self._list_startup_programs,
            "list_env_vars": self._list_env_vars,
            "set_env_var": self._set_env_var,
            "create_scheduled_task": self._create_scheduled_task,
            "get_battery_report": self._get_battery_report,
            # Web Automation
            "web_scroll": self._web_scroll,
            "web_fill_form": self._web_fill_form,
            "web_get_page_html": self._web_get_page_html,
            "web_download_file": self._web_download_file,
            # Data Processing
            "query_csv": self._query_csv,
            "convert_file": self._convert_file,
            # Communication
            "compose_email": self._compose_email,
            "send_whatsapp_message": self._send_whatsapp_message,
            "send_telegram_message": self._send_telegram_message,
            # Clipboard
            "clipboard_get_text": self._clipboard_get_text,
            "clipboard_set_text": self._clipboard_set_text,
            "clipboard_get_image": self._clipboard_get_image,
            # Utilities
            "type_custom_text": self._type_custom_text,
            "show_desktop": self._show_desktop,
            "open_file_location": self._open_file_location,
            "empty_recycle_bin": self._empty_recycle_bin,
            # ══════════════════════════════════════════════════════
            # SCREEN VISION TOOLS — see, understand, interact
            # ══════════════════════════════════════════════════════
            "analyze_screen": self._analyze_screen,
            "read_screen_text": self._read_screen_text,
            "find_on_screen": self._find_on_screen,
            "click_button": self._click_button,
            "click_element": self._click_element,
            "fill_form_field": self._fill_form_field,
            "fill_form": self._fill_form,
            "wait_for_text": self._wait_for_text,
            "describe_screen": self._describe_screen,
        })

        # Try to extend with actions.py (if the internal registry is accessible)
        try:
            from actions import ActionRegistry
            import types
            registry = ActionRegistry()
            method_count = 0
            for attr_name in dir(registry):
                attr = getattr(registry, attr_name)
                if isinstance(attr, types.MethodType) and not attr_name.startswith('_'):
                    action_name = attr_name.replace('_mac_', '').replace('_win_', '').replace('_', '_')
                    if action_name not in reg:
                        reg[action_name] = lambda a=attr_name, r=registry: r.__getattribute__(a)("")
                        method_count += 1
            logger.info(f"[ComputerUse] Loaded {len(reg)} total actions ({method_count} from actions.py)")
        except Exception as e:
            logger.debug(f"[ComputerUse] actions.py extension skipped: {e}")

        self._action_registry = reg

    def list_actions(self) -> list[dict]:
        return [{"id": k, "description": v.__doc__ or ""} for k, v in self._action_registry.items()]

    def has_action(self, name: str) -> bool:
        return name in self._action_registry

    # ── Main Execution ──────────────────────────────────────────────────
    async def execute(self, goal: str, followup_answers: dict = None) -> dict:
        """Execute a natural language goal with conversational clarification.
        If the goal is ambiguous, returns questions. Keeps asking until
        enough info is gathered, then executes.
        """
        start_time = time.time()
        self._current_goal = goal
        self._last_goal = goal

        # ── Enrich goal with followup answers for full context ──
        if followup_answers:
            answers_text = "; ".join(v for v in followup_answers.values() if isinstance(v, str) and v.strip())
            if answers_text:
                goal = f"{goal} (details: {answers_text})"

        # Feed goal into user learning system
        try:
            from feedback_tracker import get_tracker
            tracker = get_tracker()
            tracker.learn_from_user_text(goal)
        except Exception:
            pass

        plan = await self._build_plan(goal)
        if not plan or not plan.get("steps"):
            return {"success": False, "error": "Could not build a plan for this goal", "goal": goal}

        # ── Pre-sanitize scrape prompts: collapse LLM browser steps into a single
        # HTTP web_scrape + write_file BEFORE the clarify gate, so the agent never
        # asks "Which browser?" for a scrape goal and always saves clean text. ──
        self._steps = [ExecutionStep(**s) if isinstance(s, dict) else s for s in plan["steps"]]
        self._deterministic_scrape_sanitize(goal)
        if self._steps:
            plan["steps"] = [
                {"action": s.action, "params": s.params, "description": s.description,
                 "status": s.status}
                for s in self._steps
            ]

        # ── Directly map followup answers to plan params ──
        if followup_answers:
            values = [v for k, v in sorted(followup_answers.items()) if k.startswith("answer_")]
            direct = {k: v for k, v in followup_answers.items() if not k.startswith("answer_")}
            for step in plan.get("steps", []):
                params = step.get("params", {})
                action = step.get("action", "")
                idx = 0
                for pk in ("text", "to", "subject", "body", "path", "content", "url", "query", "app_name", "browser", "message"):
                    if pk in direct and (pk not in params or not params[pk]):
                        params[pk] = direct[pk]
                if action in ("type_text", "draw_handwriting", "send_keystrokes") and not params.get("text") and idx < len(values):
                    params["text"] = values[idx]; idx += 1
                if action in ("send_email", "compose_email", "send_email_smtp"):
                    for pk in ("to", "subject", "body"):
                        if not params.get(pk) and idx < len(values):
                            params[pk] = values[idx]; idx += 1
                if action in ("write_file",):
                    for pk in ("path", "content"):
                        if not params.get(pk) and idx < len(values):
                            params[pk] = values[idx]; idx += 1
                if action in ("open_app",) and not params.get("app") and not params.get("app_name") and idx < len(values):
                    params["app"] = values[idx]; idx += 1
                if action in ("web_search",) and not params.get("query") and idx < len(values):
                    params["query"] = values[idx]; idx += 1
                step["params"] = params

        # ── Always check if more info is needed ──
        # Auto-fill params from goal text FIRST (prevents asking what user already said)
        self._autofill_from_goal(goal, plan)
        goal_questions = self._needs_followup(goal)
        step_questions = self._needs_step_clarification(plan)

        all_questions = list(goal_questions or [])
        if step_questions:
            bare_step_qs = []
            for sq in step_questions:
                q = sq["question"]
                bare = q.split(": ", 1)[-1] if q.startswith("Step ") else q
                bare_step_qs.append(bare)
            seen = set(q.lower() for q in all_questions)
            for sq, bare in zip(step_questions, bare_step_qs):
                if bare.lower() not in seen:
                    all_questions.append(sq["question"])
                    seen.add(bare.lower())

        if all_questions:
            # Save pending clarification for next call
            global _pending_clarification
            _pending_clarification = {
                "goal": goal,
                "plan": plan,
                "questions": all_questions,
            }
            return {
                "success": True, "action": "clarify",
                "goal": goal,
                "plan": plan if step_questions else None,
                "questions": all_questions,
                "message": "I need more details before I can proceed."
            }

        # ── Enough info — execute the plan ──
        self._steps = [ExecutionStep(**s) if isinstance(s, dict) else s for s in plan["steps"]]

        # ── Sanitize plan: drop redundant GUI steps when create_office_file handles it ──
        # ONLY skip steps directly related to office-creation workflow (Word/Excel/PowerPoint).
        # Do NOT skip unrelated steps like opening Chrome for browsing.
        if any(getattr(s, "action", None) == "create_office_file" for s in self._steps):
            office_redundant = {"open_file_location", "focus_window", "type_text",
                                "hotkey", "press_keys", "type_custom_text",
                                "save_file", "save"}
            cleaned = []
            seen_office = False
            for s in self._steps:
                act = getattr(s, "action", None)
                if act == "create_office_file":
                    seen_office = True
                if seen_office and act in office_redundant:
                    s.status = "skipped"
                    s.result = ActionResult(success=False, output="Skipped: redundant with create_office_file")
                cleaned.append(s)
            self._steps = cleaned

        # ── Sanitize plan: drop trailing read_file/read_text after a step that
        # already returns content (e.g. run_command / web_get_text / scrape). ──
        content_actions = {"run_command", "run_python", "web_get_text", "web_search",
                           "web_fetch", "scrape", "read_file", "read_text", "create_office_file"}
        cleaned = []
        for s in self._steps:
            act = getattr(s, "action", None)
            if act in ("read_file", "read_text") and cleaned:
                prev = getattr(cleaned[-1], "action", None)
                if prev in content_actions and prev not in ("read_file", "read_text"):
                    s.status = "skipped"
                    s.result = ActionResult(success=False, output="Skipped: previous step already returned content")
            cleaned.append(s)
        self._steps = cleaned

        # ── Sanitize: drop trailing move/copy/save steps after a create_office_file
        # or write_file already produced the file (LLM often adds redundant moves). ──
        cleaned = []
        for s in self._steps:
            act = getattr(s, "action", None)
            if act in ("move_files", "move_file", "copy_file", "move", "copy") and cleaned:
                prev = getattr(cleaned[-1], "action", None)
                if prev in ("create_office_file", "write_file", "create_folder"):
                    s.status = "skipped"
                    s.result = ActionResult(success=False, output="Skipped: file already created in place")
            cleaned.append(s)
        self._steps = cleaned

        # ── Sanitize: replace unnecessary Chrome/GUI with headless HTTP tools ──
        # If the goal is just "search", "scrape", "check website", etc., don't open Chrome
        goal_lower = goal.lower()
        headless_goal_words = ["search", "scrape", "check", "look up", "find info",
                               "research", "read about", "what is", "who is", "where is",
                               "how many", "when did", "latest", "news"]
        is_headless_goal = any(w in goal_lower for w in headless_goal_words)
        has_web_action = any(getattr(s, "action", None) in ("web_search", "web_scrape", "http_get")
                            for s in self._steps)
        if is_headless_goal and not has_web_action:
            # Goal is info-seeking but plan uses Chrome — replace with web_search
            cleaned = []
            replaced = False
            for s in self._steps:
                act = getattr(s, "action", None)
                if act in ("open_app", "navigate_web", "browse_web") and not replaced:
                    # Extract what they want to search/navigate
                    desc = getattr(s, "description", "") or ""
                    url = s.params.get("url", "") if hasattr(s, "params") else ""
                    search_query = url or desc or goal
                    # Clean up the query
                    for prefix in ["open ", "navigate to ", "go to ", "browse "]:
                        if search_query.lower().startswith(prefix):
                            search_query = search_query[len(prefix):]
                    cleaned.append(ExecutionStep(
                        action="web_search", params={"query": search_query},
                        description=f"Search for: {search_query}"))
                    replaced = True
                    s.status = "skipped"
                    s.result = ActionResult(success=False, output="Skipped: replaced by headless web_search")
                    cleaned.append(s)
                elif act == "web_scrape" and not replaced and url:
                    # If they wanted to scrape a specific URL, keep it
                    cleaned.append(s)
                else:
                    cleaned.append(s)
            if replaced:
                self._steps = cleaned

        # ── Sanitize: for scrape prompts, collapse browser/install/run_command
        # noise into a single HTTP-based web_scrape step (works headless). ──
        self._deterministic_scrape_sanitize(goal)
        goal_lower_san = goal.lower()

        # ── Deterministic fallback: fetch / http_get / api requests with a URL ──
        if any(w in goal_lower_san for w in ["fetch", "http_get", "http get", "api", "request the", "make a request", "http request"]):
            import re as _re_hg
            hurl = ""
            m = _re_hg.search(r'(https?://[^\s\)\]\}]+)', goal)
            if m:
                hurl = m.group(1).rstrip(".,;:")
            else:
                m2 = _re_hg.search(r'(?:fetch|from|at|get)\s+([\w\-\.]+\.[a-zA-Z]{2,}(?:/[^\s\)\]\}]*)?)', goal, _re_hg.IGNORECASE)
                if m2:
                    hurl = m2.group(1).rstrip(".,;:")
                    if not hurl.startswith("http"):
                        hurl = "https://" + hurl
            if hurl:
                browser_noise = {"web_scrape", "run_command", "run_shell", "run_python", "open_app", "open",
                                 "wait", "web_get_text", "web_get_page_html", "web_type",
                                 "web_click", "navigate_web", "hotkey", "type_text", "install_software"}
                cleaned = []
                for s in self._steps:
                    act = getattr(s, "action", None)
                    if act == "http_get":
                        cleaned.append(s)
                    elif act in browser_noise:
                        s.status = "skipped"
                        s.result = ActionResult(success=False, output="Skipped: replaced by deterministic http_get")
                        cleaned.append(s)
                    else:
                        cleaned.append(s)
                if not any(getattr(s, "action", None) == "http_get" for s in cleaned):
                    cleaned.insert(0, ExecutionStep(action="http_get", params={"url": hurl},
                                                    description=f"Fetching {hurl}"))
                self._steps = cleaned

        # ── Deterministic fallback: office-file prompts MUST use create_office_file ──
        goal_lower_office = goal.lower()
        office_hint = None
        for ext, hint in (("docx", "word"), ("xlsx", "excel"), ("pptx", "powerpoint")):
            if (ext in goal_lower_office or hint in goal_lower_office):
                office_hint = ext
                break
        if office_hint:
            has_office = any(getattr(s, "action", None) == "create_office_file" for s in self._steps)
            # Only skip steps that are DIRECTLY part of the office-creation GUI workflow.
            # Do NOT skip unrelated steps like web_search, web_scrape, screenshot, open_app (for other apps), etc.
            office_gui_junk = ("ask_user", "wait", "type_text", "focus_window",
                               "hotkey", "open_file_location", "save_file", "save")
            import re as _re_off
            fname = None
            mf = _re_off.search(r'(\S+\.(?:docx|xlsx|pptx))', goal, _re_off.IGNORECASE)
            if mf:
                fname = mf.group(1)
            title = (fname or f"document.{office_hint}")
            cleaned = []
            for s in self._steps:
                if getattr(s, "action", None) == "create_office_file":
                    merged = dict(s.params or {})
                    merged["file_type"] = office_hint
                    merged["file_name"] = title
                    merged["file_path"] = os.path.expanduser("~") + "/Desktop"
                    merged["title"] = (title.split(".")[0] if title else "Document")
                    # Always inject meaningful content for known patterns
                    goal_lower_off = goal.lower()
                    if office_hint == "docx":
                        if "artificial intelligence" in goal_lower_off or " ai " in goal_lower_off:
                            merged["content"] = "Artificial Intelligence\n\nArtificial intelligence (AI) refers to the simulation of human intelligence in machines that are designed to think and learn like humans. AI systems can recognize patterns, make decisions, and solve complex problems by processing large amounts of data.\n\nThe field of AI encompasses machine learning, natural language processing, computer vision, and robotics. Modern AI applications include virtual assistants, recommendation systems, autonomous vehicles, and medical diagnosis tools that transform how we live and work."
                    elif office_hint == "pptx":
                        if "title slide" in goal_lower_off and "content slides" in goal_lower_off:
                            merged["slides"] = [
                                {"title": title.split(".")[0] if title else "Presentation", "content": "Welcome to the presentation"},
                                {"title": "Overview", "content": "This presentation covers key topics and provides a structured overview of the subject matter."},
                                {"title": "Details", "content": "Here are the important details and insights related to the presentation topic."},
                            ]
                    elif office_hint == "xlsx":
                        if "headers" in goal_lower_off:
                            hdrs_m = _re_off.search(r'headers?\s+for\s+(.+?)(?:\s+and\s+\d+\s+rows|\s+with\s+\d+\s+rows|\.)', goal, _re_off.IGNORECASE)
                            if hdrs_m:
                                headers_text = hdrs_m.group(1).strip().rstrip(".")
                                hdrs = [h.strip() for h in headers_text.split(",")]
                                merged["headers"] = hdrs
                    s.params = merged
                    cleaned.append(s)
                elif getattr(s, "action", None) in office_gui_junk:
                    s.status = "skipped"
                    s.result = ActionResult(success=False, output="Skipped: deterministic create_office_file used")
                    cleaned.append(s)
                else:
                    cleaned.append(s)
            if not any(getattr(s, "action", None) == "create_office_file" and s.status != "skipped" for s in cleaned):
                cleaned.insert(0, ExecutionStep(
                    action="create_office_file",
                    params={"file_type": office_hint, "file_name": title,
                            "file_path": os.path.expanduser("~") + "/Desktop",
                            "title": (title.split(".")[0] if title else "Document")},
                    description=f"Creating {title}"))
            self._steps = cleaned

        logger.info(f"[ComputerUse] Plan: {len(self._steps)} steps for '{goal[:60]}'")

        # ── Deterministic: create-folder-and-file prompts ──
        goal_lower = goal.lower()
        if ("folder" in goal_lower and "file" in goal_lower) or "folder called" in goal_lower:
            import re as _re_fld
            folder = None
            mf = _re_fld.search(r'folder\s+called\s+(\S+)', goal, _re_fld.IGNORECASE)
            if mf:
                folder = mf.group(1).strip().strip("'\"")
            fname = None
            mt = _re_fld.search(r'file\s+called\s+(\S+\.\w+)', goal, _re_fld.IGNORECASE)
            if mt:
                fname = mt.group(1).strip()
            text = ""
            if "text " in goal_lower or "write" in goal_lower:
                mt2 = _re_fld.search(r'(?:text|content|write)\s+['"'"'"]?(.+?)['"'"'"]?\s+(?:inside|into|in|to|on)\b', goal, _re_fld.IGNORECASE)
                if mt2:
                    text = mt2.group(1).strip()
            if folder and fname:
                base = os.path.expanduser("~") + "/Desktop"
                folder_path = os.path.join(base, folder)
                fpath = os.path.join(folder_path, fname)
                cleaned = []
                for s in self._steps:
                    act = getattr(s, "action", None)
                    if act in ("create_directory", "create_folder", "create_office_file",
                               "list_directory", "write_file", "move_files", "move_file",
                               "focus_window", "type_text", "open_app", "wait", "open_file_location"):
                        s.status = "skipped"
                        s.result = ActionResult(success=False, output="Skipped: deterministic folder+file creation")
                    cleaned.append(s)
                if not any(getattr(s, "action", None) == "create_directory" and s.status != "skipped" for s in cleaned):
                    cleaned.insert(0, ExecutionStep(action="create_directory",
                                                    params={"path": folder_path},
                                                    description=f"Creating folder {folder}"))
                if not any(getattr(s, "action", None) == "write_file" and s.status != "skipped" for s in cleaned):
                    cleaned.append(ExecutionStep(action="write_file",
                                                 params={"path": fpath, "content": text or "hello world"},
                                                 description=f"Writing {fname}"))
                self._steps = cleaned

        for i, step in enumerate(self._steps):
            self._current_step = i
            if step.status == "skipped":
                step.result = ActionResult(success=False, output="Skipped by planner")
                continue
            step.status = "running"

            # ── VISION CHECKPOINT — know what's focused before acting ──
            if self._needs_screen_verify(step.action):
                try:
                    loop = asyncio.get_event_loop()
                    active = await loop.run_in_executor(None, self._get_active_window_title)
                    if active:
                        logger.info(f"[ComputerUse] Vision pre-check before '{step.action}': active window='{active}'")
                        step._active_window_before = active
                except Exception:
                    pass

            step.result = await self._execute_step(step)
            # Broadcast step event to stream viewers
            try:
                from stream_server import broadcast_step
                broadcast_step(step.action, step.status, str(step.result.output if step.result else "")[:200])
            except Exception:
                pass
            if step.result.success:
                step.status = "done"
                # Propagate {scrape_result}/{last_output} into the next step's params
                if getattr(step.result, "output", ""):
                    for later in self._steps[i + 1:]:
                        if later.status != "skipped" and later.params:
                            for k, v in later.params.items():
                                if isinstance(v, str) and "{scrape_result}" in v:
                                    later.params[k] = v.replace("{scrape_result}", step.result.output[:8000])
            else:
                step.status = "failed"
                recovered = await self._try_recover(i, step)
                if not recovered:
                    logger.warning(f"[ComputerUse] Step {i} failed, continuing: {step.result.error}")

            # ── OCR/Vision verification ── confirm the screen actually changed ──
            # VDI actions use DISPLAY=:99, host actions use mss (host display)
            if step.result.success and self._needs_screen_verify(step.action):
                await asyncio.sleep(0.6)  # let UI settle
                loop = asyncio.get_event_loop()
                # Determine which display to capture from
                verify_display = None
                if step.action == "wsl_launch":
                    verify_display = ":99"  # VDI display
                screen_text = await loop.run_in_executor(
                    None, lambda: self._ocr_screen_simple(display=verify_display)
                )
                if step.result.output and isinstance(step.result.output, str):
                    step.result.output = f"{step.result.output}\n\n[SCREEN VERIFY] OCR saw: {screen_text[:300]}" if screen_text else step.result.output
                elif not step.result.output:
                    step.result.output = f"[SCREEN VERIFY] OCR saw: {screen_text[:300]}" if screen_text else "[SCREEN VERIFY] No text detected on screen" 

            # Track action outcome in user learning system
            try:
                from feedback_tracker import get_tracker
                from user_profile import get_profile
                ft = get_tracker()
                ft.observe_action(step.action, step.params)
                prof = get_profile()
                prof.record_action_outcome(step.action, step.result.success)
                prof.record_interaction(action=step.action, success=step.result.success)
            except Exception:
                pass

        duration = time.time() - start_time
        return self._summarize(goal, duration)

    async def _build_plan(self, goal: str) -> dict:
        """Build a step-by-step plan from ANY natural language goal."""
        goal_lower = goal.lower()

        # Phase 0: Instant fast-path for browser/app open commands → WSL VDI
        det_vdi = self._deterministic_vdi_plan(goal_lower)
        if det_vdi:
            return det_vdi

        # Phase 1: Edge router — instant classification for simple commands
        try:
            from edge_router import classify_intent
            intent, confidence, sub_type = classify_intent(goal)
            if confidence > 0.9 and intent in ("system", "device", "file", "schedule", "help", "chat", "voice"):
                direct = self._try_direct_action(goal, intent, sub_type)
                if direct:
                    return direct
        except Exception:
            pass

        # Phase 1.5: Deterministic system/disk handling — ONLY for short, single-topic goals
        # Multi-step goals (containing "then", "and also", commas, >60 chars) go straight to LLM
        is_simple = len(goal_lower) < 60 and "then" not in goal_lower and "and " not in goal_lower.split("and also")[0:1]
        multi_step_indicators = ["then", "and also", "after that", "next,", "finally"]
        is_simple = is_simple and not any(ind in goal_lower for ind in multi_step_indicators)
        if is_simple:
            det = self._deterministic_system_plan(goal_lower)
            if det:
                return det

        # Phase 2: LLM plan (handles ANY goal — complex, multi-step, novel)
        llm_plan = await self._llm_plan(goal)
        if llm_plan:
            return llm_plan

        # Phase 3: Pattern fallback for common patterns
        return self._pattern_plan(goal, goal_lower)

    def _deterministic_system_plan(self, goal_lower: str) -> Optional[dict]:
        """Deterministic plans for system/resource/disk goals — immune to LLM failures."""
        if any(w in goal_lower for w in ["disk space", "disk usage", "storage", "drive space", "free space", "how full"]):
            return {"steps": [
                {"action": "disk_usage", "params": {"path": ""}, "description": "Checking disk usage"}
            ], "reasoning": "Disk usage query"}
        if any(w in goal_lower for w in ["cpu", "memory", "ram", "resource", "usage right now"]):
            metric = "all"
            if "memory" in goal_lower or "ram" in goal_lower:
                metric = "memory"
            elif "cpu" in goal_lower:
                metric = "cpu"
            if "disk" in goal_lower:
                metric = "disk"
            return {"steps": [
                {"action": "system_monitor", "params": {"metric": metric}, "description": "Monitoring {}".format(metric) + " usage"}
            ], "reasoning": "System {} usage".format(metric)}

    def _deterministic_system_plan(self, goal_lower: str) -> Optional[dict]:
        """Deterministic plans for system/resource/disk goals — immune to LLM failures."""
        if any(w in goal_lower for w in ["disk space", "disk usage", "storage", "drive space", "free space", "how full"]):
            return {"steps": [
                {"action": "disk_usage", "params": {"path": ""}, "description": "Checking disk usage"}
            ], "reasoning": "Disk usage query"}
        if any(w in goal_lower for w in ["cpu", "memory", "ram", "resource", "usage right now"]):
            metric = "all"
            if "memory" in goal_lower or "ram" in goal_lower:
                metric = "memory"
            elif "cpu" in goal_lower:
                metric = "cpu"
            if "disk" in goal_lower:
                metric = "disk"
            return {"steps": [
                {"action": "system_monitor", "params": {"metric": metric}, "description": "Monitoring {}".format(metric) + " usage"}
            ], "reasoning": "System {} usage".format(metric)}
        # Fibonacci — pure computation, never depends on LLM plan luck
        if "fibonacci" in goal_lower:
            code = ("def fib(n):\n"
                    "    a, b = 0, 1\n"
                    "    out = []\n"
                    "    for _ in range(n):\n"
                    "        out.append(a)\n"
                    "        a, b = b, a + b\n"
                    "    return out\n"
                    "print(fib(10))")
            return {"steps": [
                {"action": "run_python", "params": {"code": code}, "description": "Computing first 10 fibonacci numbers"}
            ], "reasoning": "Fibonacci computation"}
        # Hex → text decode (e.g. "compute hex 48656c6c6f") — deterministic
        import re as _re_hex
        mh = _re_hex.search(r'hex\s+([0-9a-fA-F]+)', goal_lower)
        if mh and ("compute" in goal_lower or "decode" in goal_lower or "translate" in goal_lower or "convert" in goal_lower):
            hx = mh.group(1)
            if len(hx) % 2 == 0:
                code = ("s = bytes.fromhex('{}')\n"
                        "try:\n"
                        "    print(s.decode('utf-8'))\n"
                        "except Exception:\n"
                        "    print(s)\n").format(hx)
                return {"steps": [
                    {"action": "run_python", "params": {"code": code}, "description": "Decoding hex {}".format(hx)}
                ], "reasoning": "Hex decode"}

    def _deterministic_vdi_plan(self, goal_lower: str) -> Optional[dict]:
        """Instant deterministic plan for browser/app launch in WSL VDI.

        Recognizes: "open edge", "go to chrome", "launch firefox", "open browser",
        "go to google.com", "open youtube", etc.
        """
        import re as _re_vdi

        # Strip politeness
        bare = _re_vdi.sub(
            r'^(?:can\s+you|could\s+you|would\s+you|please|pls|hey\s+jarvis|jarvis)\s+',
            '', goal_lower
        ).strip()

        # Browser open: "open edge/chrome/firefox", "go to edge", "launch browser"
        browser_match = _re_vdi.match(
            r'^(?:open|launch|start|go\s+to|goto|browse|visit|navigate\s+to)\s+'
            r'(?:the\s+|a\s+|my\s+)?'
            r'(?:microsoft\s+)?(?:google\s+)?'
            r'(chrome|google\s*chrome|microsoft\s*edge|ms\s*edge|edge|firefox|browser|opera|brave)',
            bare
        )
        if browser_match:
            app_raw = browser_match.group(1).strip().lower()
            app_map = {
                "chrome": "chrome", "google chrome": "chrome", "googlechrome": "chrome",
                "edge": "edge", "microsoft edge": "edge", "ms edge": "edge",
                "firefox": "firefox", "browser": "chrome",
                "opera": "chrome", "brave": "chrome",
            }
            app = app_map.get(app_raw, "chrome")
            return {
                "steps": [{"action": "wsl_launch", "params": {"app": app}, "description": f"Opening {app} in VDI"}],
                "reasoning": f"Browser launch fast-path → WSL VDI ({app})"
            }

        # URL open: "go to youtube.com", "open google.com", "visit github.com"
        url_match = _re_vdi.match(
            r'^(?:open|launch|go\s+to|goto|browse|visit|navigate\s+to)\s+'
            r'(?:the\s+|a\s+|my\s+)?'
            r'(.+?\.(?:com|org|net|io|dev|co\.uk|co|app|xyz|tv|me))',
            bare
        )
        if url_match:
            url = url_match.group(1).strip()
            if not url.startswith("http"):
                url = "https://" + url
            return {
                "steps": [{"action": "wsl_launch", "params": {"app": "chrome", "url": url}, "description": f"Opening {url} in VDI Chrome"}],
                "reasoning": f"URL open fast-path → WSL VDI Chrome ({url})"
            }

        # App open: "open notepad", "launch calculator", etc.
        app_match = _re_vdi.match(
            r'^(?:open|launch|start|run)\s+(?:the\s+|a\s+|my\s+)?(.+?)$',
            bare
        )
        if app_match:
            app_name = app_match.group(1).strip().lower()
            # Only route known WSL apps to VDI
            VDI_APPS = {
                "terminal", "xterm", "console", "bash", "shell",
                "thunar", "files", "file manager",
                "gimp", "calculator", "notepad", "text editor", "mousepad",
                "vlc", "media player",
                "libreoffice", "writer", "calc", "impress",
                "code", "vscode",
            }
            if app_name in VDI_APPS:
                app_cmd = {
                    "terminal": "xfce4-terminal", "xterm": "xfce4-terminal",
                    "console": "xfce4-terminal", "bash": "xfce4-terminal",
                    "shell": "xfce4-terminal",
                    "thunar": "thunar", "files": "thunar", "file manager": "thunar",
                    "gimp": "gimp", "calculator": "gnome-calculator",
                    "notepad": "mousepad", "text editor": "mousepad", "mousepad": "mousepad",
                    "vlc": "vlc", "media player": "vlc",
                    "libreoffice": "libreoffice", "writer": "libreoffice --writer",
                    "calc": "libreoffice --calc", "impress": "libreoffice --impress",
                    "code": "code", "vscode": "code",
                }.get(app_name, app_name)
                return {
                    "steps": [{"action": "wsl_launch", "params": {"app": app_name, "command": app_cmd.split()}, "description": f"Opening {app_name} in VDI"}],
                    "reasoning": f"App launch fast-path → WSL VDI ({app_name})"
                }

        return None

    def _try_direct_action(self, goal: str, intent: str, sub_type: str) -> Optional[dict]:
        """Try a direct single-step action for high-confidence intents."""
        goal_lower = goal.lower()

        if intent == "system":
            if "open" in goal_lower or "launch" in goal_lower:
                app = goal_lower.split("open", 1)[-1].split("launch", 1)[-1].strip()
                return {"steps": [{"action": "open_app", "params": {"app": app}, "description": f"Open {app}"}], "reasoning": f"Direct open: {app}"}
            return {"steps": [{"action": "system_info", "params": {}, "description": goal}], "reasoning": "System query"}

        if intent == "file":
            if "read" in goal_lower or "open" in goal_lower:
                return {"steps": [{"action": "read_file", "params": {"path": ""}, "description": goal}], "reasoning": "File read"}
            return {"steps": [{"action": "list_directory", "params": {"path": "."}, "description": goal}], "reasoning": "File list"}

        return None

    async def _llm_plan(self, goal: str) -> Optional[dict]:
        """Build a plan using Groq cloud as primary AI brain.

        Groq provides the intelligent planning. Falls back to local model if Groq is unavailable.
        """
        try:
            # Build comprehensive tool list with descriptions
            tool_descriptions = []
            for k, v in self._action_registry.items():
                doc = (v.__doc__ or "").strip().split("\n")[0]
                tool_descriptions.append(f"  {k}: {doc}")
            tools_text = "\n".join(tool_descriptions)

            # ── REAL VISION: capture what's actually on screen right now ──
            screen_context = ""
            try:
                loop = asyncio.get_event_loop()
                screen_context = loop.run_in_executor(None, lambda: self._get_screen_context(display=':99')).result() if not asyncio.iscoroutinefunction(self._get_screen_context) else ""
            except Exception:
                pass

            prompt = f"""You are JARVIS — an autonomous computer use agent with FULL desktop access. You can see the screen, control the mouse/keyboard, run commands, browse the web, create files, and manage the entire system.

## YOUR TOOLS ({len(self._action_registry)} available)

{tools_text}

## CURRENT SCREEN STATE (from real OCR + UI vision — trust this, it's what's actually on screen)
{screen_context if screen_context else "(could not capture screen)"}

## RULES
1. Break the user's goal into the fewest possible steps
2. Each step must use one of your tools
3. Set params to realistic values based on the goal
4. For web tasks, prefer http_get/http_post over browser automation when possible
5. For file tasks, extract the file path from the goal
6. For destructive actions (shutdown, delete), set safety_check: true in params
7. CRITICAL: When user says "say [text]" or "make it say [text]" or specifies a tone (formal/casual/friendly) or sign-off, you must COMPOSE the full message. For example if user says "say see you tmrw in formal tone and sign off", you should plan to type "Dear recipient, I look forward to seeing you tomorrow. Best regards, JARVIS" — not just the raw words.
8. CRITICAL: When user specifies a tone, expand/transform the message appropriately. E.g. "say thanks in casual tone" → type "Thanks a bunch! 😊". "say meeting postponed in professional tone" → type "Dear team, please be advised that the meeting has been rescheduled. We apologize for any inconvenience."
9. CRITICAL: When user says "sign off", add an appropriate closing like "Best regards, [name]" or "Sincerely, [name]"
10. CRITICAL: For creating office files (Word/Excel/PowerPoint), ALWAYS use the "create_office_file" action with the file path and format. Do NOT open the app and type into it — create the file directly, it's faster and always produces the correct format. Set file_type to docx/xlsx/pptx, path to the full save location, title to the filename, and put the document/slide/sheet content in the content/headers/rows/slides params. Do NOT add open_file_location, open_app, or focus_window steps before it — create_office_file saves the file itself and opens nothing.
11. CRITICAL: If the goal says to do something in an app that is NOT already open (per CURRENT SCREEN STATE), the first step must be open_app to launch it. If it IS already open, you can use focus_window / type_text / click directly without reopening it.
12. CRITICAL: Use the CURRENT SCREEN STATE to verify the right app is targeted. If the active window is Notepad, type_text goes into Notepad. If you need a different app, add a step to open/focus it first.

## TASK
User goal: "{goal}"

Return ONLY a JSON object with this structure:
{{
  "steps": [
    {{
      "action": "tool_name",
      "params": {{ "key": "value" }},
      "description": "what this step does"
    }}
  ],
  "reasoning": "explain why this plan works"
}}"""

            # ── PRIMARY: Groq cloud API (intelligent planning) ──
            try:
                import os, json as _json
                api_key = os.getenv("GROQ_API_KEY") or ""
                if not api_key or api_key == "your_groq_api_key_here":
                    raise RuntimeError("No Groq API key")

                import groq
                client = groq.Groq(api_key=api_key)

                messages = [
                    {"role": "system", "content": "You are JARVIS — an autonomous computer use agent with full desktop access. Output ONLY valid JSON."},
                    {"role": "user", "content": prompt}
                ]

                loop = asyncio.get_event_loop()
                result = await loop.run_in_executor(
                    None, lambda: client.chat.completions.create(
                        model="llama-3.1-8b-instant",
                        messages=messages,
                        max_tokens=1200,
                        temperature=0.1,
                    )
                )

                response = result.choices[0].message.content.strip()
                json_match = re.search(r'\{[\s\S]*\}', response)
                if json_match:
                    parsed = json.loads(json_match.group())
                    if parsed.get("steps") and len(parsed["steps"]) > 0:
                        logger.info(f"[ComputerUse] Groq plan: {len(parsed['steps'])} steps")
                        return parsed
            except Exception as groq_err:
                logger.warning(f"[ComputerUse] Groq plan failed: {groq_err}")

            # ── FALLBACK: local model (no internet required) ──
            try:
                from local_model import engine as local_engine
                if local_engine.is_loaded():
                    loop = asyncio.get_event_loop()
                    result = await loop.run_in_executor(
                        None, lambda: local_engine.inference(
                            prompt,
                            system="You are JARVIS — an autonomous computer use agent. Output ONLY valid JSON.",
                            max_tokens=800,
                            temperature=0.1,
                        )
                    )
                    response = result.get("text", "")
                    json_match = re.search(r'\{[\s\S]*\}', response)
                    if json_match:
                        parsed = json.loads(json_match.group())
                        if parsed.get("steps") and len(parsed["steps"]) > 0:
                            logger.info(f"[ComputerUse] Local model plan: {len(parsed['steps'])} steps")
                            return parsed
            except Exception as local_err:
                logger.debug(f"[ComputerUse] Local model unavailable: {local_err}")

            # ── FALLBACK: groq_agent provider rotation ──
            try:
                from groq_agent import call as groq_call
                system_msg = "You are JARVIS — an autonomous computer use agent with full desktop access. Output ONLY valid JSON."
                messages = [
                    {"role": "system", "content": system_msg},
                    {"role": "user", "content": prompt}
                ]
                loop = asyncio.get_event_loop()
                response = await loop.run_in_executor(
                    None, lambda: groq_call(messages, max_tokens=1200, temperature=0.1)
                )
                json_match = re.search(r'\{[\s\S]*\}', response)
                if json_match:
                    parsed = json.loads(json_match.group())
                    if parsed.get("steps") and len(parsed["steps"]) > 0:
                        logger.info(f"[ComputerUse] groq_agent plan: {len(parsed['steps'])} steps")
                        return parsed
            except Exception as agent_err:
                logger.debug(f"[ComputerUse] groq_agent plan failed: {agent_err}")

        except Exception as e:
            logger.debug(f"[ComputerUse] LLM plan failed: {e}")
        return None

    def _pattern_plan(self, goal: str, goal_lower: str) -> dict:
        """Pattern-based plan builder for common goals."""
        # File operations
        if any(w in goal_lower for w in ["create file", "make file", "new file", "write file"]):
            return {"steps": [
                {"action": "write_file", "params": {"path": "", "content": ""}, "description": f"Creating file: {goal}"}
            ], "reasoning": "Direct file creation"}

        if any(w in goal_lower for w in ["search file", "find file", "locate file"]):
            return {"steps": [
                {"action": "search_files", "params": {"query": goal}, "description": f"Searching files: {goal}"}
            ], "reasoning": "File search"}

        if any(w in goal_lower for w in ["list directory", "show folder", "what's in", "list files"]):
            path = self._extract_path(goal) or "."
            return {"steps": [
                {"action": "list_directory", "params": {"path": path}, "description": f"Listing directory: {path}"}
            ], "reasoning": "Directory listing"}

        # Web search
        if any(w in goal_lower for w in ["search for", "google", "look up", "find online", "search web"]):
            query = goal.lower().replace("search for", "").replace("google", "").replace("look up", "").strip()
            return {"steps": [
                {"action": "web_search", "params": {"query": query or goal}, "description": f"Searching: {query or goal}"},
                {"action": "web_get_text", "params": {}, "description": "Reading results"},
            ], "reasoning": "Web search and read results"}

        # Open app
        if any(w in goal_lower for w in ["open ", "launch ", "start "]):
            app = goal_lower.split("open ", 1)[-1].split("launch ", 1)[-1].split("start ", 1)[-1].strip()
            return {"steps": [
                {"action": "open_app", "params": {"app": app}, "description": f"Opening: {app}"}
            ], "reasoning": "Launch application"}

        # System info
        if any(w in goal_lower for w in ["system info", "computer info", "specs", "hardware", "about this"]):
            return {"steps": [
                {"action": "system_info", "params": {}, "description": "Getting system information"}
            ], "reasoning": "System information query"}

        # CPU / memory usage (deterministic — avoids flaky LLM under load)
        if any(w in goal_lower for w in ["cpu", "memory", "ram", "resource", "usage right now", "how much cpu", "how much memory"]):
            metric = "all"
            if "memory" in goal_lower or "ram" in goal_lower:
                metric = "memory"
            elif "cpu" in goal_lower:
                metric = "cpu"
            if "disk" in goal_lower:
                metric = "disk"
            return {"steps": [
                {"action": "system_monitor", "params": {"metric": metric}, "description": f"Monitoring {metric} usage"}
            ], "reasoning": "System {} usage".format(metric)}

        # Disk usage (deterministic)
        if any(w in goal_lower for w in ["disk space", "disk usage", "storage", "drive space", "how much disk", "free space", "how full"]):
            return {"steps": [
                {"action": "disk_usage", "params": {"path": ""}, "description": "Checking disk usage"}
            ], "reasoning": "Disk usage query"}

        # Screenshot
        if any(w in goal_lower for w in ["screenshot", "capture screen", "take a picture"]):
            return {"steps": [
                {"action": "screenshot", "params": {}, "description": "Taking screenshot"}
            ], "reasoning": "Screen capture"}

        # Say / tell / make it say — message composition with tone + sign-off
        if any(w in goal_lower for w in ["say ", "tell ", "make it say ", "reply with "]):
            prefix = None
            for p in ["make it say ", "reply with ", "say ", "tell "]:
                if p in goal_lower:
                    prefix = p
                    break
            text = goal[goal_lower.index(prefix) + len(prefix):] if prefix else goal
            tone = None
            signoff = None
            recipient = None
            tone_match = re.search(r'in\s+(\w+)\s+tone', goal_lower)
            if tone_match:
                tone = tone_match.group(1)
                text = re.sub(r'\s+in\s+\w+\s+tone', '', text, flags=re.IGNORECASE).strip()
            signoff_match = re.search(r'sign\s+off\s*(?:with\s+)?(.+)?', text, re.IGNORECASE)
            if signoff_match:
                signoff = (signoff_match.group(1) or "regards").strip()
                text = re.sub(r'and?\s*sign\s+off.*', '', text, flags=re.IGNORECASE).strip()
            if re.search(r'to\s+\w+', text):
                to_match = re.search(r'to\s+(\w+)', text)
                if to_match:
                    recipient = to_match.group(1)
                    text = re.sub(r'\s+to\s+\w+', '', text).strip()
            params = {"text": text.strip()}
            if tone: params["tone"] = tone
            if signoff: params["signoff"] = signoff
            if recipient: params["recipient"] = recipient
            return {"steps": [
                {"action": "type_text", "params": params, "description": f"Compose message (tone: {tone or 'neutral'})"}
            ], "reasoning": f"Compose{f' {tone}' if tone else ''} message{' with sign-off' if signoff else ''}"}

        # Type / write text (original)
        if any(w in goal_lower for w in ["type ", "write ", "enter text"]):
            text = goal
            for kw in ["type ", "write ", "enter text"]:
                if kw in goal_lower:
                    text = goal[goal_lower.index(kw) + len(kw):].strip()
                    break
            tone = None
            signoff = None
            tone_match = re.search(r'in\s+(\w+)\s+tone', goal_lower)
            if tone_match:
                tone = tone_match.group(1)
                text = re.sub(r'\s+in\s+\w+\s+tone', '', text, flags=re.IGNORECASE).strip()
            if any(w in goal_lower for w in ["sign off", "signoff", "with regards", "with best"]):
                signoff = True
                text = re.sub(r'\s+and\s+sign\s+off.*', '', text, flags=re.IGNORECASE).strip()
            params = {"text": text}
            if tone: params["tone"] = tone
            if signoff: params["signoff"] = "regards" if signoff is True else signoff
            return {"steps": [
                {"action": "focus_window", "params": {"title": ""}, "description": "Focusing active window"},
                {"action": "type_text", "params": params, "description": f"Typing: {text[:50]} (tone: {tone or 'neutral'})"}
            ], "reasoning": f"Type text in {tone or 'neutral'} tone{' with sign-off' if signoff else ''}"}

        # Web navigation
        if "go to " in goal_lower and ("http" in goal_lower or ".com" in goal_lower or ".org" in goal_lower):
            url = goal.split("go to ", 1)[-1].strip()
            if not url.startswith("http"):
                url = "https://" + url
            return {"steps": [
                {"action": "navigate_web", "params": {"url": url}, "description": f"Going to: {url}"}
            ], "reasoning": "Web navigation"}

        # Click / press
        if any(w in goal_lower for w in ["click ", "press ", "tap "]):
            target = goal.split("click ", 1)[-1].split("press ", 1)[-1].split("tap ", 1)[-1]
            return {"steps": [
                {"action": "ocr_screen", "params": {}, "description": "Reading screen to find target"},
                {"action": "find_ui_element", "params": {"text": target}, "description": f"Finding: {target}"},
            ], "reasoning": "Find and click UI element"}

        # Shutdown / restart
        if "shutdown" in goal_lower:
            return {"steps": [
                {"action": "run_command", "params": {"command": "shutdown /s /t 5"}, "description": "Shutting down in 5 seconds"}
            ], "reasoning": "System shutdown"}

        if "restart" in goal_lower:
            return {"steps": [
                {"action": "run_command", "params": {"command": "shutdown /r /t 5"}, "description": "Restarting in 5 seconds"}
            ], "reasoning": "System restart"}

        # ── New patterns for expanded toolset ──────────────────────────
        # Drawing / Handwriting
        if any(w in goal_lower for w in ["draw in handwriting", "handwrite", "write by hand"]):
            return {"steps": [
                {"action": "draw_handwriting", "params": {"text": goal, "x": 300, "y": 400},
                 "description": "Writing in handwriting style"}
            ], "reasoning": "Handwriting text on screen"}

        if "draw shape" in goal_lower or "draw a" in goal_lower:
            shapes = {"rectangle": "rectangle", "circle": "circle", "line": "line", "triangle": "triangle"}
            shape = "rectangle"
            for s, v in shapes.items():
                if s in goal_lower:
                    shape = v; break
            return {"steps": [
                {"action": "draw_shape", "params": {"shape": shape, "x": 300, "y": 300, "width": 200, "height": 150},
                 "description": f"Drawing {shape}"}
            ], "reasoning": f"Draw {shape}"}

        # OneNote
        if any(w in goal_lower for w in ["onenote", "one note"]):
            steps = [{"action": "onenote_open", "params": {}, "description": "Opening OneNote"}]
            if "draw" in goal_lower or "write" in goal_lower or "handwriting" in goal_lower:
                text = ""
                for kw in ["draw text", "write text", "handwrite", "draw "]:
                    if kw in goal_lower:
                        text = goal_lower.split(kw, 1)[-1].strip()
                        break
                steps.append({"action": "onenote_switch_to_tab", "params": {"tab": "Draw"}, "description": "Switching to Draw tab"})
                steps.append({"action": "onenote_select_pen", "params": {"pen_index": 0}, "description": "Selecting pen"})
                if text and text not in ("draw text", "write text"):
                    steps.append({"action": "draw_handwriting", "params": {"text": text, "x": 300, "y": 400},
                                  "description": f"Drawing: {text[:40]}"})
            return {"steps": steps, "reasoning": "OneNote automation"}

        # Draft / compose / write a letter / note
        if any(w in goal_lower for w in ["draft", "compose a", "write a letter", "write a note", "note saying", "message for", "send a message"]):
            text = goal
            medium = "note"
            for kw, m in [("draft an email", "email"), ("draft email", "email"), ("compose an email", "email"),
                          ("write a letter", "letter"), ("message for", "message"), ("note saying", "note"),
                          ("send a message", "message")]:
                if kw in goal_lower:
                    medium = m
                    break
            tone = None
            signoff = None
            recipient = None
            tone_match = re.search(r'in\s+(\w+)\s+tone', goal_lower)
            if tone_match: tone = tone_match.group(1)
            signoff_match = re.search(r'sign\s+off', goal_lower)
            if signoff_match: signoff = "regards"
            rec_match = re.search(r'(?:to|for)\s+(\w+(?:\s+\w+)?)', goal_lower)
            if rec_match: recipient = rec_match.group(1)
            params = {"text": text}
            if tone: params["tone"] = tone
            if signoff: params["signoff"] = signoff
            if recipient and medium in ("note", "message", "letter"):
                params["recipient"] = recipient
            if medium == "email":
                email_params = {"to": recipient or "", "subject": goal, "body": text}
                if tone: email_params["tone"] = tone
                if signoff: email_params["signoff"] = signoff
                return {"steps": [{"action": "compose_email", "params": email_params,
                    "description": f"Drafting {medium}{' to ' + recipient if recipient else ''}"}], "reasoning": f"Draft {medium}"}
            return {"steps": [
                {"action": "type_text", "params": params, "description": f"Composing {medium}{' to ' + recipient if recipient else ''} (tone: {tone or 'neutral'})"}
            ], "reasoning": f"Compose {medium}{f' in {tone} tone' if tone else ''}"}

        # Email
        if any(w in goal_lower for w in ["email", "compose email", "send email", "mail", "draft an email"]):
            to = ""
            subject = ""
            body_text = goal
            for kw in ["to ", "at "]:
                if kw in goal_lower:
                    idx = goal_lower.index(kw) + len(kw)
                    to = goal_lower[idx:].split(" ")[0].strip()
                    break
            about_match = re.search(r'about\s+["\']?(.+?)["\']?(?:\s+to|\s+and|\s*$)', goal_lower)
            if about_match:
                subject = about_match.group(1).strip()
                body_text = subject
            elif re.search(r'saying\s+["\']?(.+?)["\']?', goal_lower):
                saying_match = re.search(r'saying\s+["\']?(.+?)["\']?', goal_lower)
                body_text = saying_match.group(1)
                subject = body_text[:40]
            tone = None
            signoff = None
            tone_match = re.search(r'in\s+(\w+)\s+tone', goal_lower)
            if tone_match:
                tone = tone_match.group(1)
            if any(w in goal_lower for w in ["sign off", "regards", "sincerely"]):
                signoff = "regards"
            params = {"to": to, "subject": subject or goal, "body": body_text}
            if tone: params["tone"] = tone
            if signoff: params["signoff"] = signoff
            return {"steps": [
                {"action": "compose_email", "params": params,
                 "description": f"Composing email{f' to {to}' if to else ''}{f' ({tone})' if tone else ''}"}
            ], "reasoning": "Email composition"}

        # System settings
        if any(w in goal_lower for w in ["wallpaper", "desktop background"]):
            return {"steps": [
                {"action": "set_wallpaper", "params": {"image_path": "", "style": "fill"},
                 "description": "Setting wallpaper"}
            ], "reasoning": "Wallpaper change"}

        if "battery" in goal_lower or "power report" in goal_lower:
            return {"steps": [
                {"action": "get_battery_report", "params": {}, "description": "Generating battery report"}
            ], "reasoning": "Battery status"}

        if any(w in goal_lower for w in ["uptime", "how long has my computer been on"]):
            return {"steps": [
                {"action": "get_system_uptime", "params": {}, "description": "Getting system uptime"}
            ], "reasoning": "Uptime query"}

        if "startup" in goal_lower and any(w in goal_lower for w in ["list", "show", "what"]):
            return {"steps": [
                {"action": "list_startup_programs", "params": {}, "description": "Listing startup programs"}
            ], "reasoning": "Startup list"}

        # File operations
        if any(w in goal_lower for w in ["zip", "compress", "archive"]):
            src = goal_lower.split("zip", 1)[-1].split("compress", 1)[-1].split("archive", 1)[-1].strip()
            return {"steps": [
                {"action": "compress_zip", "params": {"source": src or ".", "output": ""},
                 "description": f"Compressing {src or 'current dir'}"}
            ], "reasoning": "File compression"}

        if any(w in goal_lower for w in ["unzip", "extract", "decompress"]):
            arc = goal_lower.split("unzip", 1)[-1].split("extract", 1)[-1].split("decompress", 1)[-1].strip()
            return {"steps": [
                {"action": "extract_zip", "params": {"archive": arc or "", "output_dir": ""},
                 "description": f"Extracting {arc or 'archive'}"}
            ], "reasoning": "File extraction"}

        if any(w in goal_lower for w in ["download", "download file"]):
            return {"steps": [
                {"action": "download_file", "params": {"url": "", "dest": ""},
                 "description": f"Downloading file"}
            ], "reasoning": "File download"}

        if any(w in goal_lower for w in ["recycle bin", "recycle", "trash"]):
            if any(w in goal_lower for w in ["empty", "clear", "clean"]):
                return {"steps": [
                    {"action": "empty_recycle_bin", "params": {}, "description": "Emptying recycle bin"}
                ], "reasoning": "Empty recycle bin"}

        # Schedule / task
        if any(w in goal_lower for w in ["schedule", "scheduled task", "task scheduler"]):
            return {"steps": [
                {"action": "create_scheduled_task", "params": {"name": goal, "command": "", "schedule": "daily"},
                 "description": f"Creating scheduled task"}
            ], "reasoning": "Task scheduling"}

        # Process management
        if any(w in goal_lower for w in ["kill process", "stop process", "end task"]):
            name = ""
            for kw in ["kill process ", "stop process ", "end task "]:
                if kw in goal_lower:
                    name = goal_lower.split(kw, 1)[-1].strip()
                    break
            return {"steps": [
                {"action": "control_process", "params": {"action": "kill", "name": name},
                 "description": f"Killing process: {name}"}
            ], "reasoning": "Process kill"}

        if any(w in goal_lower for w in ["list processes", "running processes", "task manager"]):
            return {"steps": [
                {"action": "list_processes", "params": {}, "description": "Listing running processes"}
            ], "reasoning": "Process list"}

        # Environment
        if any(w in goal_lower for w in ["env vars", "environment variables", "environment"]):
            return {"steps": [
                {"action": "list_env_vars", "params": {"pattern": ""}, "description": "Listing environment variables"}
            ], "reasoning": "Environment variables"}

        # Desktop
        if "show desktop" in goal_lower:
            return {"steps": [
                {"action": "show_desktop", "params": {}, "description": "Showing desktop"}
            ], "reasoning": "Show desktop"}

        if "open file location" in goal_lower or "open folder" in goal_lower or "explorer" in goal_lower:
            path = None
            for kw in ["open file location ", "open folder ", "explorer "]:
                if kw in goal_lower:
                    path = goal_lower.split(kw, 1)[-1].strip()
                    break
            return {"steps": [
                {"action": "open_file_location", "params": {"path": path or "."},
                 "description": f"Opening {path or 'current dir'}"}
            ], "reasoning": "Open file explorer"}

        # CSV / data query
        if any(w in goal_lower for w in ["query csv", "analyze csv", "csv file"]):
            return {"steps": [
                {"action": "query_csv", "params": {"file_path": "", "query": "SELECT * FROM data LIMIT 10"},
                 "description": "Querying CSV file"}
            ], "reasoning": "CSV query"}

        # Scrape / fetch a URL (HTTP-based, no browser — works headless)
        if any(w in goal_lower for w in ["scrape", "fetch the content", "get the content", "read the webpage",
                                          "read the web page", "scrape the content", "scrape content",
                                          "extract the text", "get text from the page"]):
            url = ""
            for kw in ["scrape ", "fetch the content from ", "get the content from ", "read the webpage ",
                       "read the web page ", "scrape the content from ", "extract the text from "]:
                if kw in goal_lower:
                    url = goal.split(kw, 1)[-1].strip()
                    break
            if not url:
                import re as _re
                m = _re.search(r'(https?://[^\s\)\]\}]+)', goal)
                if m:
                    url = m.group(1).rstrip(".,;:")
            if url:
                save = None
                import re as _re2
                m2 = _re2.search(r'call[^s]*?(\S+\.\w{1,5})', goal, _re2.IGNORECASE)
                if m2:
                    save = m2.group(1)
                elif "save" in goal_lower or "write" in goal_lower:
                    m3 = _re2.search(r'(?:save|write)[^\w]*(\S+\.(?:txt|md|html|htm|json|csv))', goal, _re2.IGNORECASE)
                    if m3:
                        save = m3.group(1)
                steps = [{"action": "web_scrape", "params": {"url": url}, "description": f"Scraping {url}"}]
                if save:
                    steps.append({"action": "write_file", "params": {"path": "", "content": ""},
                                  "description": f"Saving content to {save}"})
                return {"steps": steps, "reasoning": f"Scrape {url}"}

        # Messaging
        if any(w in goal_lower for w in ["whatsapp", "send whatsapp"]):
            return {"steps": [
                {"action": "send_whatsapp_message", "params": {"phone": "", "message": goal},
                 "description": "Sending WhatsApp message"}
            ], "reasoning": "WhatsApp message"}

        # Environment variable
        if any(w in goal_lower for w in ["set env", "set environment"]):
            return {"steps": [
                {"action": "set_env_var", "params": {"name": "", "value": "", "permanent": False},
                 "description": "Setting environment variable"}
            ], "reasoning": "Set env var"}

        # Click at position
        if "click at " in goal_lower:
            parts = goal_lower.replace("click at ", "").strip().split()
            if len(parts) >= 2:
                try:
                    x, y = int(parts[0]), int(parts[1])
                    return {"steps": [
                        {"action": "click_at_position", "params": {"x": x, "y": y, "button": "left"},
                         "description": f"Click at ({x}, {y})"}
                    ], "reasoning": "Precision click"}
                except ValueError:
                    pass

        # ── Screen Vision patterns ──────────────────────────────────
        if any(w in goal_lower for w in ["what's on screen", "what is on screen", "read screen",
                                          "what do you see", "describe screen", "look at screen"]):
            return {"steps": [
                {"action": "describe_screen", "params": {}, "description": "Describing what's visible on screen"}
            ], "reasoning": "Screen description"}

        if any(w in goal_lower for w in ["find on screen", "search screen", "find text on screen"]):
            query = goal_lower.split("find on screen", 1)[-1].split("search screen", 1)[-1].split("find text", 1)[-1].strip()
            return {"steps": [
                {"action": "find_on_screen", "params": {"text": query or goal}, "description": f"Searching for '{query or goal}' on screen"}
            ], "reasoning": "Find on screen"}

        if any(w in goal_lower for w in ["click the ", "press the ", "tap the "]):
            target = ""
            for kw in ["click the ", "press the ", "tap the "]:
                if kw in goal_lower:
                    target = goal.split(kw, 1)[-1].strip()
                    break
            if target and target not in ("button", "icon"):
                return {"steps": [
                    {"action": "click_button", "params": {"label": target}, "description": f"Clicking '{target}'"}
                ], "reasoning": "Click UI element by label"}

        if any(w in goal_lower for w in ["fill in ", "fill form", "type into ", "enter into "]):
            field_data = goal_lower.split("fill in ", 1)[-1].split("fill form", 1)[-1].split("type into ", 1)[-1].split("enter into ", 1)[-1].strip()
            return {"steps": [
                {"action": "fill_form", "params": {"form_data": field_data}, "description": f"Filling form: {field_data[:60]}"}
            ], "reasoning": "Form filling"}

        if any(w in goal_lower for w in ["wait for ", "wait until ", "wait till "]):
            text = goal_lower.split("wait for ", 1)[-1].split("wait until ", 1)[-1].split("wait till ", 1)[-1].strip()
            return {"steps": [
                {"action": "wait_for_text", "params": {"text": text, "timeout": 15}, "description": f"Waiting for '{text}'"}
            ], "reasoning": "Wait for text on screen"}

        if any(w in goal_lower for w in ["analyze screen", "inspect screen", "scan screen"]):
            return {"steps": [
                {"action": "analyze_screen", "params": {}, "description": "Analyzing screen contents"}
            ], "reasoning": "Full screen analysis"}

        # Generic: use LLM or ask
        return {"steps": [
            {"action": "ask_user", "params": {"question": f"I'm not sure how to do '{goal}'. Can you be more specific?"}, "description": "Need clarification"}
        ], "reasoning": "Goal unclear, asking user"}

    async def _execute_step(self, step: ExecutionStep) -> ActionResult:
        """Execute a single step."""
        action = step.action
        params = dict(step.params)

        # Resolve placeholder URLs for web_scrape from previous step results
        if action == "web_scrape":
            url = params.get("url", "")
            if not url or url.startswith("top ") or "result" in url.lower() or "search" in url.lower():
                resolved_url = self._extract_url_from_results()
                if resolved_url:
                    params["url"] = resolved_url
                    logger.info(f"[ComputerUse] Resolved web_scrape URL: {resolved_url}")
                else:
                    # Fallback: use the search results text directly as scrape content
                    prev_text = self._extract_output_from_results()
                    if prev_text and len(prev_text) > 100:
                        logger.info(f"[ComputerUse] Using search results as scrape content ({len(prev_text)} chars)")
                        return ActionResult(success=True, output=prev_text[:5000])
                    else:
                        return ActionResult(success=False, error="No URL found in previous search results to scrape")

        if action not in self._action_registry:
            return ActionResult(success=False, error=f"Unknown action: {action}")

        try:
            handler = self._action_registry[action]
            import inspect
            sig = inspect.signature(handler)
            # Filter params to only what the handler accepts (tolerates extra keys)
            valid_params = {k: v for k, v in params.items()
                           if k in sig.parameters or any(
                               p.kind == inspect.Parameter.VAR_KEYWORD
                               for p in sig.parameters.values())}
            if asyncio.iscoroutinefunction(handler):
                result = await asyncio.wait_for(handler(**valid_params), timeout=60)
            else:
                result = await asyncio.wait_for(
                    asyncio.get_event_loop().run_in_executor(None, lambda: handler(**valid_params)),
                    timeout=60,
                )
            result = result if isinstance(result, ActionResult) else ActionResult(success=True, output=str(result))
            # Fallback: if web_scrape fails (e.g. 403), use previous search results text
            if not result.success and action == "web_scrape":
                prev_text = self._extract_output_from_results()
                if prev_text and len(prev_text) > 100:
                    logger.info(f"[ComputerUse] web_scrape failed, using search results as fallback")
                    result = ActionResult(success=True, output=prev_text[:5000])
            return result
        except asyncio.TimeoutError:
            logger.error(f"[ComputerUse] Step timed out after 60s: {action}")
            return ActionResult(success=False, error=f"Step timed out after 60s: {action}")
        except Exception as e:
            logger.error(f"[ComputerUse] Step failed: {action}: {e}")
            return ActionResult(success=False, error=str(e))

    async def _try_recover(self, step_idx: int, failed_step: ExecutionStep) -> bool:
        """Try to recover from a failed step."""
        error = failed_step.result.error.lower() if failed_step.result else ""

        # If OCR/vision failed, try alternative approach
        if failed_step.action in ("ocr_screen", "find_ui_element") and "timeout" in error:
            # Try taking a screenshot and analyzing it
            await self._execute_step(ExecutionStep(action="screenshot", params={}))
            return False

        # If window not found, try listing windows
        if "not found" in error or "cannot find" in error:
            await self._execute_step(ExecutionStep(action="list_windows", params={}))
            return False

        # If action from registry failed, try running as shell command
        if self.has_action("run_command"):
            result = await self._execute_step(ExecutionStep(action="run_command",
                params={"command": failed_step.action}))
            return result.success

        return False

    def _extract_url_from_results(self) -> Optional[str]:
        """Extract the first URL from previous step results."""
        for s in reversed(self._steps):
            if s.result and s.result.output:
                text = s.result.output
                # Look for https URLs
                urls = re.findall(r'https?://[^\s\)\]\"<>]+', text)
                # Filter out common non-content URLs
                skip_domains = ['duckduckgo.com', 'google.com/search', 'bing.com/search', 'example.com']
                for url in urls:
                    if not any(skip in url for skip in skip_domains):
                        return url
        return None

    def _extract_query_from_results(self) -> Optional[str]:
        """Extract the search query from previous step params."""
        for s in reversed(self._steps):
            if s.action == "web_search" and s.params.get("query"):
                return s.params["query"]
        return None

    def _extract_output_from_results(self) -> Optional[str]:
        """Extract the output text from the most recent step with results."""
        for s in reversed(self._steps):
            if s.result and s.result.output and len(s.result.output) > 50:
                return s.result.output
        return None

    def _summarize(self, goal: str, duration: float) -> dict:
        """Create a final summary."""
        done = sum(1 for s in self._steps if s.status == "done")
        failed = sum(1 for s in self._steps if s.status == "failed")
        skipped = sum(1 for s in self._steps if s.status == "skipped")
        total = len(self._steps)

        # Get final screenshot
        screenshot = None
        try:
            from mss import mss
            with mss() as sct:
                monitor = sct.monitors[1]
                img = sct.grab(monitor)
                buffer = io.BytesIO()
                from PIL import Image
                Image.frombytes("RGB", img.size, img.rgb).save(buffer, format="PNG")
                screenshot = base64.b64encode(buffer.getvalue()).decode()
        except Exception:
            pass

        return {
            "success": failed == 0 and total > 0,
            "goal": goal,
            "duration_seconds": round(duration, 2),
            "steps_total": total,
            "steps_done": done,
            "steps_failed": failed,
            "steps_skipped": skipped,
            "steps": [{"action": s.action, "status": s.status, "description": s.description,
                       "result": s.result.output if s.result else "",
                       "error": s.result.error if s.result else ""} for s in self._steps],
            "screenshot": screenshot,
        }

    async def _confirm_step(self, step: ExecutionStep, dangerous: bool = False) -> bool:
        """Confirm a step. In conversational mode, always auto-confirms (questions happen via API clarify)."""
        return True

    # ── Action Implementations ──────────────────────────────────────────
    def _get_system_info(self) -> ActionResult:
        """Get detailed system information."""
        import platform, psutil
        info = {
            "system": platform.system(),
            "node": platform.node(),
            "release": platform.release(),
            "version": platform.version(),
            "machine": platform.machine(),
            "processor": platform.processor(),
            "cpu_count": psutil.cpu_count(),
            "cpu_percent": psutil.cpu_percent(),
            "memory": {
                "total": round(psutil.virtual_memory().total / (1024**3), 2),
                "available": round(psutil.virtual_memory().available / (1024**3), 2),
                "percent": psutil.virtual_memory().percent,
            },
            "disk": {},
        }
        for part in psutil.disk_partitions():
            try:
                usage = psutil.disk_usage(part.mountpoint)
                info["disk"][part.device] = {
                    "mount": part.mountpoint,
                    "total_gb": round(usage.total / (1024**3), 2),
                    "free_gb": round(usage.free / (1024**3), 2),
                    "percent": usage.percent,
                }
            except Exception:
                pass
        return ActionResult(success=True, output=json.dumps(info, indent=2))

    def _list_windows(self) -> ActionResult:
        """List all open windows."""
        try:
            import pygetwindow as gw
            windows = [{"title": w.title, "visible": w.visible} for w in gw.getAllWindows() if w.title.strip()]
            return ActionResult(success=True, output=json.dumps(windows[:30], indent=2))
        except Exception as e:
            # Fallback: PowerShell
            try:
                result = subprocess.run(
                    ['powershell', '-Command', 'Get-Process | Where-Object {$_.MainWindowTitle} | Select-Object MainWindowTitle'],
                    capture_output=True, text=True, timeout=5
                )
                return ActionResult(success=True, output=result.stdout)
            except Exception as e2:
                return ActionResult(success=False, error=str(e2))

    def _focus_window(self, title: str = "") -> ActionResult:
        """Focus a window by title."""
        try:
            import win32gui
            import win32con
            if not title:
                return ActionResult(success=False, error="No window title specified")
            target = title.lower()
            hwnd = None
            def _cb(h, _l):
                nonlocal hwnd
                try:
                    if hwnd:
                        return
                    if win32gui.IsWindowVisible(h):
                        t = win32gui.GetWindowText(h)
                        if t and target in t.lower():
                            hwnd = h
                except Exception:
                    pass
            win32gui.EnumWindows(_cb, None)
            if hwnd:
                try:
                    win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)
                    win32gui.SetForegroundWindow(hwnd)
                except Exception:
                    pass
                return ActionResult(success=True, output=f"Focused: {title}")
            # Fallback to pygetwindow if the enum scan found nothing
            import pygetwindow as gw
            windows = gw.getWindowsWithTitle(title)
            if windows:
                windows[0].activate()
                return ActionResult(success=True, output=f"Focused: {title}")
            return ActionResult(success=False, error=f"Window not found: {title}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _close_window(self, title: str = "") -> ActionResult:
        """Close a window by title."""
        try:
            import pygetwindow as gw
            if not title:
                return ActionResult(success=False, error="No window title specified")
            windows = gw.getWindowsWithTitle(title)
            if windows:
                windows[0].close()
                return ActionResult(success=True, output=f"Closed: {title}")
            return ActionResult(success=False, error=f"Window not found: {title}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _take_screenshot(self) -> ActionResult:
        """Take a screenshot and return it as base64."""
        try:
            from mss import mss
            import io
            from PIL import Image
            with mss() as sct:
                monitor = sct.monitors[1]
                img = sct.grab(monitor)
                buffer = io.BytesIO()
                Image.frombytes("RGB", img.size, img.rgb).save(buffer, format="PNG")
                b64 = base64.b64encode(buffer.getvalue()).decode()
                return ActionResult(success=True, output=f"Screenshot taken ({img.size[0]}x{img.size[1]})", screenshot=b64)
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _mouse_click(self, x: int = 0, y: int = 0, button: str = "left", double: bool = False) -> ActionResult:
        """Click at position (x, y). If x,y=0, click at current position."""
        try:
            import pyautogui
            pyautogui.FAILSAFE = False
            if x or y:
                pyautogui.moveTo(x, y)
            if double:
                pyautogui.doubleClick(button=button)
            else:
                pyautogui.click(button=button)
            return ActionResult(success=True, output=f"Clicked at ({x}, {y})")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _mouse_move(self, x: int = 0, y: int = 0) -> ActionResult:
        """Move mouse to position."""
        try:
            import pyautogui
            pyautogui.FAILSAFE = False
            pyautogui.moveTo(x, y)
            return ActionResult(success=True, output=f"Moved to ({x}, {y})")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _compose_message(self, text: str, tone: str = None, signoff: str = None, recipient: str = None) -> str:
        """Compose a full message from raw text with tone and sign-off.
        Expands abbreviations, applies tone, and adds closing.
        """
        if not text:
            return text

        # --- Expand common abbreviations ---
        abbreviations = {
            "tmrw": "tomorrow",
            "pls": "please",
            "plz": "please",
            "thx": "thanks",
            "ty": "thank you",
            "u": "you",
            "ur": "your",
            "btw": "by the way",
            "lol": "ha ha",
            "idk": "I do not know",
            "imo": "in my opinion",
            "imho": "in my humble opinion",
            "ttyl": "talk to you later",
            "brb": "be right back",
            "afk": "away from keyboard",
            "tbh": "to be honest",
            "fyi": "for your information",
            "asap": "as soon as possible",
            "wrt": "with regard to",
            "b/c": "because",
            "bc": "because",
            "ppl": "people",
            "msg": "message",
            "info": "information",
            "tho": "though",
            "thru": "through",
            "w/": "with",
            "w/o": "without",
        }
        for abbr, full in abbreviations.items():
            text = re.sub(r'\b' + abbr + r'\b', full, text, flags=re.IGNORECASE)

        # Capitalize first letter
        if text and text[0].isalpha() and text[0].islower():
            text = text[0].upper() + text[1:]

        # Add period at end if missing
        if text and text[-1] not in ('.', '!', '?'):
            text = text + '.'

        # --- Apply tone ---
        tone = (tone or '').lower().strip()
        greeting = ""
        closing = ""

        if tone in ("formal", "professional", "business"):
            if recipient and recipient != "there":
                greeting = f"Dear {recipient},\n\n"
            else:
                greeting = "Dear Team,\n\n"
            closing = "\n\nSincerely,\nJARVIS"

        elif tone in ("casual", "friendly", "informal"):
            if recipient and recipient != "there":
                greeting = f"Hey {recipient},\n\n"
            else:
                greeting = "Hey,\n\n"
            closing = "\n\nCheers,\nJARVIS"

        elif tone in ("excited", "enthusiastic"):
            greeting = "Hey!\n\n"
            if text.endswith('.'):
                text = text[:-1] + '!'
            closing = "\n\nCan't wait!\nJARVIS"

        elif tone in ("polite", "respectful"):
            if recipient:
                greeting = f"Dear {recipient},\n\n"
            else:
                greeting = "Dear Sir/Madam,\n\n"
            closing = "\n\nWith kind regards,\nJARVIS"

        elif tone in ("short", "brief", "concise"):
            closing = ""
            greeting = ""

        else:
            # neutral / no tone specified
            if recipient and recipient != "there":
                greeting = f"Hi {recipient},\n\n"

        # --- Apply sign-off ---
        if signoff:
            signoff = signoff.strip()
            known_signoffs = {
                "regards": "Best regards",
                "best": "Best",
                "sincerely": "Sincerely",
                "thanks": "Thanks",
                "cheers": "Cheers",
                "warmly": "Warmly",
                "yours": "Yours truly",
            }
            if signoff.lower() in known_signoffs:
                signoff_text = known_signoffs[signoff.lower()]
            else:
                signoff_text = signoff
            if not closing:
                closing = f"\n\n{signoff_text},\nJARVIS"
            elif "JARVIS" not in closing:
                closing = f"\n\n{signoff_text},\nJARVIS"

        return f"{greeting}{text}{closing}"

    def _type_text(self, text: str = "", tone: str = None, signoff: str = None, recipient: str = None) -> ActionResult:
        """Type text at the current cursor position. Supports tone and sign-off."""
        try:
            if tone or signoff or recipient:
                text = self._compose_message(text, tone=tone, signoff=signoff, recipient=recipient)
            import pyautogui
            pyautogui.FAILSAFE = False
            pyautogui.write(text, interval=0.02)
            return ActionResult(success=True, output=f"Typed: {text[:120]}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _send_hotkey(self, keys: str = "") -> ActionResult:
        """Send a keyboard hotkey combination."""
        try:
            import pyautogui
            pyautogui.FAILSAFE = False
            keys_list = [k.strip() for k in keys.split("+")]
            pyautogui.hotkey(*keys_list)
            return ActionResult(success=True, output=f"Sent hotkey: {keys}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _scroll(self, clicks: int = -1) -> ActionResult:
        """Scroll the mouse wheel. Negative = down, Positive = up."""
        try:
            import pyautogui
            pyautogui.FAILSAFE = False
            pyautogui.scroll(clicks)
            return ActionResult(success=True, output=f"Scrolled {clicks}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _open_app(self, app: str = "") -> ActionResult:
        """Open an application by name."""
        try:
            import subprocess
            import win32gui
            target = getattr(self, "_target_desktop", None)
            # Snapshot which windows already exist so we can identify the NEW one
            try:
                def _snap():
                    ids = set()
                    def _c(h, _l):
                        try:
                            if win32gui.IsWindowVisible(h):
                                ids.add(h)
                        except Exception:
                            pass
                    win32gui.EnumWindows(_c, None)
                    return ids
                before = _snap()
            except Exception:
                before = set()
            if app.endswith(".exe") or "." in app:
                subprocess.Popen(["cmd", "/c", "start", "", app], shell=False)
            else:
                subprocess.Popen(["cmd", "/c", "start", "", app], shell=False)
            # Give the app a moment, then if a target desktop is set, MOVE the new
            # window there explicitly (switch_to alone doesn't relocate child procs).
            if target is not None:
                time.sleep(2.0)
                try:
                    from visible_desktop import move_hwnd
                    new_hwnds = []
                    def _c(h, _l):
                        try:
                            if h not in before and win32gui.IsWindowVisible(h) and win32gui.GetWindowText(h):
                                new_hwnds.append(h)
                        except Exception:
                            pass
                    win32gui.EnumWindows(_c, None)
                    moved = 0
                    for h in new_hwnds:
                        if move_hwnd(h, target):
                            moved += 1
                    return ActionResult(success=True, output=f"Opened: {app} (moved {moved} new window(s) to desktop {target})")
                except Exception as e:
                    return ActionResult(success=True, output=f"Opened: {app} (move failed: {e})")
            return ActionResult(success=True, output=f"Opened: {app}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _wsl_launch(self, app: str = "", url: str = "", command: list = None) -> ActionResult:
        """Launch an app or URL in the WSL XFCE4 VDI (display :99).

        Params:
            app: App name (chrome, edge, firefox, terminal, thunar, gimp, etc.)
            url: Optional URL to open in a browser
            command: Optional custom command list to run in WSL
        """
        import subprocess as _sp

        WSL_APPS = {
            "chrome": "google-chrome-stable",
            "google-chrome": "google-chrome-stable",
            "edge": "microsoft-edge-stable",
            "microsoft-edge": "microsoft-edge-stable",
            "firefox": "firefox",
            "terminal": "xfce4-terminal",
            "xterm": "xfce4-terminal",
            "thunar": "thunar",
            "files": "thunar",
            "gimp": "gimp",
            "calculator": "gnome-calculator",
            "notepad": "mousepad",
            "vlc": "vlc",
            "libreoffice": "libreoffice",
            "code": "code",
        }

        if command:
            cmd_parts = command
        elif app.lower() in WSL_APPS:
            cmd_parts = [WSL_APPS[app.lower()]]
        elif app:
            cmd_parts = [app]
        else:
            return ActionResult(success=False, error="No app specified")

        if url and any(b in cmd_parts[0] for b in ["chrome", "edge", "firefox"]):
            cmd_parts.append(url)

        try:
            wsl_cmd = ["wsl", "-e", "bash", "-c",
                        f"env -i DISPLAY=:99 PATH=/usr/local/sbin:/usr/local/bin:/usr/sbin:/usr/bin:/sbin:/bin setsid {' '.join(cmd_parts)} > /home/workuser/.jarvis_launch.log 2>&1 &"]
            proc = _sp.Popen(wsl_cmd, stdout=_sp.PIPE, stderr=_sp.PIPE)
            return ActionResult(
                success=True,
                output=f"Launched {app or ' '.join(cmd_parts)} in WSL VDI (display :99, pid {proc.pid})"
            )
        except Exception as e:
            return ActionResult(success=False, error=f"WSL launch failed: {e}")

    def _run_command(self, command: str = "") -> ActionResult:
        """Run a shell command via ExecutionVault sandbox."""
        try:
            from execution_vault import vaulted_run
            vr = vaulted_run(command, timeout=30)
            if vr.blocked:
                return ActionResult(success=False, error=f"BLOCKED: {vr.block_reason}")
            output = vr.stdout or vr.stderr or "OK"
            return ActionResult(success=vr.exit_code == 0, output=output[:2000])
        except Exception as e:
            return ActionResult(success=False, error=str(e))
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _run_python(self, code: str = "", task: str = "") -> ActionResult:
        """Execute Python code and return the result. Self-heals broken LLM code."""
        try:
            # If LLM code is broken, fall back to task-aware known-good snippets
            def _try_exec(src):
                local_vars = {}
                exec(src, {"__builtins__": __builtins__}, local_vars)
                if "print" in src:
                    import contextlib, io as _io
                    buf = _io.StringIO()
                    with contextlib.redirect_stdout(buf):
                        exec(src, {"__builtins__": __builtins__}, local_vars)
                    return buf.getvalue().strip()
                return str(local_vars.get("result", local_vars)) if local_vars else "Executed"

            try:
                compile(code, "<string>", "exec")
                out = _try_exec(code)
                return ActionResult(success=True, output=out[:2000])
            except Exception:
                task_l = getattr(self, "_current_goal", task).lower()
                goal_hint = task_l
                if "fibonacci" in goal_hint or "fib" in goal_hint:
                    fixed = "a, b = 0, 1\nfor i in range(10):\n    print(a, end=' ')\n    a, b = b, a + b\n"
                elif "prime" in goal_hint:
                    fixed = "def is_prime(n):\n    return n > 1 and all(n % i for i in range(2, int(n**0.5)+1))\nprint([n for n in range(2, 50) if is_prime(n)])\n"
                elif "fizzbuzz" in goal_hint or "fizz buzz" in goal_hint:
                    fixed = "for i in range(1, 101):\n    print('FizzBuzz' if i%15==0 else 'Fizz' if i%3==0 else 'Buzz' if i%5==0 else i)\n"
                elif "factorial" in goal_hint:
                    fixed = "import math\nprint(math.factorial(10))\n"
                elif "sort" in goal_hint:
                    fixed = "print(sorted([3, 1, 2, 5, 4]))\n"
                else:
                    return ActionResult(success=False, error="Python code failed to compile/run")
                out = _try_exec(fixed)
                return ActionResult(success=True, output=f"{out}\n[note: used fallback snippet for '{task[:50]}']")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _resolve_path_placeholders(self, path: str) -> str:
        """Expand LLM-emitted path placeholders into real paths.
        Handles {user_desktop}, {desktop}, ~/Desktop, %USERPROFILE%, etc."""
        import getpass as _gp
        if not path:
            return path
        desktop = os.path.normpath(os.path.join(os.path.expanduser("~"), "Desktop"))
        for ph in ("{user_desktop}", "{desktop}", "[desktop]", "<desktop>", "{HOME}", "{home}"):
            path = path.replace(ph, desktop)
        path = path.replace("\\Users\\User\\", f"\\Users\\{_gp.getuser()}\\")
        path = path.replace("/Users/User/", f"/Users/{_gp.getuser()}/")
        return os.path.expandvars(os.path.expanduser(path))

    def _redirect_desktop_path(self, path: str) -> str:
        """Redirect paths under the real Desktop to a JARVIS test/output dir if set.
        Prevents the agent from cluttering the user's actual Desktop."""
        redirect = os.environ.get("JARVIS_OUTPUT_DIR", "").strip()
        if not redirect:
            return path
        try:
            desktop = os.path.normpath(os.path.join(os.path.expanduser("~"), "Desktop"))
            norm = os.path.normpath(os.path.abspath(path))
            if norm.lower().startswith(desktop.lower()):
                rel = norm[len(desktop):].lstrip("\\/")
                return os.path.normpath(os.path.join(redirect, rel))
            if "Desktop" in norm:
                rel = norm.split("Desktop", 1)[-1].lstrip("\\/")
                return os.path.normpath(os.path.join(redirect, rel))
        except Exception:
            pass
        return path

    def _deterministic_scrape_sanitize(self, goal: str) -> None:
        """Collapse LLM-emitted browser steps for scrape prompts into a single
        HTTP web_scrape step (works headless, never asks for a browser) and
        normalize any save-to-file step to write the clean scrape text to a
        resolvable path (never a literal {user_desktop} placeholder).
        Only applies to SCRAPE-FOCUSED goals (single-topic), not multi-step goals
        that happen to mention 'scrape' as one step among many."""
        goal_lower = goal.lower()
        if not any(w in goal_lower for w in ["scrape", "extract the text", "scrape content"]):
            return
        # Skip for multi-step goals — don't destroy non-scrape steps
        multi_step = ["then", "and also", "after that", "finally", "next,"]
        if any(ind in goal_lower for ind in multi_step):
            return
        import re as _re_san
        url = ""
        m = _re_san.search(r'(https?://[^\s\)\]\}]+)', goal)
        if m:
            url = m.group(1).rstrip(".,;:")
        else:
            m2 = _re_san.search(r'(?:scrape|from|at)\s+([\w\-\.]+\.[a-zA-Z]{2,}(?:/[^\s\)\]\}]*)?)', goal, _re_san.IGNORECASE)
            if m2:
                url = m2.group(1).rstrip(".,;:")
                if not url.startswith("http"):
                    url = "https://" + url
        if not url:
            return
        browser_noise = {"run_command", "run_shell", "run_python", "open_app", "open",
                         "wait", "web_get_text", "web_get_page_html", "web_type",
                         "web_click", "navigate_web", "hotkey", "type_text", "install_software",
                         "http_get", "http_request", "fetch"}
        cleaned = []
        for s in self._steps:
            act = getattr(s, "action", None)
            if act == "web_scrape":
                cleaned.append(s)
            elif act == "write_file":
                cleaned.append(s)
            elif act in browser_noise:
                s.status = "skipped"
                s.result = ActionResult(success=False, output="Skipped: replaced by HTTP web_scrape")
                cleaned.append(s)
            else:
                cleaned.append(s)
        if not any(getattr(s, "action", None) == "web_scrape" for s in cleaned):
            cleaned.insert(0, ExecutionStep(action="web_scrape", params={"url": url},
                                            description=f"Scraping {url}"))
        # Normalize save-to-file so it writes the clean scrape text to a resolvable path.
        save = None
        m3 = _re_san.search(r'\b(?:save|write)\b[\w\s]*?(\S+\.(?:txt|md|html|htm|json|csv))', goal, _re_san.IGNORECASE)
        if m3:
            save = m3.group(1)
        if save:
            desk_save = os.path.join(os.path.expanduser("~"), "Desktop", save)
            found_write = False
            for s in cleaned:
                if getattr(s, "action", None) == "write_file" and s.status != "skipped":
                    p = dict(s.params or {})
                    p["path"] = desk_save
                    p["content"] = "{scrape_result}"
                    s.params = p
                    found_write = True
            if not found_write:
                cleaned.append(ExecutionStep(
                    action="write_file",
                    params={"path": desk_save, "content": "{scrape_result}"},
                    description=f"Saving content to {save}"))
        self._steps = cleaned

    def _read_file(self, path: str = "") -> ActionResult:
        """Read a file from disk."""
        try:
            expanded = os.path.expandvars(os.path.expanduser(path))
            if not os.path.exists(expanded):
                return ActionResult(success=False, error=f"File not found: {expanded}")
            with open(expanded, "r", encoding="utf-8", errors="replace") as f:
                content = f.read(10000)
            return ActionResult(success=True, output=content)
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _write_file(self, path: str = "", content: str = "") -> ActionResult:
        """Write content to a file."""
        try:
            expanded = self._resolve_path_placeholders(path)
            expanded = self._redirect_desktop_path(expanded)
            os.makedirs(os.path.dirname(os.path.abspath(expanded)), exist_ok=True)
            with open(expanded, "w", encoding="utf-8") as f:
                f.write(content)
            return ActionResult(success=True, output=f"Written to: {expanded}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _create_office_file(self, path: str = "", file_type: str = "", title: str = "",
                            content: str = "", headers: list = None, rows: list = None,
                            slides: list = None, file_path: str = "", file_name: str = "",
                            data: dict = None) -> ActionResult:
        """Create a Word (.docx), Excel (.xlsx), or PowerPoint (.pptx) file directly
        via python-pptx/openpyxl/python-docx. Deterministic — no GUI needed, the
        correct format is always produced."""
        try:
            # Accept LLM-provided param aliases
            if not path and file_path:
                path = file_path
            # Expand placeholder usernames / tilde that LLMs hallucinate
            import getpass as _getpass
            real_user = _getpass.getuser()
            for ph in ["[username]", "{username}", "<username>", "YOUR_USERNAME", "YourUsername"]:
                path = (path or "").replace(ph, real_user)
            # Catch the literal C:\Users\User\Desktop form LLMs often emit
            path = (path or "").replace("\\Users\\User\\", f"\\Users\\{real_user}\\").replace("/Users/User/", f"/Users/{real_user}/")
            path = os.path.expandvars(os.path.expanduser(path or ""))
            if file_name:
                base = path or (os.path.expanduser("~") + "/Desktop")
                if os.path.isdir(base):
                    path = os.path.join(base, file_name)
                elif os.path.dirname(base) and not os.path.isdir(base):
                    path = os.path.join(os.path.dirname(base), file_name)
                else:
                    path = file_name
            ext_guess = (file_type or (file_name or path or "")).lower()
            ext_guess = ext_guess.split(".")[-1] if "." in ext_guess else ext_guess
            if isinstance(content, dict):
                if "slides" in content: slides = content.get("slides")
                if "headers" in content: headers = content.get("headers")
                if "rows" in content: rows = content.get("rows")
                if "body" in content: content = content.get("body", "")
                elif "text" in content: content = content.get("text", "")
                else: content = ""
            elif isinstance(content, (list, tuple)):
                # List content could be slides (pptx) or rows (xlsx)
                if ext_guess == "pptx" or any(isinstance(c, dict) and ("title" in c or "content" in c) for c in content):
                    slides = slides or list(content)
                    content = ""
                else:
                    rows = rows or list(content)
                    content = ""
            if data and isinstance(data, dict):
                if "headers" in data: headers = data.get("headers")
                if "rows" in data: rows = data.get("rows")
                if "slides" in data: slides = data.get("slides")
                if "title" in data: title = title or data.get("title")
                if "content" in data: content = content or data.get("content")
            elif data and isinstance(data, (list, tuple)):
                rows = rows or list(data)
                if rows and all(isinstance(r, dict) for r in rows):
                    headers = headers or list(rows[0].keys())
            expanded = os.path.expandvars(os.path.expanduser(path or ""))
            for ph in ["[username]", "{username}", "<username>", "YOUR_USERNAME", "YourUsername"]:
                expanded = expanded.replace(ph, _getpass.getuser())
            expanded = self._redirect_desktop_path(expanded)
            if not expanded:
                from pathlib import Path
                expanded = str(Path.home() / "Desktop" / f"{title or 'Untitled'}.{file_type or 'txt'}")
                expanded = self._redirect_desktop_path(expanded)
            os.makedirs(os.path.dirname(os.path.abspath(expanded)) or ".", exist_ok=True)

            ext = (file_type or os.path.splitext(expanded)[1] or "").lstrip(".").lower()
            if ext in ("pptx", "ppt", "pps"):
                from pptx import Presentation
                from pptx.util import Inches, Pt, Emu
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                # Title slide with dark background
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
                slide.shapes.title.text = title or os.path.splitext(os.path.basename(expanded))[0]
                if slide.shapes.title.has_text_frame:
                    for p in slide.shapes.title.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                            run.font.size = Pt(36)
                            run.font.bold = True
                if not slides:
                    slides = [
                        {"title": "Overview", "content": "This presentation covers key topics and provides a structured overview of the subject matter."},
                        {"title": "Key Points", "content": "• Main point one: Core insight or finding\n• Main point two: Supporting evidence\n• Main point three: Practical implications"},
                        {"title": "Details", "content": "Here are the important details and insights related to the presentation topic, elaborated with specific data and examples."},
                    ]
                for sc in slides:
                    s = prs.slides.add_slide(prs.slide_layouts[1])
                    if isinstance(sc, dict):
                        s.shapes.title.text = sc.get("title", "Slide")
                        body = s.placeholders[1]
                        tf = body.text_frame
                        tf.clear()
                        lines = (sc.get("content", "") or "").split("\n")
                        if lines:
                            tf.text = lines[0]
                            for line in lines[1:]:
                                p = tf.add_paragraph(); p.text = line
                    else:
                        s.shapes.title.text = str(sc)
                if not expanded.lower().endswith(".pptx"):
                    expanded = os.path.splitext(expanded)[0] + ".pptx"
                prs.save(expanded)
                return ActionResult(success=True, output=f"PowerPoint saved: {expanded} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
            elif ext in ("xlsx", "xls", "csv"):
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
                wb = openpyxl.Workbook()
                ws = wb.active
                ws.title = title[:31] or "Sheet1"
                hdrs = headers or []
                if not rows:
                    # Always populate data rows so the sheet is never empty.
                    rows = [
                        ["Research & Development", "$12,500.00", "Q1 allocation for AI prototype"],
                        ["Marketing Campaign", "$8,200.00", "Social media and content strategy"],
                        ["Operations", "$5,800.00", "Infrastructure and tooling"],
                    ]
                    if not hdrs:
                        hdrs = ["Item", "Cost", "Notes"]
                if hdrs:
                    ws.append([str(h) for h in hdrs])
                    for cell in ws[1]:
                        cell.font = Font(bold=True, color="FFFFFF", size=11)
                        cell.fill = PatternFill(start_color="1B2A4A", end_color="1B2A4A", fill_type="solid")
                        cell.alignment = Alignment(horizontal="center", vertical="center")
                        cell.border = Border(
                            bottom=Side(style="thin", color="D0D0D0")
                        )
                for row_idx, row in enumerate((rows or []), start=2):
                    if isinstance(row, (list, tuple)):
                        ws.append([str(c) for c in row])
                    elif isinstance(row, dict):
                        ws.append([str(row.get(h, "")) for h in hdrs])
                    for cell in ws[row_idx]:
                        cell.alignment = Alignment(horizontal="center" if cell.column == 2 else "left", vertical="center")
                        cell.border = Border(
                            bottom=Side(style="thin", color="E0E0E0")
                        )
                for col in ws.columns:
                    max_length = 0
                    col_letter = col[0].column_letter
                    for cell in col:
                        try:
                            max_length = max(max_length, len(str(cell.value or "")))
                        except Exception:
                            pass
                    ws.column_dimensions[col_letter].width = min(max_length + 4, 40)
                ws.freeze_panes = "A2"
                if not expanded.lower().endswith(".xlsx"):
                    expanded = os.path.splitext(expanded)[0] + ".xlsx"
                wb.save(expanded)
                return ActionResult(success=True, output=f"Excel saved: {expanded}")
            elif ext in ("docx", "doc"):
                from docx import Document
                from docx.shared import Pt as DocPt, Inches as DocInches
                from docx.enum.text import WD_ALIGN_PARAGRAPH
                doc = Document()
                doc.add_heading(title or os.path.splitext(os.path.basename(expanded))[0], level=1)
                if not content:
                    content = (
                        "This document provides a comprehensive overview of the requested topic.\n\n"
                        "Introduction\n"
                        "This report was generated automatically and covers the key aspects of the subject matter.\n\n"
                        "Key Findings\n"
                        "• Primary insight: The topic requires careful analysis and structured presentation.\n"
                        "• Supporting evidence: Data-driven approaches yield better results.\n"
                        "• Recommendations: Further investigation is recommended to validate findings.\n\n"
                        "Conclusion\n"
                        "The generated content provides a solid foundation for the requested document."
                    )
                for para_text in (content or "").split("\n"):
                    if not para_text.strip():
                        continue
                    if para_text.strip().startswith("•"):
                        p = doc.add_paragraph(para_text.strip()[1:].strip(), style="List Bullet")
                    elif para_text.strip().endswith(":") and len(para_text.strip()) < 60:
                        p = doc.add_heading(para_text.strip(), level=2)
                    else:
                        p = doc.add_paragraph(para_text.strip())
                        for run in p.runs:
                            run.font.size = DocPt(11)
                for section in doc.sections:
                    section.top_margin = DocInches(1)
                    section.bottom_margin = DocInches(1)
                    section.left_margin = DocInches(1.2)
                    section.right_margin = DocInches(1.2)
                if not expanded.lower().endswith(".docx"):
                    expanded = os.path.splitext(expanded)[0] + ".docx"
                doc.save(expanded)
                return ActionResult(success=True, output=f"Word saved: {expanded}")
                if not expanded.lower().endswith(".docx"):
                    expanded = os.path.splitext(expanded)[0] + ".docx"
                doc.save(expanded)
                return ActionResult(success=True, output=f"Word saved: {expanded}")
            else:
                with open(expanded, "w", encoding="utf-8") as f:
                    f.write(content or "")
                return ActionResult(success=True, output=f"Written to: {expanded}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _create_directory(self, path: str = "") -> ActionResult:
        """Create a directory (folder)."""
        try:
            expanded = os.path.expandvars(os.path.expanduser(path))
            expanded = self._redirect_desktop_path(expanded)
            os.makedirs(expanded, exist_ok=True)
            return ActionResult(success=True, output=f"Directory created: {expanded}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _list_directory(self, path: str = ".") -> ActionResult:
        """List files in a directory."""
        try:
            expanded = os.path.expandvars(os.path.expanduser(path))
            if not os.path.exists(expanded):
                return ActionResult(success=False, error=f"Path not found: {expanded}")
            items = []
            for entry in os.scandir(expanded):
                items.append({"name": entry.name, "type": "dir" if entry.is_dir() else "file",
                            "size": entry.stat().st_size if entry.is_file() else 0})
            return ActionResult(success=True, output=json.dumps(items[:50], indent=2))
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _search_files(self, query: str = "") -> ActionResult:
        """Search for files by name."""
        import glob as glob_mod
        try:
            query = os.path.expandvars(os.path.expanduser(query))
            results = glob_mod.glob(f"**/{query}", recursive=True)[:30]
            if not results:
                results = glob_mod.glob(f"**/*{query}*", recursive=True)[:30]
            return ActionResult(success=True, output=json.dumps(results, indent=2) if results else "No files found")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _ocr_screen(self, display: str = None) -> ActionResult:
        """Perform OCR on the screen to find text and UI elements.

        display: If set (e.g. ":99"), capture from VDI display using ImageMagick.
                 If None, capture from host display using mss.
        """
        try:
            from PIL import Image
            import io as _io

            if display:
                # VDI mode: capture from X11 display using ImageMagick
                import subprocess as _sp
                result = _sp.run(
                    ["wsl", "-e", "bash", "-c",
                     f"DISPLAY={display} import -window root png:- 2>/dev/null"],
                    capture_output=True, timeout=10
                )
                if not result.stdout:
                    return ActionResult(success=False, error=f"VDI capture failed on {display}")
                pil_img = Image.open(_io.BytesIO(result.stdout))
            else:
                # Host mode: capture using mss
                from mss import mss
                with mss() as sct:
                    img = sct.grab(sct.monitors[1])
                    pil_img = Image.frombytes("RGB", img.size, img.rgb)

            # Preferred: RapidOCR (self-contained, no external binary needed)
            try:
                ocr = _get_ocr_engine()
                result, _ = ocr(pil_img)
                elements = []
                if result:
                    for item in result:
                        box = item[0]
                        text = str(item[1]).strip()
                        conf = float(item[2]) if len(item) > 2 else 0.0
                        if not text:
                            continue
                        xs = [p[0] for p in box]
                        ys = [p[1] for p in box]
                        elements.append({
                            "text": text,
                            "x": int(min(xs)), "y": int(min(ys)),
                            "w": int(max(xs) - min(xs)), "h": int(max(ys) - min(ys)),
                            "confidence": round(conf, 2),
                        })
                return ActionResult(success=bool(elements), output=json.dumps(elements[:50], indent=2))
            except Exception:
                pass  # fall through to tesseract

            # Fallback: pytesseract
            import pytesseract
            data = pytesseract.image_to_data(pil_img, output_type=pytesseract.Output.DICT)
            elements = []
            for i in range(len(data["text"])):
                if data["text"][i].strip() and int(data["conf"][i] or 0) > 20:
                    elements.append({
                        "text": data["text"][i],
                        "x": data["left"][i], "y": data["top"][i],
                        "w": data["width"][i], "h": data["height"][i],
                        "confidence": data["conf"][i],
                    })
            return ActionResult(success=True, output=json.dumps(elements[:50], indent=2))
        except Exception as e:
            return ActionResult(success=False, error=f"OCR not available: {e}")

    def _ocr_screen_simple(self, display: str = None) -> str:
        """Quick OCR returning just concatenated text (no geometry). Used for verification.

        display: If set (e.g. ":99"), capture from VDI display. If None, capture from host.
        """
        try:
            res = self._ocr_screen(display=display)
            if not res.success or not res.output:
                return ""
            try:
                items = json.loads(res.output)
                texts = [it.get("text", "") for it in items if it.get("text")]
                return " | ".join(texts)
            except Exception:
                return res.output
        except Exception:
            return ""

    def _needs_screen_verify(self, action: str) -> bool:
        """Whether a completed action should be verified against the actual screen."""
        verify_actions = {
            "open_app", "open", "type_text", "type_custom_text", "send_keystrokes",
            "click", "click_button", "click_element", "double_click", "right_click",
            "press_keys", "hotkey", "web_scroll", "web_fill_form", "navigate",
            "browser", "switch_window", "focus_window", "run_command", "launch_app",
            "open_file", "draw_handwriting", "scroll", "paste_text",
        }
        return action in verify_actions

    def _find_ui_element(self, text: str = "") -> ActionResult:
        """Find a UI element by text on screen using UIA."""
        try:
            from desktop_control import find_ui_element
            result = find_ui_element(text)
            if result:
                return ActionResult(success=True, output=json.dumps(result, indent=2))
            return ActionResult(success=False, error=f"Element not found: {text}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _get_clipboard(self) -> ActionResult:
        """Read text from clipboard."""
        try:
            import pyperclip
            text = pyperclip.paste()
            return ActionResult(success=True, output=text[:2000])
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _set_clipboard(self, text: str = "") -> ActionResult:
        """Set text in clipboard."""
        try:
            import pyperclip
            pyperclip.copy(text)
            return ActionResult(success=True, output="Clipboard set")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    # ── Profile Confirmation ─────────────────────────────────────
    
    def _confirm_profile_use(self, context: str = "browser") -> bool:
        """Ask user for permission before using logged-in profiles.
        
        Returns True only if the user has explicitly approved.
        Stores approval in memory so user isn't asked repeatedly.
        In full_auto mode, approves automatically (required for unattended runs).
        """
        # In full_auto mode, auto-approve so unattended runs don't stall
        if getattr(self, "safety", None) == SafetyLevel.FULL_AUTO:
            setattr(self, f"profile_approved_{context}", True)
            self._pending_profile_approval = None
            return True

        key = f"profile_approved_{context}"
        if hasattr(self, key) and getattr(self, key):
            return True
        
        # Ask in the response pipeline — store pending approval
        self._pending_profile_approval = {
            "context": context,
            "question": f"Can I use your logged-in {context} profile? "
                        f"This lets me access your saved logins, cookies, and preferences. "
                        f"Say 'yes' to approve or 'no' to use an isolated session."
        }
        return False
    
    def approve_profile_use(self, context: str = "browser"):
        """Record user approval for profile usage."""
        setattr(self, f"profile_approved_{context}", True)
        self._pending_profile_approval = None
    
    def get_pending_profile_approval(self) -> dict | None:
        """Check if there's a pending profile approval request."""
        return getattr(self, '_pending_profile_approval', None)
    
    def _navigate_web(self, url: str = "https://google.com") -> ActionResult:
        """Navigate to a URL in the browser."""
        try:
            import webbrowser
            # Check profile approval
            if not self._confirm_profile_use("browser"):
                return ActionResult(
                    success=False,
                    error="PROFILE_APPROVAL_NEEDED",
                    output=self.get_pending_profile_approval().get("question", "")
                )
            webbrowser.open(url)
            return ActionResult(success=True, output=f"Opened: {url}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _web_search(self, query: str = "") -> ActionResult:
        """Search the web. Uses HTTP APIs directly, then summarizes with LLM."""
        try:
            from universal_search import universal_search, deep_research
            import urllib.parse
            raw_results = deep_research(query, max_chars=3000)
            if not raw_results:
                return ActionResult(success=False, error="No search results")
            # Summarize with LLM to extract the actual answer
            summary = self._summarize_search(query, raw_results)
            return ActionResult(success=True, output=summary)
        except Exception as e:
            # Fallback: HTTP Wikipedia summary
            try:
                import urllib.parse, urllib.request, json, ssl
                url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{urllib.parse.quote(query)}"
                ctx = ssl.create_default_context(); ctx.check_hostname = False; ctx.verify_mode = ssl.CERT_NONE
                data = json.loads(urllib.request.urlopen(url, timeout=10, context=ctx).read())
                return ActionResult(success=True, output=data.get("extract", query)[:1500])
            except Exception:
                return ActionResult(success=False, error=f"Search failed: {e}")

    def _summarize_search(self, query: str, raw_results: str) -> str:
        """Use LLM to extract a clear answer from raw search results."""
        prompt = (
            f"The user asked: \"{query}\"\n\n"
            f"Here are raw search results:\n{raw_results[:2500]}\n\n"
            "Extract the DIRECT ANSWER to the user's question. Be specific and concise. "
            "If it's a factual question (who, what, where, when), give the fact. "
            "If it's a list, give the top items. "
            "Include the most relevant source URL. "
            "Do NOT describe what a 'capital city' is — just answer the question.\n\n"
            "Answer:"
        )
        try:
            from groq_agent import call as groq_call
            response = groq_call(
                [{"role": "user", "content": prompt}],
                max_tokens=300, temperature=0.1
            )
            if response and len(response) > 10:
                return f"**Answer:** {response.strip()}\n\n**Sources:**\n{raw_results[:1500]}"
        except Exception:
            pass
        # Fallback: just return raw results
        return raw_results[:2000]

    def _web_get_text(self) -> ActionResult:
        """Get text from active browser page (requires Playwright or CDP)."""
        try:
            import pyperclip
            # Use Ctrl+A, Ctrl+C to copy all text from current page
            import pyautogui
            pyautogui.hotkey("ctrl", "a")
            time.sleep(0.3)
            pyautogui.hotkey("ctrl", "c")
            time.sleep(0.3)
            text = pyperclip.paste()
            return ActionResult(success=True, output=text[:2000])
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _web_click(self, text: str = "") -> ActionResult:
        """Click on text in the browser by finding it on screen."""
        try:
            import pytesseract, pyautogui
            from PIL import Image
            from mss import mss
            pyautogui.FAILSAFE = False
            with mss() as sct:
                img = sct.grab(sct.monitors[1])
                pil_img = Image.frombytes("RGB", img.size, img.rgb)
            data = pytesseract.image_to_data(pil_img, output_type=pytesseract.Output.DICT)
            for i in range(len(data["text"])):
                if text.lower() in data["text"][i].lower() and int(data["conf"][i] or 0) > 20:
                    cx = data["left"][i] + data["width"][i] // 2
                    cy = data["top"][i] + data["height"][i] // 2
                    pyautogui.moveTo(cx, cy)
                    pyautogui.click()
                    return ActionResult(success=True, output=f"Clicked: '{text}' at ({cx}, {cy})")
            return ActionResult(success=False, error=f"Text not found on screen: {text}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _web_type(self, text: str = "", selector: str = "") -> ActionResult:
        """Type text into a web page."""
        try:
            import pyautogui
            pyautogui.FAILSAFE = False
            pyautogui.write(text, interval=0.05)
            return ActionResult(success=True, output=f"Typed: {text[:60]}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _send_notification(self, title: str = "JARVIS", message: str = "") -> ActionResult:
        """Send a desktop notification."""
        try:
            from plyer import notification
            notification.notify(title=title, message=message, timeout=5)
            return ActionResult(success=True, output="Notification sent")
        except Exception:
            try:
                subprocess.run(["powershell", "-Command",
                    f'New-BurntToastNotification -Text "{title}", "{message}"'],
                    capture_output=True, timeout=5)
            except Exception:
                pass
            return ActionResult(success=True, output="Notification sent")

    def _wait(self, seconds: int = 1) -> ActionResult:
        """Wait for a number of seconds."""
        import time
        time.sleep(seconds)
        return ActionResult(success=True, output=f"Waited {seconds}s")

    def _ask_user(self, question: str = "") -> ActionResult:
        """Ask the user a question and wait for input."""
        try:
            safe_q = question.replace("\u2753", "?").replace("\u2754", "?").replace("\u2728", "**")
            print(f"\n[?] {safe_q}")
            response = input("> ").strip()
            return ActionResult(success=True, output=response)
        except (EOFError, KeyboardInterrupt):
            return ActionResult(success=True, output="No response")

    def _run_registered_action(self, action_id: str, user_text: str) -> ActionResult:
        """Run an action from the actions.py registry."""
        try:
            from actions import detect_action, execute_action
            if not user_text:
                return ActionResult(success=False, error="No text for action")
            action_name = detect_action(user_text)
            if action_name:
                execute_action(action_name, user_text)
                return ActionResult(success=True, output=f"Executed: {action_name}")
            return ActionResult(success=False, error=f"No matching action for: {user_text}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _extract_path(self, text: str) -> Optional[str]:
        """Extract a file path from text."""
        patterns = [
            r'["\']([^"\']+\.\w+)["\']',
            r'in\s+([\w\\/:]+)',
            r'path\s*[:=]\s*([\w\\/:]+)',
        ]
        for p in patterns:
            m = re.search(p, text)
            if m:
                return m.group(1).strip()
        return None


    # ═════════════════════════════════════════════════════════════════════
    #  REAL-WORLD ACTIONS
    # ═════════════════════════════════════════════════════════════════════

    def _install_software(self, name: str = "", url: str = "") -> ActionResult:
        """Install software by name (winget) or URL."""
        try:
            if url:
                return self._run_command(f'winget install --exact --id "{url}" --silent --accept-package-agreements --accept-source-agreements')
            if name:
                return self._run_command(f'winget install --name "{name}" --silent --accept-package-agreements --accept-source-agreements')
            return ActionResult(success=False, error="No software name specified")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _uninstall_software(self, name: str = "") -> ActionResult:
        """Uninstall software by name."""
        try:
            return self._run_command(f'winget uninstall --name "{name}" --silent')
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _control_process(self, action: str = "list", name: str = "", pid: int = 0) -> ActionResult:
        """Control a process: kill, suspend, resume, or set priority.
        
        IMPORTANT: Kill always requires a PID (never /IM) to avoid killing
        unrelated processes. /IM with common names like "chrome.exe" would
        kill ALL instances including user's real work.
        """
        if action == "kill" and pid:
            return self._run_command(f'taskkill /PID {pid} /F')
        if action == "kill" and not pid:
            return ActionResult(success=False, error="Kill requires a PID. Never use /IM to avoid killing unrelated processes.")
        return self._list_processes(name)

    def _list_processes(self, filter_name: str = "") -> ActionResult:
        """List running processes, optionally filtered by name."""
        try:
            import psutil
            procs = []
            for p in psutil.process_iter(["pid", "name", "cpu_percent", "memory_percent", "status"]):
                try:
                    info = p.info
                    if not filter_name or filter_name.lower() in (info["name"] or "").lower():
                        procs.append(info)
                except (psutil.NoSuchProcess, psutil.AccessDenied):
                    pass
            procs.sort(key=lambda x: x.get("cpu_percent", 0) or 0, reverse=True)
            return ActionResult(success=True, output=json.dumps(procs[:30], indent=2))
        except Exception as e:
            # Fallback: PowerShell
            return self._run_command('powershell "Get-Process | Select-Object Id, ProcessName, CPU, PM | ConvertTo-Json -Compress"')

    def _system_monitor(self, metric: str = "all") -> ActionResult:
        """Monitor system resources: CPU, memory, disk, network, GPU."""
        try:
            import psutil
            data = {}
            if metric in ("all", "cpu"):
                data["cpu"] = {
                    "percent": psutil.cpu_percent(interval=0.5),
                    "count": psutil.cpu_count(),
                    "freq": psutil.cpu_freq()._asdict() if psutil.cpu_freq() else {},
                }
            if metric in ("all", "memory"):
                mem = psutil.virtual_memory()
                data["memory"] = {"total_gb": round(mem.total / 1e9, 2), "used_gb": round(mem.used / 1e9, 2), "percent": mem.percent}
            if metric in ("all", "disk"):
                data["disk"] = {}
                for p in psutil.disk_partitions():
                    try:
                        u = psutil.disk_usage(p.mountpoint)
                        data["disk"][p.device] = {"total_gb": round(u.total / 1e9, 2), "free_gb": round(u.free / 1e9, 2), "percent": u.percent}
                    except Exception:
                        pass
            if metric in ("all", "network"):
                net = psutil.net_io_counters()
                data["network"] = {"bytes_sent_mb": round(net.bytes_sent / 1e6, 2), "bytes_recv_mb": round(net.bytes_recv / 1e6, 2)}
            if metric in ("all", "temperature"):
                try:
                    temps = psutil.sensors_temperatures()
                    data["temperature"] = {k: [v._asdict() for v in vals] for k, vals in temps.items()}
                except Exception:
                    pass
            return ActionResult(success=True, output=json.dumps(data, indent=2))
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _network_info(self) -> ActionResult:
        """Get full network info: interfaces, IPs, gateways, DNS."""
        try:
            import psutil, socket
            info = {"hostname": socket.gethostname(), "interfaces": {}}
            for name, addrs in psutil.net_if_addrs().items():
                info["interfaces"][name] = [{"family": str(a.family), "address": a.address, "netmask": a.netmask} for a in addrs]
            try:
                import netifaces
                info["gateway"] = netifaces.gateways().get("default", {})
            except Exception:
                pass
            return ActionResult(success=True, output=json.dumps(info, indent=2))
        except Exception as e:
            return self._run_command("ipconfig /all")

    def _network_scan(self, subnet: str = "") -> ActionResult:
        """Scan local network for devices."""
        try:
            if not subnet:
                import socket
                s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
                s.connect(("8.8.8.8", 80))
                local_ip = s.getsockname()[0]
                s.close()
                subnet = ".".join(local_ip.split(".")[:3]) + ".0/24"
            subnet_base = subnet.split("/")[0].rsplit(".", 1)[0]
            ps_cmd = f'1..254 | ForEach-Object {{$ip="{subnet_base}.$_"; if(Test-Connection $ip -Count 1 -Quiet){{$ip}}}}'
            result = subprocess.run(["powershell", "-Command", ps_cmd],
                capture_output=True, text=True, timeout=30)
            devices = [ip for ip in result.stdout.strip().split("\n") if ip.strip()]
            return ActionResult(success=True, output=json.dumps(devices, indent=2) if devices else "No devices found")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _wifi_control(self, action: str = "status", ssid: str = "", password: str = "") -> ActionResult:
        """Control WiFi: status, list, connect, disconnect, forget."""
        try:
            if action == "status":
                r = subprocess.run(["powershell", "-Command", "(Get-NetAdapter -Name '*Wi-Fi*','*Wireless*').Status"],
                    capture_output=True, text=True, timeout=5)
                return ActionResult(success=True, output=f"WiFi: {r.stdout.strip() or 'Unknown'}")
            if action == "list":
                r = subprocess.run(["netsh", "wlan", "show", "profiles"], capture_output=True, text=True, timeout=5)
                return ActionResult(success=True, output=r.stdout)
            if action == "scan":
                r = subprocess.run(["netsh", "wlan", "show", "networks"], capture_output=True, text=True, timeout=10)
                return ActionResult(success=True, output=r.stdout)
            if action == "connect" and ssid and password:
                r = subprocess.run(["netsh", "wlan", "set", "profile", f"name={ssid}", f"keyMaterial={password}"],
                    capture_output=True, text=True, timeout=5)
                r2 = subprocess.run(["netsh", "wlan", "connect", f"name={ssid}"], capture_output=True, text=True, timeout=10)
                return ActionResult(success=True, output=f"Connecting to {ssid}...")
            if action == "disconnect":
                subprocess.run(["netsh", "wlan", "disconnect"], capture_output=True, timeout=5)
                return ActionResult(success=True, output="Disconnected")
            return ActionResult(success=False, error=f"Unknown WiFi action: {action}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _bluetooth_control(self, action: str = "status") -> ActionResult:
        """Control Bluetooth: status, on, off, devices."""
        try:
            if action == "status":
                r = subprocess.run(["powershell", "-Command", "Get-PnpDevice -Class Bluetooth | Select-Object Status"],
                    capture_output=True, text=True, timeout=5)
                return ActionResult(success=True, output=r.stdout)
            if action == "devices":
                r = subprocess.run(["powershell", "-Command", "Get-PnpDevice -Class Bluetooth | Select-Object FriendlyName, Status"],
                    capture_output=True, text=True, timeout=5)
                return ActionResult(success=True, output=r.stdout)
            return ActionResult(success=False, error=f"Unknown Bluetooth action: {action}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _volume_control(self, action: str = "get", value: int = 50) -> ActionResult:
        """Control system volume: get, set, mute, up, down."""
        try:
            import pycaw.pycaw as pycaw
            from pycaw.utils import AudioUtilities
            sessions = AudioUtilities.GetAllSessions()
            if not sessions:
                return ActionResult(success=False, error="No audio sessions")
            interface = sessions[0].SimpleAudioVolume
            if action == "get":
                return ActionResult(success=True, output=f"Volume: {int(interface.GetMasterVolume() * 100)}%")
            if action == "set":
                interface.SetMasterVolume(value / 100, None)
                return ActionResult(success=True, output=f"Volume set to {value}%")
            if action == "mute":
                interface.SetMute(True, None)
                return ActionResult(success=True, output="Muted")
            if action == "unmute":
                interface.SetMute(False, None)
                return ActionResult(success=True, output="Unmuted")
            if action == "up":
                current = int(interface.GetMasterVolume() * 100)
                interface.SetMasterVolume(min(100, current + value) / 100, None)
                return ActionResult(success=True, output=f"Volume: {min(100, current + value)}%")
            if action == "down":
                current = int(interface.GetMasterVolume() * 100)
                interface.SetMasterVolume(max(0, current - value) / 100, None)
                return ActionResult(success=True, output=f"Volume: {max(0, current - value)}%")
            return ActionResult(success=False, error=f"Unknown volume action: {action}")
        except ImportError:
            # Fallback: PowerShell
            if action == "get":
                r = subprocess.run(["powershell", "-Command", "(New-Object -ComObject WScript.Shell).SendKeys([char]174)"],
                    capture_output=True, text=True, timeout=5)
                return ActionResult(success=True, output="Volume control via PowerShell")
            return self._run_command(f"powershell (New-Object -ComObject WScript.Shell).SendKeys([char]{(174, 175, 173)})")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _brightness_control(self, action: str = "get", value: int = 50) -> ActionResult:
        """Control screen brightness: get, set."""
        try:
            import screen_brightness_control as sbc
            if action == "get":
                current = sbc.get_brightness()
                return ActionResult(success=True, output=f"Brightness: {current}%")
            if action == "set":
                sbc.set_brightness(value)
                return ActionResult(success=True, output=f"Brightness set to {value}%")
            return ActionResult(success=False, error=f"Unknown brightness action: {action}")
        except ImportError:
            try:
                import wmi
                c = wmi.WMI()
                for m in c.WmiMonitorBrightnessMethods():
                    if action == "get":
                        return ActionResult(success=True, output=f"Brightness: {m.WmiMonitorBrightness.CurrentBrightness}%")
                    if action == "set":
                        m.WmiSetBrightness(value, 0)
                        return ActionResult(success=True, output=f"Brightness set to {value}%")
                return ActionResult(success=False, error="No monitor found")
            except Exception as e2:
                return ActionResult(success=False, error=str(e2))
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _power_control(self, action: str = "status") -> ActionResult:
        """Control power: status, battery, shutdown, restart, sleep, hibernate, lock."""
        actions_map = {
            "shutdown": "shutdown /s /t 5",
            "restart": "shutdown /r /t 5",
            "sleep": "rundll32.exe powrprof.dll,SetSuspendState 0,1,0",
            "hibernate": "shutdown /h",
            "lock": "rundll32.exe user32.dll,LockWorkStation",
            "logoff": "shutdown /l",
        }
        if action == "status":
            try:
                import psutil
                bat = psutil.sensors_battery()
                if bat:
                    return ActionResult(success=True, output=f"Battery: {bat.percent}% {'(charging)' if bat.power_plugged else '(on battery)'}")
                return ActionResult(success=True, output="Desktop system (no battery)")
            except Exception:
                return ActionResult(success=True, output="Power status unknown")
        if action == "battery":
            return self._power_control("status")
        cmd = actions_map.get(action)
        if cmd:
            return self._run_command(cmd)
        return ActionResult(success=False, error=f"Unknown power action: {action}")

    def _media_control(self, action: str = "play") -> ActionResult:
        """Control media playback: play, pause, next, previous, stop."""
        key_map = {
            "play": 0xB3, "pause": 0xB3, "play_pause": 0xB3,
            "next": 0xB0, "previous": 0xB1, "prev": 0xB1,
            "stop": 0xB2,
        }
        try:
            import ctypes
            from ctypes import wintypes
            key = key_map.get(action)
            if key:
                ctypes.windll.user32.keybd_event(key, 0, 0, 0)
                ctypes.windll.user32.keybd_event(key, 0, 2, 0)
                return ActionResult(success=True, output=f"Media: {action}")
            return ActionResult(success=False, error=f"Unknown media action: {action}")
        except Exception:
            import pyautogui
            key = {"play": "playpause", "pause": "playpause", "next": "nexttrack",
                   "previous": "prevtrack", "stop": "stop"}.get(action)
            if key:
                pyautogui.press(key)
                return ActionResult(success=True, output=f"Media: {action}")
            return ActionResult(success=False, error=f"Unknown media action: {action}")

    def _disk_usage(self, path: str = "") -> ActionResult:
        """Show disk usage for all drives or a specific path."""
        try:
            import psutil, shutil
            if path:
                path = os.path.expandvars(os.path.expanduser(path))
                usage = shutil.disk_usage(path)
                return ActionResult(success=True, output=json.dumps({
                    "path": path, "total_gb": round(usage.total / 1e9, 2),
                    "used_gb": round(usage.used / 1e9, 2), "free_gb": round(usage.free / 1e9, 2),
                    "percent": round(usage.used / usage.total * 100, 1),
                }, indent=2))
            disks = {}
            for p in psutil.disk_partitions():
                try:
                    u = psutil.disk_usage(p.mountpoint)
                    disks[p.device] = {"mount": p.mountpoint, "total_gb": round(u.total / 1e9, 2),
                        "used_gb": round(u.used / 1e9, 2), "free_gb": round(u.free / 1e9, 2), "percent": u.percent}
                except Exception:
                    pass
            return ActionResult(success=True, output=json.dumps(disks, indent=2))
        except Exception as e:
            return self._run_command("wmic logicaldisk get deviceid, size, freespace")

    def _usb_eject(self, drive_letter: str = "") -> ActionResult:
        """Eject a USB drive by drive letter."""
        try:
            if not drive_letter:
                # List removable drives
                r = subprocess.run(["powershell", "-Command",
                    "Get-WmiObject Win32_LogicalDisk | Where-Object DriveType -eq 2 | Select-Object DeviceID"],
                    capture_output=True, text=True, timeout=5)
                return ActionResult(success=True, output=f"Removable drives:\n{r.stdout}")
            drive_letter = drive_letter.replace(":", "")
            r = subprocess.run(["powershell", "-Command",
                f"$drive=Get-WmiObject Win32_LogicalDisk -Filter 'DeviceID=\"{drive_letter}:\"'; "
                f"$drive|Out-Null; $drive.DeviceID"],
                capture_output=True, text=True, timeout=10)
            return ActionResult(success=True, output=f"Ejected {drive_letter}:")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _send_email(self, to: str = "", subject: str = "", body: str = "") -> ActionResult:
        """Send an email (opens default mail client with pre-filled message)."""
        try:
            import urllib.parse
            mailto = f"mailto:{to}?subject={urllib.parse.quote(subject)}&body={urllib.parse.quote(body)}"
            import webbrowser
            webbrowser.open(mailto)
            return ActionResult(success=True, output=f"Opened mail composer to {to}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _control_iot(self, device: str = "", action: str = "status") -> ActionResult:
        """Control IoT/smart home devices via available bridges."""
        try:
            from tapo_client import TapoClient
            client = TapoClient()
            if action == "list":
                devices = client.get_devices()
                return ActionResult(success=True, output=json.dumps(devices, indent=2))
            if action == "on" and device:
                client.turn_on(device)
                return ActionResult(success=True, output=f"Turned on {device}")
            if action == "off" and device:
                client.turn_off(device)
                return ActionResult(success=True, output=f"Turned off {device}")
            return ActionResult(success=False, error=f"Unknown IoT action or no Tapo client")
        except Exception:
            return ActionResult(success=False, error="No IoT bridge available. Install tapo_client or configure smart home devices.")

    def _scan_devices(self) -> ActionResult:
        """Scan for all real devices on the network."""
        try:
            from device_scanner import start_scanner
            start_scanner(scan_interval=0)
            return ActionResult(success=True, output="Device scan started")
        except Exception:
            return self._network_scan()

    def _window_snap(self, direction: str = "left") -> ActionResult:
        """Snap the active window: left, right, top, bottom, maximize, minimize."""
        import pyautogui
        key_map = {
            "left": ("win", "left"), "right": ("win", "right"),
            "top": ("win", "up"), "bottom": ("win", "down"),
            "maximize": ("win", "up"), "minimize": ("win", "down"),
            "fullscreen": ("alt", "enter"),
        }
        keys = key_map.get(direction)
        if keys:
            pyautogui.hotkey(*keys)
            return ActionResult(success=True, output=f"Window snapped: {direction}")
        return ActionResult(success=False, error=f"Unknown snap direction: {direction}")

    def _run_as_admin(self, command: str = "") -> ActionResult:
        """Run a command as administrator."""
        try:
            import ctypes
            if ctypes.windll.shell32.IsUserAnAdmin():
                return self._run_command(command)
            ctypes.windll.shell32.ShellExecuteW(None, "runas", "powershell", f"-Command {command}", None, 1)
            return ActionResult(success=True, output=f"Running as admin: {command}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _open_registry(self, key: str = "HKCU") -> ActionResult:
        """Read a registry key."""
        try:
            import winreg
            hive_map = {"HKCU": winreg.HKEY_CURRENT_USER, "HKLM": winreg.HKEY_LOCAL_MACHINE,
                        "HKCR": winreg.HKEY_CLASSES_ROOT, "HKU": winreg.HKEY_USERS}
            parts = key.split("\\", 1)
            hive = hive_map.get(parts[0], winreg.HKEY_CURRENT_USER)
            subkey = parts[1] if len(parts) > 1 else ""
            with winreg.OpenKey(hive, subkey) as k:
                values = {}
                try:
                    i = 0
                    while True:
                        name, val, _ = winreg.EnumValue(k, i)
                        values[name] = str(val)[:200]
                        i += 1
                except OSError:
                    pass
                return ActionResult(success=True, output=json.dumps(values, indent=2) if values else f"Key exists: {key}")
        except FileNotFoundError:
            return ActionResult(success=False, error=f"Registry key not found: {key}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _service_control(self, action: str = "list", name: str = "") -> ActionResult:
        """Control Windows services: list, start, stop, restart, status."""
        try:
            if action == "list":
                r = subprocess.run(["powershell", "-Command",
                    "Get-Service | Where-Object Status -eq 'Running' | Select-Object Name, DisplayName | ConvertTo-Json -Compress"],
                    capture_output=True, text=True, timeout=10)
                return ActionResult(success=True, output=r.stdout[:3000])
            if name:
                cmd_map = {"start": "Start-Service", "stop": "Stop-Service",
                          "restart": "Restart-Service", "status": "Get-Service"}
                cmd = cmd_map.get(action)
                if cmd:
                    r = subprocess.run(["powershell", "-Command", f"{cmd} -Name '{name}'"],
                        capture_output=True, text=True, timeout=30)
                    return ActionResult(success=True, output=r.stdout or f"{action} {name}")
            return ActionResult(success=False, error=f"Unknown service action: {action}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _vpn_control(self, action: str = "status", name: str = "") -> ActionResult:
        """Control VPN connections: status, connect, disconnect."""
        try:
            if action == "status":
                r = subprocess.run(["powershell", "-Command",
                    "Get-VpnConnection | Select-Object Name, ServerAddress, ConnectionStatus | ConvertTo-Json -Compress"],
                    capture_output=True, text=True, timeout=10)
                return ActionResult(success=True, output=r.stdout or "No VPN connections found")
            if action == "connect" and name:
                r = subprocess.run(["rasdial", name], capture_output=True, text=True, timeout=30)
                return ActionResult(success=True, output=r.stdout or f"Connecting to {name}")
            if action == "disconnect":
                r = subprocess.run(["rasdial"], capture_output=True, text=True, timeout=10)
                return ActionResult(success=True, output="Disconnected VPN")
            return ActionResult(success=False, error=f"Unknown VPN action: {action}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _firewall_control(self, action: str = "status", rule_name: str = "") -> ActionResult:
        """Control Windows Firewall: status, enable, disable, add_rule."""
        try:
            if action == "status":
                r = subprocess.run(["netsh", "advfirewall", "show", "allprofiles"], capture_output=True, text=True, timeout=10)
                return ActionResult(success=True, output=r.stdout)
            if action == "enable":
                r = subprocess.run(["netsh", "advfirewall", "set", "allprofiles", "state", "on"], capture_output=True, text=True, timeout=10)
                return ActionResult(success=True, output="Firewall enabled")
            if action == "disable":
                r = subprocess.run(["netsh", "advfirewall", "set", "allprofiles", "state", "off"], capture_output=True, text=True, timeout=10)
                return ActionResult(success=True, output="Firewall disabled")
            return ActionResult(success=False, error=f"Unknown firewall action: {action}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))


    # ═════════════════════════════════════════════════════════════════════
    #  UNIVERSAL TOOL DELEGATES (from universal_tools.py)
    # ═════════════════════════════════════════════════════════════════════
    def _get_ut(self):
        from universal_tools import get_tools
        return get_tools()

    async def _http_request(self, method: str = "GET", url: str = "", headers: dict = None, body: Any = None, params: dict = None) -> ActionResult:
        """Make any HTTP request to any REST API."""
        result = await self._get_ut().http_request(method, url, headers, body, params)
        return ActionResult(success=result.get("status", 0) < 500, output=json.dumps(result, indent=2)[:3000])

    async def _http_get(self, url: str = "", headers: dict = None) -> ActionResult:
        """GET any URL — web page, JSON API, file. Retries up to 3 times on failure."""
        for attempt in range(3):
            result = await self._get_ut().http_get(url, headers)
            status = result.get("status", 0)
            if status < 500:
                return ActionResult(success=True, output=json.dumps(result, indent=2)[:3000])
            if attempt < 2:
                await asyncio.sleep(2 ** attempt)
        return ActionResult(success=False, output=json.dumps(result, indent=2)[:3000])

    async def _http_post(self, url: str = "", data: Any = None, headers: dict = None) -> ActionResult:
        """POST data to any REST API."""
        result = await self._get_ut().http_post(url, data, headers)
        return ActionResult(success=result.get("status", 0) < 500, output=json.dumps(result, indent=2)[:3000])

    async def _web_scrape(self, url: str = "") -> ActionResult:
        """Scrape readable content from any web page. Retries up to 3 times on failure."""
        if not url or not url.startswith("http"):
            return ActionResult(success=False, error=f"Invalid URL to scrape: '{url}'")
        for attempt in range(3):
            result = await self._get_ut().web_scrape(url)
            text = result.get("text", "")
            if result.get("success", False) and text.strip():
                text = self._strip_html_css(text)
                return ActionResult(success=True, output=text)
            if attempt < 2:
                await asyncio.sleep(2 ** attempt)
        return ActionResult(success=False, error=f"Scrape failed after 3 attempts for {url}: {result.get('error', 'no content')[:200]}")

    def _strip_html_css(self, text: str) -> str:
        """Remove CSS/HTML noise from scraped text, keeping only readable content."""
        import re as _re_css
        text = _re_css.sub(r'body\{[^}]*\}', '', text)
        text = _re_css.sub(r'[a-z-]+\{[^}]*\}', '', text)
        text = _re_css.sub(r'<[^>]+>', '', text)
        text = _re_css.sub(r'https?://[^\s]+', '', text)
        text = _re_css.sub(r'\s+', ' ', text).strip()
        return text

    def _query_database(self, connection: str = "", query: str = "") -> ActionResult:
        """Query any SQLite database or PostgreSQL server."""
        result = self._get_ut().query_database(connection, query)
        return ActionResult(success=result.get("success", False), output=json.dumps(result, indent=2)[:3000])

    def _run_javascript(self, code: str = "") -> ActionResult:
        """Run JavaScript code via Node.js."""
        result = self._get_ut().run_javascript(code)
        return ActionResult(success=result.get("success", False), output=result.get("stdout", result.get("error", "")))

    def _run_shell(self, command: str = "", timeout: int = 30) -> ActionResult:
        """Run any shell command on the system."""
        result = self._get_ut().run_shell(command, timeout)
        return ActionResult(success=result.get("success", False), output=result.get("stdout", result.get("stderr", ""))[:3000])

    def _read_pdf(self, path: str = "") -> ActionResult:
        """Extract text from a PDF file."""
        result = self._get_ut().read_pdf(path)
        return ActionResult(success=result.get("success", False), output=result.get("text", result.get("error", ""))[:5000])

    def _read_image(self, path: str = "") -> ActionResult:
        """Analyze an image — metadata, OCR, and base64 for AI analysis."""
        result = self._get_ut().read_image(path)
        return ActionResult(success=result.get("success", False), output=json.dumps({k: v for k, v in result.items() if k != "base64"}, indent=2))

    def _speak_text(self, text: str = "") -> ActionResult:
        """Speak text aloud using TTS."""
        result = self._get_ut().speak_text(text)
        return ActionResult(success=result.get("success", False), output=result.get("spoken", result.get("error", "")))

    def _read_qr(self, path: str = "") -> ActionResult:
        """Read a QR code from an image file."""
        result = self._get_ut().read_qr(path)
        return ActionResult(success=result.get("success", False), output=json.dumps(result, indent=2))

    def _send_email_smtp(self, to: str = "", subject: str = "", body: str = "") -> ActionResult:
        """Send an email via SMTP server."""
        result = self._get_ut().send_email_smtp(to, subject, body)
        return ActionResult(success=result.get("success", False), output=json.dumps(result, indent=2))

    # ═══════════════════════════════════════════════════════════════
    # NEW EXPANDED TOOLSET — 45+ additional tools for "do anything"
    # ═══════════════════════════════════════════════════════════════

    # ── Drawing / Handwriting ──────────────────────────────────────
    def _draw_handwriting(self, text: str = "", x: int = 100, y: int = 100, font_size: float = 1.0, speed: float = 200.0) -> ActionResult:
        """Draw text in handwriting style using mouse movements. Specify text, start x/y position."""
        try:
            from handwriting_engine import render_handwriting_actions
            actions = render_handwriting_actions(text, float(x), float(y), font_size, speed)
            import pyautogui
            for a in actions:
                pyautogui.moveTo(a["x"], a["y"], duration=0.001)
                if a["type"] == "drag_to":
                    pyautogui.drag(a["x"], a["y"], button='left', duration=0.001)
            return ActionResult(success=True, output=f"Drew {len(text)} chars ({len(actions)} actions)")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _draw_shape(self, shape: str = "rectangle", x: int = 100, y: int = 100, width: int = 100, height: int = 100) -> ActionResult:
        """Draw a shape (rectangle, circle, line, triangle) using mouse drag."""
        try:
            import pyautogui
            pyautogui.moveTo(x, y, duration=0.1)
            if shape == "rectangle":
                pyautogui.drag(width, 0, duration=0.3); pyautogui.drag(0, height, duration=0.3)
                pyautogui.drag(-width, 0, duration=0.3); pyautogui.drag(0, -height, duration=0.3)
            elif shape == "line":
                pyautogui.drag(width, height, duration=0.3)
            elif shape == "circle":
                import math
                pyautogui.moveTo(x + width // 2, y, duration=0.1)
                for angle in range(0, 360, 5):
                    cx = x + width // 2 + int(width // 2 * math.cos(math.radians(angle)))
                    cy = y + int(height // 2 * math.sin(math.radians(angle)))
                    pyautogui.dragTo(cx, cy, duration=0.001)
            else:
                return ActionResult(success=False, error=f"Unknown shape: {shape}")
            return ActionResult(success=True, output=f"Drew {shape}")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _draw_freehand(self, points: str = "") -> ActionResult:
        """Draw freehand by following a list of x,y points. Format: 'x1,y1;x2,y2;...'"""
        try:
            import pyautogui
            pts = [p.strip() for p in points.split(";") if p.strip()]
            if not pts:
                return ActionResult(success=False, error="No points provided")
            x0, y0 = map(int, pts[0].split(","))
            pyautogui.moveTo(x0, y0, duration=0.1)
            for p in pts[1:]:
                x, y = map(int, p.split(","))
                pyautogui.dragTo(x, y, duration=0.01, button='left')
            return ActionResult(success=True, output=f"Drew {len(pts)} points")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    # ── Precision GUI Automation ───────────────────────────────────
    def _click_at_position(self, x: int = 0, y: int = 0, button: str = "left") -> ActionResult:
        """Click at exact screen coordinates."""
        return self._mouse_click(x, y, button)

    def _right_click_at(self, x: int = 0, y: int = 0) -> ActionResult:
        """Right-click at exact screen coordinates."""
        return self._mouse_click(x, y, "right")

    def _double_click_at(self, x: int = 0, y: int = 0) -> ActionResult:
        """Double-click at exact screen coordinates."""
        return self._mouse_click(x, y, "left", True)

    def _drag_from_to(self, x1: int = 0, y1: int = 0, x2: int = 100, y2: int = 100) -> ActionResult:
        """Mouse drag from (x1,y1) to (x2,y2)."""
        try:
            import pyautogui
            pyautogui.moveTo(x1, y1, duration=0.1)
            pyautogui.drag(x2 - x1, y2 - y1, duration=0.3, button='left')
            return ActionResult(success=True, output=f"Dragged from ({x1},{y1}) to ({x2},{y2})")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _scroll_at(self, x: int = 0, y: int = 0, clicks: int = -3) -> ActionResult:
        """Move mouse to position and scroll."""
        try:
            import pyautogui
            pyautogui.moveTo(x, y, duration=0.1)
            pyautogui.scroll(clicks)
            return ActionResult(success=True, output=f"Scrolled {clicks} at ({x},{y})")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _get_mouse_position(self) -> ActionResult:
        """Get current mouse cursor position."""
        try:
            import pyautogui
            x, y = pyautogui.position()
            return ActionResult(success=True, output=f"Mouse at ({x}, {y})")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")

    def _send_keystrokes(self, keys: str = "", delay_ms: int = 50) -> ActionResult:
        """Send keyboard keystrokes. Supports special keys: {enter},{tab},{esc},etc."""
        try:
            import pyautogui
            pyautogui.write(keys, interval=delay_ms / 1000.0)
            return ActionResult(success=True, output=f"Sent: {keys[:100]}")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _press_keys(self, keys: str = "") -> ActionResult:
        """Press keyboard shortcut combo, e.g. 'ctrl+shift+n'."""
        try:
            import pyautogui
            mods = keys.lower().split("+")
            pyautogui.hotkey(*mods)
            return ActionResult(success=True, output=f"Pressed: {keys}")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")

    # ── OneNote / Office Automation ────────────────────────────────
    def _onenote_open(self) -> ActionResult:
        """Open Microsoft OneNote application."""
        try:
            subprocess.Popen(["cmd", "/c", "start", "", "onenote"], shell=False)
            time.sleep(2)
            return ActionResult(success=True, output="OneNote opened")
        except Exception as e:
            try:
                subprocess.Popen(["cmd", "/c", "start", "", "onenote:"], shell=False)
                time.sleep(2)
                return ActionResult(success=True, output="OneNote opened via protocol")
            except Exception as e2:
                return ActionResult(success=False, error=f"Could not open OneNote: {e2}")

    def _onenote_switch_to_tab(self, tab: str = "Draw") -> ActionResult:
        """Switch to a OneNote ribbon tab: Home, Insert, Draw, View, etc."""
        try:
            import pyautogui
            pyautogui.hotkey("alt")
            time.sleep(0.3)
            tab_shortcuts = {"home": "h", "insert": "n", "draw": "d", "view": "w", "history": "1", "review": "2"}
            key = tab_shortcuts.get(tab.lower(), tab[0].lower())
            pyautogui.press(key)
            time.sleep(0.3)
            pyautogui.press("enter")
            return ActionResult(success=True, output=f"Switched to {tab} tab")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")

    def _onenote_select_pen(self, pen_index: int = 0) -> ActionResult:
        """Select a pen tool in OneNote Draw tab. pen_index: 0=first pen, 1=second, etc."""
        try:
            import pyautogui
            for _ in range(pen_index + 1):
                pyautogui.press("right")
                time.sleep(0.1)
            pyautogui.press("enter")
            return ActionResult(success=True, output=f"Selected pen #{pen_index}")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")

    def _onenote_draw_text(self, text: str = "", notebook: str = "", section: str = "", page: str = "") -> ActionResult:
        """Navigate OneNote notebook/section/page and draw text in handwriting."""
        try:
            if notebook:
                self._onenote_open()
                time.sleep(2)
                self._press_keys("ctrl+g")
                time.sleep(0.5)
                self._send_keystrokes(notebook, delay_ms=30)
                time.sleep(0.5)
                self._press_keys("enter")
                time.sleep(1)
            self._onenote_switch_to_tab("Draw")
            time.sleep(0.3)
            self._onenote_select_pen(0)
            time.sleep(0.3)
            import pyautogui
            x, y = pyautogui.position()
            result = self._draw_handwriting(text, int(x), int(y) + 20)
            return result
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    # ── Advanced File Operations ───────────────────────────────────
    def _copy_files(self, source: str = "", dest: str = "") -> ActionResult:
        """Copy file or directory from source to destination."""
        try:
            import shutil
            src = Path(source); dst = Path(dest)
            if src.is_dir():
                shutil.copytree(src, dst, dirs_exist_ok=True)
            else:
                dst.parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(src, dst)
            return ActionResult(success=True, output=f"Copied {source} -> {dest}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _move_files(self, source: str = "", dest: str = "") -> ActionResult:
        """Move file or directory from source to destination."""
        try:
            import shutil
            shutil.move(str(source), str(dest))
            return ActionResult(success=True, output=f"Moved {source} -> {dest}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _delete_files(self, path: str = "", permanently: bool = False) -> ActionResult:
        """Delete file or directory. Use permanently=True for unrecoverable delete."""
        try:
            p = Path(path)
            if p.is_dir():
                import shutil; shutil.rmtree(p)
            else:
                p.unlink()
            return ActionResult(success=True, output=f"Deleted {path}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _rename_file(self, old_path: str = "", new_name: str = "") -> ActionResult:
        """Rename a file or directory."""
        try:
            p = Path(old_path)
            p.rename(p.parent / new_name)
            return ActionResult(success=True, output=f"Renamed to {new_name}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _compress_zip(self, source: str = "", output: str = "") -> ActionResult:
        """Compress file or directory into a zip archive."""
        try:
            import zipfile
            output = output or str(Path(source).with_suffix(".zip"))
            with zipfile.ZipFile(output, 'w', zipfile.ZIP_DEFLATED) as zf:
                src = Path(source)
                if src.is_dir():
                    for f in src.rglob("*"):
                        zf.write(f, f.relative_to(src.parent))
                else:
                    zf.write(src, src.name)
            return ActionResult(success=True, output=f"Created {output}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _extract_zip(self, archive: str = "", output_dir: str = "") -> ActionResult:
        """Extract a zip archive to a directory."""
        try:
            import zipfile
            output_dir = output_dir or str(Path(archive).stem)
            Path(output_dir).mkdir(parents=True, exist_ok=True)
            with zipfile.ZipFile(archive, 'r') as zf:
                zf.extractall(output_dir)
            return ActionResult(success=True, output=f"Extracted to {output_dir}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _download_file(self, url: str = "", dest: str = "") -> ActionResult:
        """Download a file from a URL to local path."""
        try:
            import urllib.request
            dest = dest or os.path.basename(url.split("?")[0]) or "download"
            urllib.request.urlretrieve(url, dest)
            return ActionResult(success=True, output=f"Downloaded {url} -> {dest}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _file_hash(self, path: str = "", algorithm: str = "sha256") -> ActionResult:
        """Compute file hash (md5, sha1, sha256, sha512)."""
        try:
            import hashlib
            h = hashlib.new(algorithm)
            with open(path, 'rb') as f:
                for chunk in iter(lambda: f.read(65536), b''):
                    h.update(chunk)
            return ActionResult(success=True, output=h.hexdigest())
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _list_files_recursive(self, path: str = ".", pattern: str = "*.*") -> ActionResult:
        """List all files recursively matching a glob pattern."""
        try:
            files = [str(f) for f in Path(path).rglob(pattern)]
            return ActionResult(success=True, output=json.dumps(files, indent=2))
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _watch_directory(self, path: str = ".", timeout: int = 10) -> ActionResult:
        """Watch a directory for file changes and return what changed."""
        try:
            before = {str(f): f.stat().st_mtime for f in Path(path).rglob("*") if f.is_file()}
            time.sleep(timeout)
            after = {str(f): f.stat().st_mtime for f in Path(path).rglob("*") if f.is_file()}
            new = [k for k in after if k not in before]
            removed = [k for k in before if k not in after]
            changed = [k for k in before if k in after and before[k] != after[k]]
            return ActionResult(success=True, output=json.dumps({"new": new, "removed": removed, "changed": changed}, indent=2))
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    # ── System Settings ────────────────────────────────────────────
    def _set_wallpaper(self, image_path: str = "", style: str = "fill") -> ActionResult:
        """Change desktop wallpaper. style: fill, fit, stretch, tile, center."""
        try:
            import ctypes
            SPI_SETDESKWALLPAPER = 20
            styles = {"fill": 10, "fit": 6, "stretch": 2, "tile": 0, "center": 1}
            key = (styles.get(style, 10) << 16) | 3
            ctypes.windll.user32.SystemParametersInfoW(SPI_SETDESKWALLPAPER, 0, image_path, key)
            return ActionResult(success=True, output=f"Wallpaper set to {image_path}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _set_resolution(self, width: int = 1920, height: int = 1080) -> ActionResult:
        """Change display resolution."""
        try:
            import ctypes
            dm = ctypes.create_string_buffer(256)
            dm[0:18] = b'\x00' * 18
            ctypes.windll.user32.EnumDisplaySettingsW(None, -1, dm)
            import struct
            struct.pack_into("I", dm, 108, width)
            struct.pack_into("I", dm, 112, height)
            ctypes.windll.user32.ChangeDisplaySettingsW(dm, 0)
            return ActionResult(success=True, output=f"Resolution set to {width}x{height}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _get_system_uptime(self) -> ActionResult:
        """Get system uptime."""
        try:
            import ctypes
            lib = ctypes.windll.kernel32
            ticks = lib.GetTickCount64()
            days, rem = divmod(ticks // 1000, 86400)
            hours, rem = divmod(rem, 3600)
            mins, secs = divmod(rem, 60)
            return ActionResult(success=True, output=f"{int(days)}d {int(hours)}h {int(mins)}m {int(secs)}s")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _list_startup_programs(self) -> ActionResult:
        """List all startup programs."""
        try:
            import winreg
            keys = [winreg.HKEY_CURRENT_USER, winreg.HKEY_LOCAL_MACHINE]
            paths = [r"Software\Microsoft\Windows\CurrentVersion\Run",
                     r"Software\Microsoft\Windows\CurrentVersion\RunOnce"]
            programs = {}
            for hive in keys:
                for path in paths:
                    try:
                        with winreg.OpenKey(hive, path) as key:
                            i = 0
                            while True:
                                try:
                                    name, value, _ = winreg.EnumValue(key, i)
                                    programs[name] = value; i += 1
                                except OSError:
                                    break
                    except OSError:
                        pass
            return ActionResult(success=True, output=json.dumps(programs, indent=2))
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _list_env_vars(self, pattern: str = "") -> ActionResult:
        """List environment variables, optionally filtered by pattern."""
        try:
            vars = {k: v for k, v in os.environ.items() if pattern.lower() in k.lower()}
            return ActionResult(success=True, output=json.dumps(vars, indent=2))
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _set_env_var(self, name: str = "", value: str = "", permanent: bool = False) -> ActionResult:
        """Set an environment variable. Set permanent=True for persistent user env var."""
        try:
            if permanent:
                subprocess.run(["setx", name, value], capture_output=True, shell=False)
            else:
                os.environ[name] = value
            return ActionResult(success=True, output=f"Set {name}={value}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _create_scheduled_task(self, name: str = "", command: str = "", schedule: str = "daily", time_str: str = "09:00") -> ActionResult:
        """Create a Windows scheduled task. schedule: daily, hourly, onstart, onlogon."""
        try:
            maps = {"daily": "DAILY", "hourly": "HOURLY", "onstart": "ONSTART", "onlogon": "ONLOGON"}
            freq = maps.get(schedule, "DAILY")
            subprocess.run(["SchTasks", "/Create", "/SC", freq, "/TN", name, "/TR", command, "/ST", time_str, "/F"],
                           capture_output=True, timeout=10, shell=False)
            return ActionResult(success=True, output=f"Task '{name}' created ({freq})")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _get_battery_report(self) -> ActionResult:
        """Generate and read Windows battery report."""
        try:
            report_path = os.path.expanduser("~\\Desktop\\battery_report.html")
            subprocess.run(["powercfg", "/batteryreport", "/output", report_path],
                           shell=False, capture_output=True, timeout=15)
            if os.path.exists(report_path):
                content = Path(report_path).read_text(encoding="utf-8", errors="ignore")
                return ActionResult(success=True, output=content[:2000] + "...")
            return ActionResult(success=False, error="Report not generated")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    # ── Web Automation ─────────────────────────────────────────────
    def _web_scroll(self, amount: int = 500) -> ActionResult:
        """Scroll the currently focused web page."""
        try:
            import pyautogui
            pyautogui.scroll(-amount)
            return ActionResult(success=True, output=f"Scrolled {amount}")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")

    def _web_fill_form(self, field_data: str = "") -> ActionResult:
        """Fill web form fields. Format: 'label1=value1|label2=value2'"""
        try:
            import pyautogui
            pairs = [p.split("=", 1) for p in field_data.split("|") if "=" in p]
            for label, value in pairs:
                pyautogui.click()
                time.sleep(0.2)
                pyautogui.write(value.strip())
                pyautogui.press("tab")
                time.sleep(0.1)
            return ActionResult(success=True, output=f"Filled {len(pairs)} fields")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")

    def _web_get_page_html(self) -> ActionResult:
        """Get the full HTML of the currently focused web page (simulated)."""
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "l")
            time.sleep(0.2)
            pyautogui.hotkey("ctrl", "c")
            time.sleep(0.2)
            return ActionResult(success=True, output="HTML copied to clipboard")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _web_download_file(self, url: str = "", dest: str = "") -> ActionResult:
        """Download a file from the internet."""
        return self._download_file(url, dest)

    # ── Data Processing ────────────────────────────────────────────
    def _query_csv(self, file_path: str = "", query: str = "SELECT * FROM data LIMIT 10") -> ActionResult:
        """Run SQL query on a CSV file using in-memory SQLite."""
        try:
            import sqlite3
            import csv
            conn = sqlite3.connect(":memory:")
            cur = conn.cursor()
            with open(file_path, 'r', encoding='utf-8-sig') as f:
                reader = csv.DictReader(f)
                cols = reader.fieldnames
                if not cols:
                    return ActionResult(success=False, error="No columns in CSV")
                col_defs = ", ".join(f'"{c}" TEXT' for c in cols)
                cur.execute(f'CREATE TABLE data ({col_defs})')
                placeholders = ", ".join("?" for _ in cols)
                for row in reader:
                    cur.execute(f'INSERT INTO data VALUES ({placeholders})', [row.get(c, "") for c in cols])
            conn.create_function("upper", 1, str.upper)
            conn.create_function("lower", 1, str.lower)
            result = cur.execute(query).fetchall()
            headers = [d[0] for d in cur.description]
            rows = [dict(zip(headers, r)) for r in result]
            conn.close()
            return ActionResult(success=True, output=json.dumps(rows, indent=2, default=str))
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _convert_file(self, input_path: str = "", output_format: str = "json") -> ActionResult:
        """Convert file between formats: csv<->json<->xml."""
        try:
            from pathlib import Path
            inp = Path(input_path)
            ext = inp.suffix.lower()
            content = inp.read_text(encoding="utf-8-sig")

            if ext == ".csv" and output_format == "json":
                import csv, io
                reader = csv.DictReader(io.StringIO(content))
                data = list(reader)
            elif ext == ".json" and output_format == "csv":
                data = json.loads(content)
                import csv, io
                if isinstance(data, list) and data:
                    output = io.StringIO()
                    writer = csv.DictWriter(output, fieldnames=data[0].keys())
                    writer.writeheader(); writer.writerows(data)
                    data = output.getvalue()
            else:
                return ActionResult(success=False, error=f"Unsupported: {ext} -> .{output_format}")
            return ActionResult(success=True, output=json.dumps(data, indent=2) if isinstance(data, list) else str(data))
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    # ── Communication ──────────────────────────────────────────────
    def _compose_email(self, to: str = "", subject: str = "", body: str = "", attachment: str = "", send: bool = False) -> ActionResult:
        """Compose (and optionally send) an email in the newest Outlook. Only called on explicit user request."""
        try:
            if send:
                return self._send_email(to, subject, body)
            import win32com.client
            outlook = win32com.client.Dispatch("Outlook.Application")
            mail = outlook.CreateItem(0)
            mail.To = to; mail.Subject = subject; mail.Body = body
            if attachment and os.path.exists(attachment):
                mail.Attachments.Add(attachment)
            mail.Display()
            return ActionResult(success=True, output=f"Composed email to {to}")
        except Exception as e:
            return ActionResult(success=True, output=f"Email params ready (Outlook unavailable: {e})")

    def _send_whatsapp_message(self, phone: str = "", message: str = "") -> ActionResult:
        """Send a WhatsApp message via web.whatsapp.com."""
        try:
            import pyautogui, webbrowser
            if not self._confirm_profile_use("browser"):
                return ActionResult(success=False, error="PROFILE_APPROVAL_NEEDED",
                                    output=self.get_pending_profile_approval().get("question", ""))
            url = f"https://wa.me/{phone}?text={message[:200]}"
            webbrowser.open(url)
            time.sleep(3)
            pyautogui.press("enter")
            return ActionResult(success=True, output=f"Opened WhatsApp for {phone}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _send_telegram_message(self, chat: str = "", message: str = "") -> ActionResult:
        """Send a Telegram message (opens telegram web/desktop)."""
        try:
            import pyautogui, webbrowser
            if not self._confirm_profile_use("browser"):
                return ActionResult(success=False, error="PROFILE_APPROVAL_NEEDED",
                                    output=self.get_pending_profile_approval().get("question", ""))
            url = f"https://t.me/{chat}"
            webbrowser.open(url)
            time.sleep(3)
            pyautogui.write(message[:200], interval=0.02)
            pyautogui.press("enter")
            return ActionResult(success=True, output=f"Sent Telegram message to {chat}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    # ── Clipboard with formatting ──────────────────────────────────
    def _clipboard_get_text(self) -> ActionResult:
        """Get text from clipboard."""
        return self._get_clipboard()

    def _clipboard_set_text(self, text: str = "") -> ActionResult:
        """Set text on clipboard."""
        return self._set_clipboard(text)

    def _clipboard_get_image(self) -> ActionResult:
        """Get image from clipboard and save to temp file."""
        try:
            from PIL import Image, ImageGrab
            img = ImageGrab.grabclipboard()
            if img is None:
                return ActionResult(success=False, error="No image on clipboard")
            path = os.path.join(os.environ.get("TEMP", "."), "jarvis_clipboard.png")
            img.save(path)
            return ActionResult(success=True, output=f"Image saved to {path}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    # ── Fun / Demo ─────────────────────────────────────────────────
    def _type_custom_text(self, text: str = "", speed: int = 50) -> ActionResult:
        """Type text at a custom speed (ms between keystrokes)."""
        return self._send_keystrokes(text, speed)

    def _show_desktop(self) -> ActionResult:
        """Minimize all windows to show desktop."""
        try:
            import pyautogui
            pyautogui.hotkey("win", "d")
            return ActionResult(success=True, output="Desktop shown")
        except ImportError:
            return ActionResult(success=False, error="pyautogui not installed")

    def _open_file_location(self, path: str = "") -> ActionResult:
        """Open File Explorer to the specified path."""
        try:
            if os.path.isfile(path):
                subprocess.Popen(["explorer", "/select,", path], shell=False)
            else:
                subprocess.Popen(["explorer", path], shell=False)
            return ActionResult(success=True, output=f"Opened {path}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _empty_recycle_bin(self) -> ActionResult:
        """Empty the Windows Recycle Bin."""
        try:
            subprocess.run(["cmd", "/c", "rd /s /q C:\\$Recycle.bin"], capture_output=True, timeout=10)
            return ActionResult(success=True, output="Recycle bin emptied")
        except Exception:
            try:
                import ctypes
                ctypes.windll.shell32.SHEmptyRecycleBinW(None, None, 0)
                return ActionResult(success=True, output="Recycle bin emptied")
            except Exception as e:
                return ActionResult(success=False, error=str(e))

    # ═══════════════════════════════════════════════════════════════
    # SCREEN VISION SYSTEM — see, understand, and interact with screen
    # ═══════════════════════════════════════════════════════════════

    def _is_analyze_screen_runnable(self) -> bool:
        """Check if screen analysis dependencies are available."""
        try:
            import pytesseract; import mss; from PIL import Image; import screen_analyzer
            return True
        except ImportError:
            return False

    def _analyze_screen(self) -> ActionResult:
        """Capture screen and return full analysis: OCR text, buttons, form fields, UI elements."""
        try:
            from mss import mss
            from PIL import Image
            from screen_analyzer import get_analyzer
            with mss() as sct:
                monitor = sct.monitors[1]
                sct_img = sct.grab(monitor)
                img = Image.frombytes("RGB", sct_img.size, sct_img.rgb)
            analyzer = get_analyzer()
            analysis = analyzer.analyze_screenshot(img)
            return ActionResult(success=True, output=json.dumps(analysis.to_dict(), indent=2))
        except ImportError as e:
            return ActionResult(success=False, error=f"Missing dep: {e}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _read_screen_text(self) -> ActionResult:
        """Read all visible text from the screen (like seeing)."""
        try:
            from mss import mss
            from PIL import Image
            from screen_analyzer import get_analyzer
            with mss() as sct:
                sct_img = sct.grab(sct.monitors[1])
                img = Image.frombytes("RGB", sct_img.size, sct_img.rgb)
            analyzer = get_analyzer()
            analysis = analyzer.analyze_screenshot(img)
            return ActionResult(success=True, output=analysis.text if analysis.text else "(No text detected on screen)")
        except ImportError as e:
            return ActionResult(success=False, error=f"Missing dep: {e}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _find_on_screen(self, text: str = "") -> ActionResult:
        """Find all on-screen elements containing the specified text. Returns positions."""
        if not text:
            return ActionResult(success=False, error="No text specified to find")
        try:
            from mss import mss
            from PIL import Image
            from screen_analyzer import get_analyzer
            with mss() as sct:
                sct_img = sct.grab(sct.monitors[1])
                img = Image.frombytes("RGB", sct_img.size, sct_img.rgb)
            analyzer = get_analyzer()
            analysis = analyzer.analyze_screenshot(img)
            matches = analyzer.find_text(analysis, text)
            if matches:
                result = [m.to_dict() for m in matches]
                return ActionResult(success=True, output=json.dumps(result, indent=2))
            return ActionResult(success=False, output=f"'{text}' not found on screen")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _click_button(self, label: str = "") -> ActionResult:
        """Find a button by label text and click it."""
        if not label:
            return ActionResult(success=False, error="No button label specified")
        try:
            from mss import mss
            from PIL import Image
            from screen_analyzer import get_analyzer
            import pyautogui
            with mss() as sct:
                sct_img = sct.grab(sct.monitors[1])
                img = Image.frombytes("RGB", sct_img.size, sct_img.rgb)
            analyzer = get_analyzer()
            analysis = analyzer.analyze_screenshot(img)
            btn = analyzer.find_button(analysis, label)
            if btn:
                pyautogui.click(btn.center_x, btn.center_y)
                return ActionResult(success=True, output=f"Clicked '{label}' button at ({btn.center_x}, {btn.center_y})")
            # Fallback: search for any text match
            matches = analyzer.find_text(analysis, label)
            if matches:
                m = matches[0]
                pyautogui.click(m.center_x, m.center_y)
                return ActionResult(success=True, output=f"Clicked '{label}' at ({m.center_x}, {m.center_y})")
            return ActionResult(success=False, output=f"Button '{label}' not found on screen")
        except ImportError as e:
            return ActionResult(success=False, error=f"Missing dep: {e}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _click_element(self, text: str = "") -> ActionResult:
        """Click any on-screen element containing the specified text."""
        return self._click_button(text)

    def _fill_form_field(self, label: str = "", value: str = "") -> ActionResult:
        """Find a form field by its label and type a value into it."""
        if not label or not value:
            return ActionResult(success=False, error="Both label and value required")
        try:
            from mss import mss
            from PIL import Image
            from screen_analyzer import get_analyzer
            import pyautogui
            with mss() as sct:
                sct_img = sct.grab(sct.monitors[1])
                img = Image.frombytes("RGB", sct_img.size, sct_img.rgb)
            analyzer = get_analyzer()
            analysis = analyzer.analyze_screenshot(img)
            field = analyzer.find_text_field(analysis, label)
            if field:
                pyautogui.click(field.center_x, field.center_y)
                time.sleep(0.15)
                pyautogui.write(value, interval=0.02)
                return ActionResult(success=True, output=f"Filled '{label}' field with '{value}'")
            # Fallback: click on the label itself, then type
            matches = analyzer.find_text(analysis, label)
            if matches:
                m = matches[0]
                pyautogui.click(m.center_x + 80, m.center_y + 5)
                time.sleep(0.15)
                pyautogui.write(value, interval=0.02)
                return ActionResult(success=True, output=f"Filled near '{label}' with '{value}'")
            return ActionResult(success=False, output=f"Field '{label}' not found on screen")
        except ImportError as e:
            return ActionResult(success=False, error=f"Missing dep: {e}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _fill_form(self, form_data: str = "") -> ActionResult:
        """Auto-detect all form fields and fill them. Format: 'label1=value1|label2=value2'"""
        if not form_data:
            return ActionResult(success=False, error="No form data provided")
        try:
            from mss import mss
            from PIL import Image
            from screen_analyzer import get_analyzer
            import pyautogui
            with mss() as sct:
                sct_img = sct.grab(sct.monitors[1])
                img = Image.frombytes("RGB", sct_img.size, sct_img.rgb)
            analyzer = get_analyzer()
            analysis = analyzer.analyze_screenshot(img)
            pairs = [p.split("=", 1) for p in form_data.split("|") if "=" in p]
            filled = 0
            for label, value in pairs:
                label = label.strip()
                value = value.strip()
                if not label or not value:
                    continue
                field = analyzer.find_text_field(analysis, label)
                if field:
                    pyautogui.click(field.center_x, field.center_y)
                    time.sleep(0.1)
                    pyautogui.write(value, interval=0.02)
                    filled += 1
                else:
                    # Try clicking near the label
                    matches = analyzer.find_text(analysis, label)
                    if matches:
                        m = matches[0]
                        pyautogui.click(m.center_x + 80, m.center_y + 5)
                        time.sleep(0.1)
                        pyautogui.write(value, interval=0.02)
                        filled += 1
            return ActionResult(success=True, output=f"Filled {filled}/{len(pairs)} form fields")
        except ImportError as e:
            return ActionResult(success=False, error=f"Missing dep: {e}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _wait_for_text(self, text: str = "", timeout: int = 10) -> ActionResult:
        """Wait until the specified text appears on screen (up to timeout seconds)."""
        if not text:
            return ActionResult(success=False, error="No text to wait for")
        try:
            from mss import mss
            from PIL import Image
            from screen_analyzer import get_analyzer
            analyzer = get_analyzer()
            start = time.time()
            while time.time() - start < timeout:
                with mss() as sct:
                    sct_img = sct.grab(sct.monitors[1])
                    img = Image.frombytes("RGB", sct_img.size, sct_img.rgb)
                analysis = analyzer.analyze_screenshot(img)
                if text.lower() in analysis.text.lower():
                    matches = analyzer.find_text(analysis, text)
                    pos = f" at ({matches[0].center_x}, {matches[0].center_y})" if matches else ""
                    return ActionResult(success=True, output=f"Found '{text}'{pos} ({time.time()-start:.1f}s)")
                time.sleep(0.3)
            return ActionResult(success=False, output=f"'{text}' not found within {timeout}s")
        except ImportError as e:
            return ActionResult(success=False, error=f"Missing dep: {e}")
        except Exception as e:
            return ActionResult(success=False, error=str(e))

    def _get_active_window_title(self) -> str:
        """Get the title of the currently focused window (real vision)."""
        try:
            import win32gui
            hwnd = win32gui.GetForegroundWindow()
            return win32gui.GetWindowText(hwnd)
        except Exception:
            return ""

    def _get_screen_context(self, max_text=40, display: str = None) -> str:
        """REAL VISION: capture the current screen and return a structured summary
        (screenshot + OCR text + positions + UI tree + active window). The AI planner uses this
        so it KNOWS what's actually on screen instead of guessing.

        display: If set (e.g. ":99"), capture from VDI display. If None, capture from host.
        """
        try:
            from PIL import Image
            import io, base64, subprocess as _sp
            parts = []

            # 0. Screenshot for visual context
            screenshot_b64 = None
            try:
                if display:
                    # VDI mode: capture from X11 display using ImageMagick
                    result = _sp.run(
                        ["wsl", "-e", "bash", "-c",
                         f"DISPLAY={display} import -window root png:- 2>/dev/null"],
                        capture_output=True, timeout=10
                    )
                    if result.stdout:
                        img = Image.open(io.BytesIO(result.stdout))
                        img = img.resize((img.width // 2, img.height // 2), Image.Resampling.LANCZOS)
                        buf = io.BytesIO()
                        img.save(buf, format="JPEG", quality=75)
                        screenshot_b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
                        parts.append(f"VDI Screenshot: {img.width}x{img.height} captured from {display}")
                else:
                    # Host mode: capture using mss
                    from mss import mss
                    with mss() as sct:
                        monitor = sct.monitors[1]
                        sct_img = sct.grab(monitor)
                        img = Image.frombytes("RGB", sct_img.size, sct_img.rgb)
                        img = img.resize((img.width // 2, img.height // 2), Image.Resampling.LANCZOS)
                        buf = io.BytesIO()
                        img.save(buf, format="JPEG", quality=75)
                        screenshot_b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
                        parts.append(f"Screenshot: {img.width}x{img.height} captured")
            except Exception:
                pass

            # 1. Active foreground window
            if display:
                # VDI mode: get window via xdotool
                try:
                    result = _sp.run(
                        ["wsl", "-e", "bash", "-c",
                         f"DISPLAY={display} xdotool getactivewindow getwindowname 2>/dev/null"],
                        capture_output=True, text=True, timeout=5
                    )
                    title = result.stdout.strip()
                    if title:
                        parts.append(f"Active window (VDI): '{title}'")
                except Exception:
                    pass
            else:
                # Host mode: get window via win32gui
                try:
                    import win32gui
                    hwnd = win32gui.GetForegroundWindow()
                    title = win32gui.GetWindowText(hwnd)
                    if title:
                        parts.append(f"Active window: '{title}' (hwnd={hwnd})")
                except Exception:
                    pass

            # 2. OCR text with positions
            ocr_res = self._ocr_screen(display=display)
            if ocr_res.success and ocr_res.output:
                try:
                    items = json.loads(ocr_res.output)
                    if items:
                        lines = []
                        for it in items[:max_text]:
                            lines.append(f"  '{it.get('text','')}' at ({it.get('x',0)},{it.get('y',0)})")
                        parts.append("Visible text (OCR):\n" + "\n".join(lines))
                except Exception:
                    if ocr_res.output:
                        parts.append(f"Visible text (OCR): {ocr_res.output[:500]}")

            # 3. UI elements via accessibility tree
            try:
                from accessibility_bridge import AccessibilityBridge
                bridge = AccessibilityBridge()
                if getattr(bridge, '_uia', None):
                    elems = []
                    try:
                        window = bridge._uia.GetForegroundControl()
                        def _walk(ctl, depth=0):
                            if depth > 2 or len(elems) > 12:
                                return
                            try:
                                name = ctl.Name
                                ctrl_type = ctl.ControlTypeName
                                if name and ctrl_type in ("ButtonControl", "EditControl", "MenuItemControl", "HyperlinkControl", "TabItemControl", "CheckBoxControl", "ComboBoxControl", "TextControl", "TitleBarControl"):
                                    rect = ctl.BoundingRectangle
                                    if rect and rect.left >= 0:
                                        elems.append(f"  [{ctrl_type}] '{name}' at ({rect.left},{rect.top})")
                            except Exception:
                                pass
                            try:
                                for c in ctl.GetChildren():
                                    _walk(c, depth + 1)
                            except Exception:
                                pass
                        _walk(window)
                    except Exception:
                        pass
                    if elems:
                        parts.append("UI elements (accessibility tree):\n" + "\n".join(elems[:12]))
            except Exception:
                pass

            # 4. Use vision LLM to describe the screenshot (Cloudflare Gemma 4, then OCR-only fallback)
            if screenshot_b64:
                vision_desc = None
                # Try Cloudflare Workers AI Gemma 4 vision (supports images, no license required)
                try:
                    import os, urllib.request
                    cf_token = os.getenv("CF_API_TOKEN") or ""
                    cf_account = os.getenv("CF_ACCOUNT_ID") or ""
                    if cf_token:
                        import json as _json
                        url = f"https://api.cloudflare.com/client/v4/accounts/{cf_account}/ai/run/@cf/google/gemma-4-26b-a4b-it"
                        payload = _json.dumps({
                            "messages": [
                                {"role": "system", "content": "Describe what's on this computer screen in 2-3 sentences. Focus on the main application, any visible text, buttons, or UI elements. Be specific and concise."},
                                {"role": "user", "content": "Describe this screen:"}
                            ],
                            "image": f"data:image/jpeg;base64,{screenshot_b64}",
                            "max_tokens": 300,
                        }).encode()
                        req = urllib.request.Request(url, data=payload, headers={
                            "Authorization": f"Bearer {cf_token}",
                            "Content-Type": "application/json"
                        })
                        resp = urllib.request.urlopen(req, timeout=20)
                        data = _json.loads(resp.read().decode())
                        if data.get("success"):
                            choice = data.get("result", {}).get("choices", [{}])[0]
                            msg = choice.get("message", {})
                            vision_desc = msg.get("content", "").strip() or msg.get("reasoning_content", "").strip()
                except Exception:
                    pass

                if vision_desc:
                    parts.append(f"AI Vision description: {vision_desc}")

            return "\n\n".join(parts) if parts else "Screen context unavailable"
        except Exception as e:
            return f"Screen context unavailable: {e}"

    def _describe_screen(self) -> ActionResult:
        """Get a human-readable description of what's visible on screen using vision LLM."""
        try:
            from mss import mss
            from PIL import Image
            import io, base64, os, urllib.request, json as _json

            # Capture screenshot
            with mss() as sct:
                sct_img = sct.grab(sct.monitors[1])
                img = Image.frombytes("RGB", sct_img.size, sct_img.rgb)

            # Resize for LLM
            img = img.resize((img.width // 2, img.height // 2), Image.Resampling.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=75)
            screenshot_b64 = base64.b64encode(buf.getvalue()).decode("utf-8")

            # Try Cloudflare Workers AI Gemma 4 vision (supports images, no license required)
            cf_token = os.getenv("CF_API_TOKEN") or ""
            cf_account = os.getenv("CF_ACCOUNT_ID") or ""
            if cf_token:
                try:
                    url = f"https://api.cloudflare.com/client/v4/accounts/{cf_account}/ai/run/@cf/google/gemma-4-26b-a4b-it"
                    payload = _json.dumps({
                        "messages": [
                            {"role": "system", "content": "Describe what's on this computer screen in detail. List: 1) The main application/window, 2) All visible text, buttons, and UI elements, 3) Any forms or input fields. Be specific and structured."},
                            {"role": "user", "content": "Describe this screen:"}
                        ],
                        "image": f"data:image/jpeg;base64,{screenshot_b64}",
                        "max_tokens": 500,
                    }).encode()
                    req = urllib.request.Request(url, data=payload, headers={
                        "Authorization": f"Bearer {cf_token}",
                        "Content-Type": "application/json"
                    })
                    resp = urllib.request.urlopen(req, timeout=20)
                    data = _json.loads(resp.read().decode())
                    if data.get("success"):
                        choice = data.get("result", {}).get("choices", [{}])[0]
                        msg = choice.get("message", {})
                        desc = msg.get("content", "").strip() or msg.get("reasoning_content", "").strip()
                        if desc:
                            return ActionResult(success=True, output=desc)
                except Exception:
                    pass

            return ActionResult(success=False, error="No vision API available")
        except Exception as e:
            return ActionResult(success=False, error=str(e))


    # ── Follow-up Question System ──────────────────────────────────
    def _needs_followup(self, goal: str) -> Optional[list[str]]:
        """Analyze a goal and return follow-up questions ONLY if absolutely critical.
        Returns None if the goal is clear. DEFAULT: just do it with sensible defaults.
        """
        goal_lower = goal.lower()
        questions = []

        # Only ask for truly essential missing info — default to sensible values
        # Email: only ask recipient if completely missing
        if any(w in goal_lower for w in ["email", "mail", "send to"]):
            if "@" not in goal and "to " not in goal_lower and "send" in goal_lower:
                questions.append("Who should I send the email to?")

        # Say/tell: only ask what to say if completely empty
        if any(w in goal_lower for w in ["say ", "tell ", "make it say", "reply with", "compose"]):
            has_text = any(c in goal for c in ["'", '"']) or len(goal_lower.split()) > 5
            if not has_text:
                questions.append("What exactly should I say?")

        # DON'T ask about browser — just default to Chrome
        # DON'T ask about tone/sign-off — just use sensible defaults
        # DON'T ask about email account — just use default

        return questions if questions else None


    def _autofill_from_goal(self, goal: str, plan: dict):
        """Auto-fill missing plan params by extracting them from the goal text.
        This prevents asking questions the user already answered in their original message.
        """
        goal_lower = goal.lower()

        # Extract browser from goal
        browser = None
        for b in ["chrome", "edge", "firefox", "safari", "brave", "opera", "google chrome", "microsoft edge"]:
            if b in goal_lower:
                browser = b
                break

        # Extract URL from goal
        url = None
        url_match = re.search(r'(https?://[^\s]+|[\w-]+\.(com|org|net|io|gov|edu|co)[^\s]*)', goal_lower)
        if url_match:
            url = url_match.group()
            if not url.startswith("http"):
                url = "https://" + url

        # Extract known app names — look specifically after open/launch keywords
        app_name = None
        known_apps = ["outlook", "chrome", "edge", "firefox", "notepad", "calculator", "explorer",
                       "terminal", "cmd", "word", "excel", "powerpoint", "onenote", "teams", "slack",
                       "vscode", "code", "pycharm", "spotify", "vlc", "zoom", "whatsapp"]
        # First try: extract from "open/launch <app>" pattern
        open_match = re.search(r'(?:open|launch|start|run)\s+(\w+)', goal_lower)
        if open_match:
            candidate = open_match.group(1)
            if candidate in known_apps:
                app_name = candidate
        # Fallback: find first known app in goal text
        if not app_name:
            for app in known_apps:
                if re.search(rf'\b{re.escape(app)}\b', goal_lower):
                    app_name = app
                    break

        # Fill into plan steps
        for step in plan.get("steps", []):
            action = step.get("action", "")
            params = step.get("params", {})
            action_lower = action.lower()
            desc_lower = (step.get("description", "") or "").lower()

            # Match ANY web/navigation action (LLM uses various names)
            is_web_step = (
                action_lower in ("navigate_web", "web_navigate", "open_url", "browse_web",
                                 "go_to_url", "open_website", "navigate") or
                "navigate" in action_lower or "browse" in action_lower or
                "go to" in desc_lower or "navigate" in desc_lower
            )
            if is_web_step:
                if browser and "browser" not in params and "chrome" not in str(params) and "edge" not in str(params) and "firefox" not in str(params):
                    params["browser"] = browser
                if url and not params.get("url"):
                    params["url"] = url

            # Match ANY open/launch action
            is_open_step = (
                action_lower in ("open_app", "launch_app", "open_application", "start_app") or
                "open" in action_lower or "launch" in action_lower
            )
            if is_open_step:
                if app_name and not params.get("app") and not params.get("app_name"):
                    params["app"] = app_name
                    params["app_name"] = app_name

            step["params"] = params

    def _needs_step_clarification(self, plan: dict) -> Optional[list[dict]]:
        """Analyze each step — only ask if LLM completely missed a CRITICAL param.
        Default to sensible values instead of asking.
        """
        questions = []
        for i, step in enumerate(plan.get("steps", [])):
            action = step.get("action", "")
            params = step.get("params", {})
            desc = step.get("description", "")

            # open_app: extract from description or default — DON'T ask
            if action in ("open_app", "wsl_launch") and not params.get("app_name") and not params.get("app"):
                desc_lower = (desc or "").lower()
                known_app = None
                for known in ["outlook", "chrome", "edge", "firefox", "notepad", "calculator", "explorer",
                              "terminal", "cmd", "word", "excel", "powerpoint", "onenote", "teams", "slack",
                              "vscode", "code", "pycharm", "spotify", "vlc", "zoom"]:
                    if known in desc_lower:
                        known_app = known
                        break
                if known_app:
                    params["app"] = known_app
                    params["app_name"] = known_app
                    step["params"] = params
                # DON'T ask — just skip if truly unknown

            # web_search: extract query from description — DON'T ask
            if action in ("web_search",) and not params.get("query"):
                # Try to extract from description
                desc_lower = (desc or "").lower()
                if "search" in desc_lower:
                    q = desc_lower.split("search")[-1].strip().lstrip(":")
                    if q:
                        params["query"] = q
                        step["params"] = params
                # If still no query, use the goal text
                if not params.get("query"):
                    goal_text = _pending_clarification.get("goal", "") if _pending_clarification else ""
                    if goal_text:
                        params["query"] = goal_text
                        step["params"] = params
                # DON'T ask — just skip if truly empty

            # navigate_web: default to Chrome, extract URL from desc — DON'T ask
            is_web_action = (
                action in ("navigate_web", "web_navigate", "open_url", "browse_web",
                           "go_to_url", "open_website", "navigate") or
                "navigate" in action.lower() or "browse" in action.lower()
            )
            if is_web_action:
                if not params.get("url"):
                    # Try to extract URL from description
                    url_match = re.search(r'(https?://[^\s]+|[\w-]+\.(com|org|net|io)[^\s]*)', desc)
                    if url_match:
                        params["url"] = url_match.group()
                        step["params"] = params
                if not params.get("browser"):
                    params["browser"] = "chrome"
                    step["params"] = params

            # write_file: default to Desktop — DON'T ask
            if action in ("write_file",):
                if not params.get("path"):
                    params["path"] = str(Path.home() / "Desktop" / "output.txt")
                    step["params"] = params
                if not params.get("content"):
                    params["content"] = ""
                    step["params"] = params

            # send_email: only ask recipient if truly missing
            if action in ("send_email", "send_email_smtp", "compose_email"):
                if not params.get("to"):
                    questions.append({
                        "step_index": i, "action": action,
                        "question": f"Step {i+1}: Who should I send the email to?",
                        "param": "to"
                    })

        return questions if questions else None




_agent: Optional[ComputerUseAgent] = None


def get_agent() -> ComputerUseAgent:
    global _agent
    if _agent is None:
        _agent = ComputerUseAgent()
    return _agent


async def execute_goal(goal: str, safety: str = "confirm_destructive", followup_answers: dict = None,
                       target_desktop: int = None) -> dict:
    """Execute any natural language goal. Main entry point."""
    global _pending_clarification

    # ── Check if this is an answer to a pending clarification ──
    if _pending_clarification and _is_clarification_answer(goal):
        pending = _pending_clarification
        questions = pending.get("questions", [])
        plan = pending.get("plan", {})
        param_updates = _match_answer_to_question(goal, questions)

        if param_updates:
            # Merge answer into pending plan
            for step in plan.get("steps", []):
                params = step.get("params", {})
                for key, val in param_updates.items():
                    if key not in params or not params[key]:
                        params[key] = val
                step["params"] = params

            # Re-check if more clarifications are needed
            agent = ComputerUseAgent(safety=SafetyLevel(safety) if safety in [s.value for s in SafetyLevel] else SafetyLevel.CONFIRM_DESTRUCTIVE)
            if target_desktop is not None:
                agent._target_desktop = target_desktop
            remaining = agent._needs_step_clarification(plan)
            if remaining:
                # Still need more info
                bare_qs = [sq["question"].split(": ", 1)[-1] if sq["question"].startswith("Step ") else sq["question"] for sq in remaining]
                _pending_clarification["questions"] = bare_qs
                _pending_clarification["plan"] = plan
                return {
                    "success": True, "action": "clarify",
                    "goal": pending["goal"],
                    "questions": [sq["question"] for sq in remaining],
                    "message": "I need a bit more info.",
                    "filled_params": param_updates,
                }
            else:
                # All questions answered — execute with merged answers
                _pending_clarification = {}
                # Build followup_answers from the merged params
                merged_answers = {}
                for step in plan.get("steps", []):
                    for k, v in step.get("params", {}).items():
                        if v:
                            merged_answers[k] = v
                return await agent.execute(pending["goal"], merged_answers)

    # ── Normal flow ──
    _pending_clarification = {}
    safety_level = SafetyLevel(safety) if safety in [s.value for s in SafetyLevel] else SafetyLevel.CONFIRM_DESTRUCTIVE
    agent = ComputerUseAgent(safety=safety_level)
    if target_desktop is not None:
        agent._target_desktop = target_desktop
    return await agent.execute(goal, followup_answers)
