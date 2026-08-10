"""JARVIS Artifact Engine — Generate and validate real files.

JARVIS doesn't just output text — it creates real, usable files:
- Documents: .docx, .pdf, .txt, .md
- Spreadsheets: .xlsx, .csv
- Presentations: .pptx
- Images: .png, .jpg, .svg
- Code: .py, .js, .ts, .html, .css
- Archives: .zip
- Video: .mp4 (via ffmpeg)
- Audio: .wav (via ffmpeg)
"""

import os
import json
import time
import hashlib
import logging
import subprocess
from typing import Optional, Dict, List, Any
from dataclasses import dataclass, field
from pathlib import Path

log = logging.getLogger("artifact_engine")

ARTIFACTS_DIR = Path.home() / ".jarvis" / "artifacts"
ARTIFACTS_DIR.mkdir(parents=True, exist_ok=True)


@dataclass
class Artifact:
    id: str
    name: str
    path: str
    type: str
    size_bytes: int = 0
    mime_type: str = ""
    checksum: str = ""
    created_at: float = 0
    validated: bool = False
    validation_results: Dict[str, Any] = field(default_factory=dict)
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict:
        return {
            "id": self.id, "name": self.name, "path": self.path,
            "type": self.type, "size_bytes": self.size_bytes,
            "mime_type": self.mime_type, "checksum": self.checksum,
            "created_at": self.created_at, "validated": self.validated,
            "validation_results": self.validation_results,
            "metadata": self.metadata,
        }


class ArtifactEngine:
    """Generates and validates real files."""

    def __init__(self):
        self._artifacts: Dict[str, Artifact] = {}
        self._load_index()

    def _load_index(self):
        """Load artifact index from disk."""
        index_file = ARTIFACTS_DIR / "index.json"
        if index_file.exists():
            try:
                data = json.loads(index_file.read_text())
                for item in data:
                    artifact = Artifact(**item)
                    self._artifacts[artifact.id] = artifact
            except Exception as e:
                log.error(f"[ARTIFACT] Index load failed: {e}")

    def _save_index(self):
        """Save artifact index to disk."""
        index_file = ARTIFACTS_DIR / "index.json"
        data = [a.to_dict() for a in self._artifacts.values()]
        index_file.write_text(json.dumps(data, indent=2))

    def _register_artifact(self, file_path: Path, artifact_type: str) -> Artifact:
        """Register a new artifact."""
        import uuid
        stat = file_path.stat()
        sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                sha256.update(chunk)

        artifact = Artifact(
            id=str(uuid.uuid4())[:8],
            name=file_path.stem,
            path=str(file_path),
            type=artifact_type,
            size_bytes=stat.st_size,
            mime_type=self._guess_mime(file_path.suffix),
            checksum=sha256.hexdigest(),
            created_at=time.time(),
        )
        self._artifacts[artifact.id] = artifact
        self._save_index()
        return artifact

    def _guess_mime(self, extension: str) -> str:
        """Guess MIME type from file extension."""
        mime_map = {
            ".txt": "text/plain", ".md": "text/markdown", ".json": "application/json",
            ".csv": "text/csv", ".py": "text/x-python", ".js": "application/javascript",
            ".html": "text/html", ".css": "text/css", ".pdf": "application/pdf",
            ".png": "image/png", ".jpg": "image/jpeg", ".svg": "image/svg+xml",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".zip": "application/zip", ".mp4": "video/mp4", ".wav": "audio/wav",
        }
        return mime_map.get(extension.lower(), "application/octet-stream")

    def create_text(self, name: str, content: str, ext: str = "txt") -> Artifact:
        """Create a plain text file."""
        file_path = ARTIFACTS_DIR / f"{name}.{ext}"
        file_path.write_text(content, encoding="utf-8")
        return self._register_artifact(file_path, "document")

    def create_markdown(self, name: str, content: str) -> Artifact:
        """Create a Markdown file."""
        return self.create_text(name, content, "md")

    def create_json(self, name: str, data: Any) -> Artifact:
        """Create a JSON file."""
        content = json.dumps(data, indent=2, default=str)
        file_path = ARTIFACTS_DIR / f"{name}.json"
        file_path.write_text(content, encoding="utf-8")
        return self._register_artifact(file_path, "data")

    def create_csv(self, name: str, headers: List[str], rows: List[List[Any]]) -> Artifact:
        """Create a CSV file."""
        import csv
        file_path = ARTIFACTS_DIR / f"{name}.csv"
        with open(file_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(headers)
            writer.writerows(rows)
        return self._register_artifact(file_path, "spreadsheet")

    def create_python(self, name: str, code: str) -> Artifact:
        """Create a Python script."""
        return self.create_text(name, code, "py")

    def create_html(self, name: str, content: str) -> Artifact:
        """Create an HTML file."""
        return self.create_text(name, content, "html")

    def create_docx(self, name: str, content: str, title: str = "") -> Artifact:
        """Create a Word document."""
        try:
            from docx import Document
            doc = Document()
            if title:
                doc.add_heading(title, level=1)
            for paragraph in content.split("\n\n"):
                doc.add_paragraph(paragraph)
            file_path = ARTIFACTS_DIR / f"{name}.docx"
            doc.save(str(file_path))
            return self._register_artifact(file_path, "document")
        except ImportError:
            return self.create_markdown(name, content)

    def create_xlsx(self, name: str, data: Dict[str, List[Dict]]) -> Artifact:
        """Create an Excel spreadsheet."""
        try:
            from openpyxl import Workbook
            wb = Workbook()
            for sheet_name, rows in data.items():
                ws = wb.create_sheet(title=sheet_name[:30])
                if rows:
                    headers = list(rows[0].keys())
                    ws.append(headers)
                    for row in rows:
                        ws.append([row.get(h, "") for h in headers])
            if "Sheet" in wb.sheetnames:
                del wb["Sheet"]
            file_path = ARTIFACTS_DIR / f"{name}.xlsx"
            wb.save(str(file_path))
            return self._register_artifact(file_path, "spreadsheet")
        except ImportError:
            return self.create_csv(name, [], [])

    def create_pptx(self, name: str, slides: List[Dict[str, str]]) -> Artifact:
        """Create a PowerPoint presentation."""
        try:
            from pptx import Presentation
            prs = Presentation()
            for slide_data in slides:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = slide_data.get("title", "")
                slide.placeholders[1].text = slide_data.get("content", "")
            file_path = ARTIFACTS_DIR / f"{name}.pptx"
            prs.save(str(file_path))
            return self._register_artifact(file_path, "presentation")
        except ImportError:
            md_content = "\n\n---\n\n".join([f"# {s.get('title', '')}\n\n{s.get('content', '')}" for s in slides])
            return self.create_markdown(name, md_content)

    def create_pdf(self, name: str, content: str, title: str = "") -> Artifact:
        """Create a PDF document."""
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            if title:
                pdf.set_font("Arial", "B", 16)
                pdf.cell(0, 10, title, ln=True, align="C")
                pdf.ln(10)
                pdf.set_font("Arial", size=12)
            for line in content.split("\n"):
                pdf.cell(0, 5, line, ln=True)
            file_path = ARTIFACTS_DIR / f"{name}.pdf"
            pdf.output(str(file_path))
            return self._register_artifact(file_path, "document")
        except ImportError:
            return self.create_text(name, content, "txt")

    def create_image(self, name: str, width: int = 800, height: int = 600,
                    color: str = "#4A90D9") -> Artifact:
        """Create an image (PIL or SVG fallback)."""
        try:
            from PIL import Image, ImageDraw, ImageFont
            img = Image.new("RGB", (width, height), color)
            draw = ImageDraw.Draw(img)
            try:
                font = ImageFont.truetype("arial.ttf", 24)
            except Exception:
                font = ImageFont.load_default()
            draw.text((width // 4, height // 2), name, fill="white", font=font)
            file_path = ARTIFACTS_DIR / f"{name}.png"
            img.save(str(file_path))
            return self._register_artifact(file_path, "image")
        except ImportError:
            svg = f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}"><rect width="100%" height="100%" fill="{color}"/><text x="50%" y="50%" fill="white" text-anchor="middle" font-size="24">{name}</text></svg>'
            return self.create_text(name, svg, "svg")

    def create_zip(self, name: str, files: Dict[str, str]) -> Artifact:
        """Create a ZIP archive."""
        import zipfile
        file_path = ARTIFACTS_DIR / f"{name}.zip"
        with zipfile.ZipFile(str(file_path), "w", zipfile.ZIP_DEFLATED) as zf:
            for filename, content in files.items():
                zf.writestr(filename, content)
        return self._register_artifact(file_path, "archive")

    def create_video(self, name: str) -> Artifact:
        """Create a video file (requires ffmpeg)."""
        file_path = ARTIFACTS_DIR / f"{name}.mp4"
        try:
            subprocess.run([
                "ffmpeg", "-y", "-f", "lavfi", "-i",
                "color=c=blue:s=320x240:d=3",
                "-c:v", "libx264", str(file_path)
            ], capture_output=True, timeout=30)
            return self._register_artifact(file_path, "video")
        except Exception:
            return self.create_text(name, "[Video placeholder - ffmpeg not available]", "txt")

    def create_audio(self, name: str) -> Artifact:
        """Create an audio file (requires ffmpeg)."""
        file_path = ARTIFACTS_DIR / f"{name}.wav"
        try:
            subprocess.run([
                "ffmpeg", "-y", "-f", "lavfi", "-i",
                "sine=frequency=440:duration=3",
                str(file_path)
            ], capture_output=True, timeout=30)
            return self._register_artifact(file_path, "audio")
        except Exception:
            return self.create_text(name, "[Audio placeholder - ffmpeg not available]", "txt")

    def validate(self, artifact_id: str) -> Dict[str, Any]:
        """Validate an artifact exists and is usable."""
        artifact = self._artifacts.get(artifact_id)
        if not artifact:
            return {"valid": False, "error": "Artifact not found"}

        results = {
            "valid": True,
            "exists": False,
            "size": 0,
            "checksum_match": False,
            "format_valid": False,
            "checks": [],
        }

        file_path = Path(artifact.path)
        if file_path.exists():
            results["exists"] = True
            results["size"] = file_path.stat().st_size
            results["checks"].append({"check": "File exists", "passed": True})
        else:
            results["valid"] = False
            results["checks"].append({"check": "File exists", "passed": False})
            return results

        if results["size"] > 0:
            results["checks"].append({"check": "Non-empty file", "passed": True, "detail": f"{results['size']} bytes"})
        else:
            results["valid"] = False
            results["checks"].append({"check": "Non-empty file", "passed": False})
            return results

        # Verify checksum
        sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                sha256.update(chunk)
        actual_checksum = sha256.hexdigest()
        results["checksum_match"] = actual_checksum == artifact.checksum
        results["checks"].append({"check": "Checksum valid", "passed": results["checksum_match"]})

        # Format-specific validation
        ext = file_path.suffix.lower()
        if ext == ".json":
            try:
                json.loads(file_path.read_text())
                results["format_valid"] = True
                results["checks"].append({"check": "Valid JSON", "passed": True})
            except Exception:
                results["format_valid"] = False
                results["checks"].append({"check": "Valid JSON", "passed": False})
        elif ext == ".csv":
            try:
                import csv
                with open(file_path, "r") as f:
                    reader = csv.reader(f)
                    rows = list(reader)
                results["format_valid"] = len(rows) > 0
                results["checks"].append({"check": "Valid CSV", "passed": results["format_valid"], "detail": f"{len(rows)} rows"})
            except Exception:
                results["checks"].append({"check": "Valid CSV", "passed": False})
        elif ext in (".py", ".js", ".ts", ".html", ".css", ".md", ".txt"):
            results["format_valid"] = True
            results["checks"].append({"check": "Text file readable", "passed": True})
        elif ext == ".zip":
            try:
                import zipfile
                with zipfile.ZipFile(str(file_path), "r") as zf:
                    names = zf.namelist()
                results["format_valid"] = True
                results["checks"].append({"check": "Valid ZIP", "passed": True, "detail": f"{len(names)} files"})
            except Exception:
                results["checks"].append({"check": "Valid ZIP", "passed": False})
        elif ext == ".docx":
            try:
                from docx import Document
                doc = Document(str(file_path))
                results["format_valid"] = True
                results["checks"].append({"check": "Valid DOCX", "passed": True, "detail": f"{len(doc.paragraphs)} paragraphs"})
            except Exception:
                results["checks"].append({"check": "Valid DOCX", "passed": False})
        elif ext == ".xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(str(file_path))
                results["format_valid"] = True
                results["checks"].append({"check": "Valid XLSX", "passed": True, "detail": f"{len(wb.sheetnames)} sheets"})
            except Exception:
                results["checks"].append({"check": "Valid XLSX", "passed": False})
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                results["format_valid"] = True
                results["checks"].append({"check": "Valid PPTX", "passed": True, "detail": f"{len(prs.slides)} slides"})
            except Exception:
                results["checks"].append({"check": "Valid PPTX", "passed": False})
        elif ext == ".pdf":
            try:
                with open(file_path, "rb") as f:
                    header = f.read(5)
                results["format_valid"] = header == b"%PDF-"
                results["checks"].append({"check": "Valid PDF", "passed": results["format_valid"]})
            except Exception:
                results["checks"].append({"check": "Valid PDF", "passed": False})
        else:
            results["format_valid"] = True
            results["checks"].append({"check": "Format check skipped", "passed": True})

        artifact.validated = results["valid"]
        artifact.validation_results = results
        self._save_index()
        return results

    def list_artifacts(self, artifact_type: str = None) -> List[dict]:
        """List all artifacts, optionally filtered by type."""
        artifacts = self._artifacts.values()
        if artifact_type:
            artifacts = [a for a in artifacts if a.type == artifact_type]
        return [a.to_dict() for a in sorted(artifacts, key=lambda a: a.created_at, reverse=True)]

    def get_artifact(self, artifact_id: str) -> Optional[Artifact]:
        """Get an artifact by ID."""
        return self._artifacts.get(artifact_id)

    def delete_artifact(self, artifact_id: str) -> bool:
        """Delete an artifact."""
        artifact = self._artifacts.get(artifact_id)
        if not artifact:
            return False
        try:
            Path(artifact.path).unlink(missing_ok=True)
        except Exception:
            pass
        del self._artifacts[artifact_id]
        self._save_index()
        return True

    def get_stats(self) -> Dict[str, Any]:
        """Get artifact statistics."""
        total_size = sum(a.size_bytes for a in self._artifacts.values())
        by_type = {}
        for a in self._artifacts.values():
            by_type[a.type] = by_type.get(a.type, 0) + 1
        return {
            "total_artifacts": len(self._artifacts),
            "total_size_bytes": total_size,
            "total_size_mb": round(total_size / (1024 * 1024), 2),
            "by_type": by_type,
            "validated_count": sum(1 for a in self._artifacts.values() if a.validated),
        }


_engine = None


def get_artifact_engine() -> ArtifactEngine:
    global _engine
    if _engine is None:
        _engine = ArtifactEngine()
    return _engine
