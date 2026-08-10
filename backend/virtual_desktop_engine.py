"""
Virtual Desktop Engine — REAL isolation, not fake hotkey desktops.

3-Tier Architecture:
  Tier 1: Python APIs (python-pptx, openpyxl, python-docx) — instant, no GUI
  Tier 2: COM automation (Word, Excel, Chrome) — works from any desktop
  Tier 3: TrueDesktop Win32 isolation + subprocess — invisible to user

Blender/CAD: Headless scripted execution via subprocess (bpy --background --python)
"""

import subprocess
import time
import os
import sys
import json
import tempfile
import traceback
from typing import Optional
from dataclasses import dataclass


class WorkResult:
    def __init__(self, success: bool, message: str, details: dict = None):
        self.success = success
        self.message = message
        self.details = details or {}

    def to_dict(self):
        return {"success": self.success, "message": self.message, "details": self.details}


def _ps(cmd: str, timeout: float = 15.0) -> str:
    try:
        r = subprocess.run(
            ["powershell", "-NoProfile", "-Command", cmd],
            capture_output=True, text=True, timeout=timeout,
        )
        return r.stdout.strip()
    except Exception:
        return ""


# TIER 1: Python APIs — No GUI needed, instant execution

class Tier1_PythonAPIs:
    """Create documents, spreadsheets, presentations directly via Python libraries."""

    @staticmethod
    def create_powerpoint(title: str, slides: list, save_path: str = "") -> WorkResult:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx"], capture_output=True)
            from pptx import Presentation
            from pptx.util import Inches

        if not save_path:
            save_path = os.path.join(os.path.expanduser("~"), "Desktop", f"{title.replace(' ', '_')}.pptx")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title

        for sc in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sc.get("title", "")
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, line in enumerate(sc.get("content", "").split("\n")):
                if i == 0:
                    tf.text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line

        prs.save(save_path)
        return WorkResult(True, f"Presentation saved: {save_path}", {"path": save_path, "slides": len(slides)})

    @staticmethod
    def create_word(title: str, content: str, save_path: str = "") -> WorkResult:
        try:
            from docx import Document
        except ImportError:
            subprocess.run([sys.executable, "-m", "pip", "install", "python-docx"], capture_output=True)
            from docx import Document

        if not save_path:
            save_path = os.path.join(os.path.expanduser("~"), "Desktop", f"{title.replace(' ', '_')}.docx")

        doc = Document()
        doc.add_heading(title, 0)
        for para in content.split("\n"):
            if para.strip():
                doc.add_paragraph(para)
        doc.save(save_path)
        return WorkResult(True, f"Document saved: {save_path}", {"path": save_path})

    @staticmethod
    def create_excel(title: str, headers: list, rows: list, save_path: str = "") -> WorkResult:
        try:
            from openpyxl import Workbook
        except ImportError:
            subprocess.run([sys.executable, "-m", "pip", "install", "openpyxl"], capture_output=True)
            from openpyxl import Workbook

        if not save_path:
            save_path = os.path.join(os.path.expanduser("~"), "Desktop", f"{title.replace(' ', '_')}.xlsx")

        wb = Workbook()
        ws = wb.active
        ws.title = title[:31]
        if headers:
            ws.append(headers)
        for row in rows:
            ws.append(row)
        wb.save(save_path)
        return WorkResult(True, f"Spreadsheet saved: {save_path}", {"path": save_path, "rows": len(rows)})


# TIER 2: COM Automation — Works from ANY desktop, no GUI needed

class Tier2_COMAutomation:
    """Control Office apps via COM — works regardless of which desktop is active."""

    def word_create(self, title="", content="") -> WorkResult:
        t = title.replace('"', '""')
        c = content.replace('"', '""').replace("\n", "\\n")
        r = _ps(f'try {{ $w = New-Object -ComObject Word.Application; $w.Visible = $true; $doc = $w.Documents.Add(); if ("{t}") {{ $w.Selection.Style = "Heading 1"; $w.Selection.TypeText("{t}"); $w.Selection.TypeParagraph() }}; if ("{c}") {{ $w.Selection.Style = "Normal"; $w.Selection.TypeText("{c}") }}; "OK" }} catch {{ "FAIL:$_" }}')
        return WorkResult("OK" in r, "Word document created" if "OK" in r else f"Failed: {r}")

    def word_save(self, path="") -> WorkResult:
        if not path:
            path = os.path.expanduser("~/Desktop/document.docx")
        r = _ps(f'try {{ $w = [Runtime.InteropServices.Marshal]::GetActiveObject("Word.Application"); $w.ActiveDocument.SaveAs("{path}"); "OK" }} catch {{ "FAIL:$_" }}')
        return WorkResult("OK" in r, f"Saved to {path}" if "OK" in r else f"Failed: {r}")

    def excel_create(self, title="") -> WorkResult:
        t = title.replace('"', '""')
        r = _ps(f'try {{ $e = New-Object -ComObject Excel.Application; $e.Visible = $true; $wb = $e.Workbooks.Add(); if ("{t}") {{ $wb.Worksheets(1).Name = "{t[:31]}" }}; "OK" }} catch {{ "FAIL:$_" }}')
        return WorkResult("OK" in r, "Excel created" if "OK" in r else f"Failed: {r}")

    def excel_cell(self, sheet, row, col, value) -> WorkResult:
        v = value.replace('"', '""')
        r = _ps(f'try {{ $e = [Runtime.InteropServices.Marshal]::GetActiveObject("Excel.Application"); $e.Worksheets({sheet}).Cells({row},{col}).Value = "{v}"; "OK" }} catch {{ "FAIL:$_" }}')
        return WorkResult("OK" in r, "Cell set" if "OK" in r else f"Failed: {r}")

    def excel_save(self, path="") -> WorkResult:
        if not path:
            path = os.path.expanduser("~/Desktop/workbook.xlsx")
        r = _ps(f'try {{ $e = [Runtime.InteropServices.Marshal]::GetActiveObject("Excel.Application"); $e.ActiveWorkbook.SaveAs("{path}"); "OK" }} catch {{ "FAIL:$_" }}')
        return WorkResult("OK" in r, f"Saved to {path}" if "OK" in r else f"Failed: {r}")

    def outlook_send(self, to, subject, body) -> WorkResult:
        t = to.replace('"', '""')
        s = subject.replace('"', '""')
        b = body.replace('"', '""').replace("\n", "\\n")
        r = _ps(f'try {{ $o = New-Object -ComObject Outlook.Application; $m = $o.CreateItem(0); $m.To = "{t}"; $m.Subject = "{s}"; $m.Body = "{b}"; $m.Send(); "OK" }} catch {{ "FAIL:$_" }}')
        return WorkResult("OK" in r, f"Email sent to {to}" if "OK" in r else f"Failed: {r}")

    def chrome_navigate(self, url) -> WorkResult:
        u = url.replace('"', '""')
        r = _ps(f'try {{ $w = [Runtime.InteropServices.Marshal]::GetActiveObject("Chrome.Application"); $w.ActiveWindow.LocationURL = "{u}"; "OK" }} catch {{ Start-Process "chrome.exe" -ArgumentList "{u}"; "OK:launched" }}')
        return WorkResult("OK" in r, f"Navigated to {url}" if "OK" in r else f"Failed: {r}")


# TIER 3: True Desktop Isolation + Subprocess

class Tier3_DesktopIsolation:
    """Launch apps on real isolated Windows desktops via CreateDesktopW."""

    def __init__(self):
        self._td = None

    def _get_td(self):
        if self._td is None:
            from true_desktop import get_true_desktop
            self._td = get_true_desktop()
        return self._td

    def launch_app(self, app_name, args=None, desktop_name="ai_work") -> WorkResult:
        td = self._get_td()
        if desktop_name not in td.list_desktops():
            td.create(desktop_name)
        exe = self._find_exe(app_name)
        if not exe:
            return WorkResult(False, f"App not found: {app_name}")
        proc = td.launch_on(desktop_name, exe, args or [])
        if proc:
            return WorkResult(True, f"Launched {app_name} on desktop '{desktop_name}'", {"pid": proc.pid, "desktop": desktop_name})
        return WorkResult(False, f"Failed to launch {app_name}")

    def browse(self, url, browser="chrome", desktop_name="ai_browse") -> WorkResult:
        exe = self._find_browser(browser)
        if not exe:
            return WorkResult(False, f"Browser not found: {browser}")
        td = self._get_td()
        if desktop_name not in td.list_desktops():
            td.create(desktop_name)
        proc = td.launch_on(desktop_name, exe, [url])
        if proc:
            return WorkResult(True, f"Browsing {url} on isolated desktop", {"url": url, "desktop": desktop_name, "pid": proc.pid})
        return WorkResult(False, "Failed to launch browser")

    def run_headless(self, command, env=None, timeout=300) -> WorkResult:
        try:
            full_env = os.environ.copy()
            if env:
                full_env.update(env)
            r = subprocess.run(command, capture_output=True, text=True, timeout=timeout, env=full_env)
            output = r.stdout + r.stderr
            return WorkResult(r.returncode == 0, output[:2000] if output else "Completed", {"returncode": r.returncode})
        except subprocess.TimeoutExpired:
            return WorkResult(False, f"Timed out after {timeout}s")
        except Exception as e:
            return WorkResult(False, f"Error: {e}")

    def blender_script(self, script_path, background=True) -> WorkResult:
        blender = self._find_blender()
        if not blender:
            return WorkResult(False, "Blender not found")
        cmd = [blender]
        if background:
            cmd.append("--background")
        cmd.extend(["--python", script_path])
        return self.run_headless(cmd, timeout=600)

    def blender_command(self, python_code) -> WorkResult:
        blender = self._find_blender()
        if not blender:
            return WorkResult(False, "Blender not found")
        script = os.path.join(tempfile.gettempdir(), "jarvis_blender_script.py")
        with open(script, "w") as f:
            f.write(python_code)
        return self.blender_script(script)

    def _find_exe(self, name):
        name_lower = name.lower().strip()
        known = {
            "notepad": "notepad.exe", "calculator": "calc.exe",
            "word": "winword.exe", "excel": "excel.exe",
            "powerpoint": "powerpnt.exe", "cmd": "cmd.exe",
            "terminal": "wt.exe", "powershell": "pwsh.exe",
            "chrome": self._find_browser("chrome"),
            "edge": self._find_browser("edge"),
            "blender": self._find_blender(),
        }
        if name_lower in known:
            return known[name_lower]
        if os.path.exists(name):
            return name
        return name_lower

    def _find_browser(self, browser="chrome"):
        paths = {
            "chrome": [r"C:\Program Files\Google\Chrome\Application\chrome.exe", r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"],
            "edge": [r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe", r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"],
        }
        for p in paths.get(browser, []):
            if os.path.exists(p):
                return p
        return None

    def _find_blender(self):
        paths = [
            r"C:\Program Files\Blender Foundation\Blender 4.2\blender.exe",
            r"C:\Program Files\Blender Foundation\Blender 4.1\blender.exe",
            r"C:\Program Files\Blender Foundation\Blender 4.0\blender.exe",
            r"C:\Program Files\Blender Foundation\Blender 3.6\blender.exe",
        ]
        for p in paths:
            if os.path.exists(p):
                return p
        try:
            r = subprocess.run(["where", "blender"], capture_output=True, text=True, timeout=5)
            if r.stdout.strip():
                return r.stdout.strip().split("\n")[0]
        except Exception:
            pass
        return None


# MAIN ENGINE

class VirtualDesktopEngine:
    """The main engine. Routes tasks to the right tier automatically."""

    def __init__(self):
        self.tier1 = Tier1_PythonAPIs()
        self.tier2 = Tier2_COMAutomation()
        self.tier3 = Tier3_DesktopIsolation()
        self._sessions = {}

    def browse(self, url, browser="chrome") -> WorkResult:
        """Open URL in browser on user's PRIMARY desktop (visible to user)."""
        # Try COM first (if Chrome is already open)
        try:
            result = self.tier2.chrome_navigate(url)
            if result.success:
                return result
        except Exception:
            pass
        # Launch browser directly on primary desktop
        exe = self.tier3._find_browser(browser)
        if not exe:
            return WorkResult(False, f"Browser not found: {browser}")
        try:
            proc = subprocess.Popen([exe, url], shell=False)
            return WorkResult(True, f"Browsing {url} (PID {proc.pid})", {"url": url, "pid": proc.pid})
        except Exception as e:
            return WorkResult(False, f"Failed to launch browser: {e}")

    def create_document(self, title, content="", save_path="") -> WorkResult:
        try:
            return self.tier1.create_word(title, content, save_path)
        except Exception:
            pass
        result = self.tier2.word_create(title, content)
        if result.success and save_path:
            self.tier2.word_save(save_path)
        return result

    def create_spreadsheet(self, title, headers=None, rows=None, save_path="") -> WorkResult:
        try:
            return self.tier1.create_excel(title, headers or [], rows or [], save_path)
        except Exception:
            pass
        result = self.tier2.excel_create(title)
        if result.success:
            if headers:
                for ci, h in enumerate(headers):
                    self.tier2.excel_cell(1, 1, ci + 1, str(h))
            for ri, row in enumerate(rows or []):
                for ci, val in enumerate(row):
                    self.tier2.excel_cell(1, ri + 2, ci + 1, str(val))
            if save_path:
                self.tier2.excel_save(save_path)
        return result

    def create_presentation(self, title, slides=None, save_path="") -> WorkResult:
        return self.tier1.create_powerpoint(title, slides or [], save_path)

    def send_email(self, to, subject, body) -> WorkResult:
        return self.tier2.outlook_send(to, subject, body)

    def research(self, query) -> WorkResult:
        import urllib.parse
        url = f"https://www.google.com/search?q={urllib.parse.quote(query)}"
        return self.browse(url)

    def open_app(self, app_name, args=None) -> WorkResult:
        """Launch app on the user's PRIMARY desktop (visible to user)."""
        exe = self.tier3._find_exe(app_name)
        if not exe:
            return WorkResult(False, f"App not found: {app_name}")
        try:
            cmd = [exe] + (args or [])
            proc = subprocess.Popen(cmd, shell=False)
            return WorkResult(True, f"Launched {app_name} (PID {proc.pid})", {"pid": proc.pid})
        except Exception as e:
            return WorkResult(False, f"Failed to launch {app_name}: {e}")

    def run_blender(self, script_or_code) -> WorkResult:
        if os.path.isfile(script_or_code):
            return self.tier3.blender_script(script_or_code)
        return self.tier3.blender_command(script_or_code)

    def run_command(self, command, timeout=300) -> WorkResult:
        return self.tier3.run_headless(command, timeout=timeout)

    def do_task(self, task_description) -> WorkResult:
        """Do ANYTHING the user asks. Uses LLM planner for complex tasks.

        Simple tasks (single action) -> fast regex path
        Complex tasks (multi-step) -> LLM planner with observe-act feedback
        """
        import re
        import urllib.parse
        desc = task_description.lower()

        # Quick single-action detection (fast path)
        is_simple = (
            (desc.startswith("browse ") and " " not in desc[7:].strip().split()[0:1]) or
            (desc.startswith("open ") and len(desc.split()) <= 3) or
            (desc in ["lock", "screenshot", "volume up", "volume down"])
        )

        # For simple single-action tasks, use fast regex path
        if is_simple:
            return self._do_simple_task(task_description)

        # For complex tasks, use the LLM planner
        try:
            from task_planner import get_planner
            planner = get_planner()
            planner._engine = self
            result_text = planner.do_anything(task_description)
            return WorkResult(True, result_text)
        except Exception as e:
            # Fallback to simple path
            return self._do_simple_task(task_description)

    def _do_simple_task(self, task_description) -> WorkResult:
        """Fast regex-based single-action dispatch."""
        import re
        import urllib.parse
        desc = task_description.lower()

        if any(w in desc for w in ["browse", "open website", "go to http", "open http", "open google", "open youtube"]):
            url_match = re.search(r"(https?://\S+|www\.\S+)", task_description)
            if url_match:
                return self.browse(url_match.group(0))
            q = re.sub(r".*?(?:browse|open|go to)\s+", "", task_description, flags=re.I).strip()
            if q:
                return self.browse(f"https://www.google.com/search?q={urllib.parse.quote(q)}")
            return WorkResult(False, "What should I browse?")

        if any(w in desc for w in ["create document", "word", "essay", "report", "letter", "write document"]):
            title = re.search(r"(?:create|write|make)\s+(?:a\s+)?(?:document|doc|word|essay|report|letter)\s+(?:about|on|titled|called)?\s*(.+?)$", task_description, re.I)
            t = title.group(1).strip() if title else "Document"
            content = re.search(r"(?:about|on|containing|with)\s+(.+)$", task_description, re.I)
            c = content.group(1).strip() if content else ""
            return self.create_document(t, c)

        if any(w in desc for w in ["excel", "spreadsheet", "table"]):
            title = re.search(r"(?:create|make|new)\s+(?:excel|spreadsheet|table)\s+(?:called|named)?\s*(.+?)$", task_description, re.I)
            t = title.group(1).strip() if title else "Spreadsheet"
            return self.create_spreadsheet(t)

        if any(w in desc for w in ["powerpoint", "presentation", "ppt", "slides"]):
            title = re.search(r"(?:create|make|new)\s+(?:presentation|ppt|powerpoint)\s+(?:called|named)?\s*(.+?)$", task_description, re.I)
            t = title.group(1).strip() if title else "Presentation"
            return self.create_presentation(t)

        if any(w in desc for w in ["email", "send mail", "compose"]):
            to_match = re.search(r"(\S+@\S+)", task_description)
            to = to_match.group(1) if to_match else ""
            body = re.search(r"(?:saying|body|content|message)\s+(.+)$", task_description, re.I)
            b = body.group(1).strip() if body else ""
            if to:
                return self.send_email(to, "Message from JARVIS", b)
            return WorkResult(False, "Who should I email?")

        if any(w in desc for w in ["research", "search", "find", "look up", "investigate"]):
            query = re.sub(r".*?(?:research|search|find|look up|investigate)\s+", "", task_description, flags=re.I).strip()
            if query:
                return self.research(query)
            return WorkResult(False, "What should I research?")

        if any(w in desc for w in ["blender", "3d model", "render"]):
            code = 'import bpy; bpy.ops.mesh.primitive_cube_add(); print("Cube created")'
            return self.run_blender(code)

        if any(w in desc for w in ["notepad", "text file", "notes"]):
            return self.open_app("notepad")

        # Unknown: try LLM planner as last resort
        try:
            from task_planner import get_planner
            planner = get_planner()
            planner._engine = self
            result_text = planner.do_anything(task_description)
            return WorkResult(True, result_text)
        except Exception:
            pass

        return self.open_app(task_description.split()[0] if task_description.split() else "cmd")


# Singleton

_engine = None

def get_engine() -> VirtualDesktopEngine:
    global _engine
    if _engine is None:
        _engine = VirtualDesktopEngine()
    return _engine
