"""JARVIS Specialized Agents — Domain-Specific Execution.

Each agent knows how to execute tasks in its domain using the
Execution Fabric. The Mission Engine coordinates them.

Agents don't "think" — they execute verified capabilities.
"""

import os, sys, json, logging, time, re
from pathlib import Path

log = logging.getLogger("specialized_agents")

sys.path.insert(0, os.path.dirname(__file__))


class AgentResult:
    """Result from a specialized agent execution."""
    def __init__(self, success: bool, output: str = "", error: str = "",
                 artifacts: list = None, data: dict = None):
        self.success = success
        self.output = output
        self.error = error
        self.artifacts = artifacts or []
        self.data = data or {}

    def to_dict(self) -> dict:
        return {
            "success": self.success,
            "output": self.output[:500],
            "error": self.error[:200],
            "artifacts": self.artifacts,
            "data": self.data,
        }


class BrowserAgent:
    """Executes browser-related tasks via CDP or mouse/keyboard."""

    def __init__(self):
        self.fabric = None

    def _get_fabric(self):
        if self.fabric is None:
            from execution_fabric import get_execution_fabric
            self.fabric = get_execution_fabric()
        return self.fabric

    def execute(self, action: str, params: dict) -> AgentResult:
        fabric = self._get_fabric()

        if action == "launch_browser":
            result = fabric.execute("launch_browser", params)
            return AgentResult(result.success, result.output, result.error)

        if action == "open_tab":
            url = params.get("url", "")
            result = fabric.execute("vdi_open_tab", {"url": url})
            return AgentResult(result.success, result.output, result.error)

        if action == "search_sites":
            sites = params.get("sites", [])
            query = params.get("query", "")
            opened = 0
            for site in sites[:3]:
                url = self._build_search_url(site, query)
                result = fabric.execute("vdi_open_tab", {"url": url})
                if result.success:
                    opened += 1
                time.sleep(1)
            return AgentResult(True, output=f"Opened {opened}/{len(sites)} sites")

        if action == "extract_clipboard":
            result = fabric.execute("extract_clipboard", params)
            if result.success:
                prices = self._extract_prices(result.output)
                return AgentResult(True, output=result.output[:200],
                                  data={"prices": prices, "text_length": len(result.output)})
            return AgentResult(False, error=result.error)

        if action == "scroll_page":
            direction = params.get("direction", "down")
            amount = params.get("amount", 5)
            result = fabric.execute("vdi_scroll", {"direction": direction, "amount": amount})
            return AgentResult(result.success, result.output, result.error)

        if action == "click":
            x = params.get("x", 0)
            y = params.get("y", 0)
            result = fabric.execute("vdi_click", {"x": x, "y": y})
            return AgentResult(result.success, result.output, result.error)

        if action == "type_text":
            text = params.get("text", "")
            result = fabric.execute("vdi_type", {"text": text})
            return AgentResult(result.success, result.output, result.error)

        if action == "press_key":
            key = params.get("key", "Return")
            result = fabric.execute("vdi_key", {"key": key})
            return AgentResult(result.success, result.output, result.error)

        if action == "list_tabs":
            result = fabric.execute("list_windows", params)
            if result.success:
                tabs = self._parse_tabs(result.output)
                return AgentResult(True, data={"tabs": tabs})
            return AgentResult(False, error=result.error)

        if action == "focus_tab":
            wid = params.get("wid", "")
            result = fabric.execute("vdi_focus_tab", {"wid": wid})
            return AgentResult(result.success, result.output, result.error)

        return AgentResult(False, error=f"Unknown browser action: {action}")

    def _build_search_url(self, site: str, query: str) -> str:
        import urllib.parse
        q = urllib.parse.quote(query)
        urls = {
            "google": f"https://www.google.com/search?q={q}",
            "skyscanner": f"https://www.skyscanner.net/transport/flights/lond/alas20dec/?adults=5",
            "booking.com": f"https://www.booking.com/searchresults.html?ss={q}",
            "kayak": f"https://www.kayak.com/flights",
        }
        return urls.get(site, f"https://www.google.com/search?q={q}")

    def _extract_prices(self, text: str) -> list:
        prices = []
        for m in re.finditer(r'[£$€₹¥](\d[\d,]*(?:\.\d{2})?)', text):
            try:
                val = float(m.group(1).replace(',', ''))
                if 10 < val < 500000:
                    prices.append({"amount": val, "raw": m.group(0)})
            except ValueError:
                pass
        return prices

    def _parse_tabs(self, output: str) -> list:
        tabs = []
        for line in output.strip().split('\n'):
            if any(w in line.lower() for w in ['chrome', 'google', 'firefox']):
                parts = line.split()
                if len(parts) >= 4:
                    tabs.append({"wid": parts[0], "title": ' '.join(parts[3:])})
        return tabs


class TravelAgent:
    """Executes travel search and price comparison tasks."""

    def __init__(self):
        self.browser = BrowserAgent()
        self.fabric = None

    def _get_fabric(self):
        if self.fabric is None:
            from execution_fabric import get_execution_fabric
            self.fabric = get_execution_fabric()
        return self.fabric

    def execute(self, action: str, params: dict) -> AgentResult:
        if action == "search_flights":
            return self._search_flights(params)
        elif action == "search_hotels":
            return self._search_hotels(params)
        elif action == "extract_prices":
            return self._extract_prices(params)
        elif action == "compare_currencies":
            return self._compare_currencies(params)
        return AgentResult(False, error=f"Unknown travel action: {action}")

    def _search_flights(self, params: dict) -> AgentResult:
        dest = params.get("destination", "alaska")
        group = params.get("group_size", 5)
        checkin = params.get("checkin", "2025-12-20")
        checkout = params.get("checkout", "2026-01-03")

        sites = [
            ("Google Flights", f"https://www.google.com/travel/flights?q=flights+to+{dest}+{checkin}+{group}+people"),
            ("Skyscanner", f"https://www.skyscanner.net/transport/flights/lond/{dest[:3]}20dec/?adults={group}"),
            ("Kayak", f"https://www.kayak.com/flights/LON-{dest.upper()}/{checkin}/{checkout}?sort=price_a"),
        ]

        all_prices = []
        for name, url in sites:
            # Open tab
            self.browser.execute("open_tab", {"url": url})
            time.sleep(3)

            # Scroll to load
            self.browser.execute("scroll_page", {"direction": "down", "amount": 8})
            time.sleep(1)

            # Extract prices
            result = self.browser.execute("extract_clipboard", {})
            if result.success:
                prices = result.data.get("prices", [])
                for p in prices:
                    p["source"] = name
                    p["destination"] = dest
                all_prices.extend(prices)

        return AgentResult(True, output=f"Found {len(all_prices)} prices",
                          data={"prices": all_prices})

    def _search_hotels(self, params: dict) -> AgentResult:
        dest = params.get("destination", "alaska")
        checkin = params.get("checkin", "2025-12-20")
        checkout = params.get("checkout", "2026-01-03")
        group = params.get("group_size", 5)

        url = f"https://www.booking.com/searchresults.html?ss={dest}&checkin={checkin}&checkout={checkout}&group_adults={group}"
        self.browser.execute("open_tab", {"url": url})
        time.sleep(4)
        self.browser.execute("scroll_page", {"direction": "down", "amount": 10})
        time.sleep(1)

        result = self.browser.execute("extract_clipboard", {})
        if result.success:
            prices = result.data.get("prices", [])
            for p in prices:
                p["source"] = "Booking.com"
                p["destination"] = dest
            return AgentResult(True, output=f"Found {len(prices)} hotel prices",
                              data={"prices": prices})
        return AgentResult(False, error=result.error)

    def _extract_prices(self, params: dict) -> AgentResult:
        result = self.browser.execute("extract_clipboard", {})
        if result.success:
            return AgentResult(True, data=result.data)
        return AgentResult(False, error=result.error)

    def _compare_currencies(self, params: dict) -> AgentResult:
        amounts = params.get("amounts", [])
        if not amounts:
            return AgentResult(False, error="No amounts to compare")

        fabric = self._get_fabric()
        from vdi_agent import CURRENCIES

        normalized = []
        for item in amounts:
            amount = item.get("amount", 0)
            currency = item.get("currency", "GBP")

            if currency == "GBP":
                gbp = amount
            else:
                result = fabric.execute("exchange_rate", {"from": currency, "to": "GBP"})
                rate = float(result.output) if result.success else 1.0
                gbp = amount * rate

            normalized.append({**item, "gbp_equivalent": round(gbp, 2)})

        normalized.sort(key=lambda x: x.get("gbp_equivalent", 999999))
        return AgentResult(True, data={"normalized": normalized, "cheapest": normalized[0] if normalized else None})


class OSAgent:
    """Executes OS-level tasks: file management, process control, system info."""

    def __init__(self):
        self.fabric = None

    def _get_fabric(self):
        if self.fabric is None:
            from execution_fabric import get_execution_fabric
            self.fabric = get_execution_fabric()
        return self.fabric

    def execute(self, action: str, params: dict) -> AgentResult:
        fabric = self._get_fabric()

        if action == "check_display":
            result = fabric.execute("check_display", params)
            return AgentResult(result.success, result.output, result.error)

        if action == "screenshot":
            result = fabric.execute("screenshot", params)
            return AgentResult(result.success, result.output, result.error,
                              artifacts=result.artifacts)

        if action == "kill_chrome":
            try:
                import subprocess
                subprocess.run(
                    ["sudo", "-u", "#1001", "bash", "-c",
                     "ps -eo pid,comm | grep -wi chrome | awk '{print $1}' | xargs -r kill -9 2>/dev/null"],
                    timeout=5, capture_output=True
                )
                return AgentResult(True, output="Chrome killed")
            except Exception as e:
                return AgentResult(False, error=str(e))

        if action == "list_processes":
            try:
                import subprocess
                result = subprocess.run(
                    ["ps", "-eo", "pid,comm,%cpu,%mem"],
                    capture_output=True, text=True, timeout=5
                )
                return AgentResult(True, output=result.stdout)
            except Exception as e:
                return AgentResult(False, error=str(e))

        if action == "system_stats":
            try:
                import subprocess
                result = subprocess.run(
                    ["free", "-m"],
                    capture_output=True, text=True, timeout=5
                )
                return AgentResult(True, output=result.stdout)
            except Exception as e:
                return AgentResult(False, error=str(e))

        return AgentResult(False, error=f"Unknown OS action: {action}")


class CodingAgent:
    """Executes code-related tasks: write, test, lint, run."""

    def __init__(self):
        self.fabric = None

    def _get_fabric(self):
        if self.fabric is None:
            from execution_fabric import get_execution_fabric
            self.fabric = get_execution_fabric()
        return self.fabric

    def execute(self, action: str, params: dict) -> AgentResult:
        fabric = self._get_fabric()

        if action == "write_file":
            filepath = params.get("path", "")
            content = params.get("content", "")
            try:
                from pathlib import Path
                Path(filepath).parent.mkdir(parents=True, exist_ok=True)
                Path(filepath).write_text(content)
                return AgentResult(True, output=f"Written {len(content)} bytes to {filepath}",
                                  artifacts=[filepath])
            except Exception as e:
                return AgentResult(False, error=str(e))

        if action == "run_python":
            code = params.get("code", "")
            result = fabric.execute("run_script", {"code": code, "timeout": 30})
            return AgentResult(result.success, result.output[:500], result.error)

        if action == "run_bash":
            command = params.get("command", "")
            result = fabric.execute("run_command", {"command": command, "timeout": 30})
            return AgentResult(result.success, result.output[:500], result.error)

        if action == "run_tests":
            test_path = params.get("path", ".")
            result = fabric.execute("run_command", {
                "command": f"cd {test_path} && python3 -m pytest --tb=short -q 2>&1 | head -50",
                "timeout": 60,
            })
            passed = "passed" in result.output.lower() or result.returncode == 0
            return AgentResult(passed, output=result.output[:500], error=result.error)

        if action == "lint":
            filepath = params.get("path", "")
            result = fabric.execute("run_command", {
                "command": f"python3 -m py_compile {filepath} 2>&1",
                "timeout": 10,
            })
            return AgentResult(result.success, output=result.output, error=result.error)

        if action == "install_package":
            package = params.get("package", "")
            result = fabric.execute("run_command", {
                "command": f"pip3 install --break-system-packages {package} 2>&1 | tail -5",
                "timeout": 60,
            })
            return AgentResult(result.success, output=result.output[:300], error=result.error)

        return AgentResult(False, error=f"Unknown coding action: {action}")


class DocumentAgent:
    """Creates and manipulates documents: PPTX, DOCX, XLSX, CSV, PDF."""

    def __init__(self):
        self.fabric = None

    def _get_fabric(self):
        if self.fabric is None:
            from execution_fabric import get_execution_fabric
            self.fabric = get_execution_fabric()
        return self.fabric

    def execute(self, action: str, params: dict) -> AgentResult:
        fabric = self._get_fabric()

        if action == "create_pptx":
            title = params.get("title", "Presentation")
            content = params.get("content", [])
            output = params.get("output", f"/home/workuser/Desktop/{title}.pptx")
            try:
                from pptx import Presentation
                from pptx.util import Inches
                prs = Presentation()
                for section in content:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = section.get("title", "")
                    if len(slide.placeholders) > 1:
                        slide.placeholders[1].text = section.get("content", "")
                prs.save(output)
                return AgentResult(True, output=f"Created {output}", artifacts=[output])
            except Exception as e:
                return AgentResult(False, error=str(e))

        if action == "create_docx":
            title = params.get("title", "Document")
            text = params.get("text", "")
            output = params.get("output", f"/home/workuser/Desktop/{title}.docx")
            try:
                from docx import Document
                doc = Document()
                doc.add_heading(title, 0)
                for para in text.split("\n"):
                    if para.strip():
                        doc.add_paragraph(para)
                doc.save(output)
                return AgentResult(True, output=f"Created {output}", artifacts=[output])
            except Exception as e:
                return AgentResult(False, error=str(e))

        if action == "create_xlsx":
            title = params.get("title", "Spreadsheet")
            headers = params.get("headers", [])
            rows = params.get("rows", [])
            output = params.get("output", f"/home/workuser/Desktop/{title}.xlsx")
            try:
                from openpyxl import Workbook
                wb = Workbook()
                ws = wb.active
                ws.title = title[:31]
                if headers:
                    ws.append(headers)
                for row in rows:
                    ws.append(row)
                wb.save(output)
                return AgentResult(True, output=f"Created {output}", artifacts=[output])
            except Exception as e:
                return AgentResult(False, error=str(e))

        if action == "create_csv":
            headers = params.get("headers", [])
            rows = params.get("rows", [])
            output = params.get("output", "/home/workuser/Desktop/data.csv")
            try:
                import csv
                with open(output, 'w', newline='') as f:
                    writer = csv.writer(f)
                    if headers:
                        writer.writerow(headers)
                    writer.writerows(rows)
                return AgentResult(True, output=f"Created {output}", artifacts=[output])
            except Exception as e:
                return AgentResult(False, error=str(e))

        return AgentResult(False, error=f"Unknown document action: {action}")


class DataAgent:
    """Data processing: fetch, parse, transform, analyze."""

    def __init__(self):
        self.fabric = None

    def _get_fabric(self):
        if self.fabric is None:
            from execution_fabric import get_execution_fabric
            self.fabric = get_execution_fabric()
        return self.fabric

    def execute(self, action: str, params: dict) -> AgentResult:
        fabric = self._get_fabric()

        if action == "fetch_url":
            url = params.get("url", "")
            try:
                import requests
                r = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
                return AgentResult(True, output=r.text[:5000],
                                  data={"status": r.status_code, "length": len(r.text)})
            except Exception as e:
                return AgentResult(False, error=str(e))

        if action == "parse_json":
            text = params.get("text", "")
            try:
                import json
                data = json.loads(text)
                return AgentResult(True, output=json.dumps(data, indent=2)[:3000],
                                  data={"parsed": True})
            except Exception as e:
                return AgentResult(False, error=str(e))

        if action == "extract_table":
            html = params.get("html", "")
            try:
                import re
                rows = re.findall(r'<tr[^>]*>(.*?)</tr>', html, re.DOTALL)
                tables = []
                for row in rows:
                    cells = re.findall(r'<t[dh][^>]*>(.*?)</t[dh]>', row, re.DOTALL)
                    if cells:
                        tables.append([re.sub(r'<[^>]+>', '', c).strip() for c in cells])
                return AgentResult(True, output=json.dumps(tables[:20]),
                                  data={"rows": len(tables)})
            except Exception as e:
                return AgentResult(False, error=str(e))

        if action == "analyze_data":
            data = params.get("data", [])
            try:
                if not data:
                    return AgentResult(False, error="No data to analyze")
                # Basic stats
                if all(isinstance(d, (int, float)) for d in data):
                    import statistics
                    stats = {
                        "count": len(data),
                        "mean": statistics.mean(data),
                        "median": statistics.median(data),
                        "min": min(data),
                        "max": max(data),
                    }
                    return AgentResult(True, output=json.dumps(stats), data=stats)
                return AgentResult(True, output=f"Got {len(data)} items")
            except Exception as e:
                return AgentResult(False, error=str(e))

        return AgentResult(False, error=f"Unknown data action: {action}")


# ── Agent Registry ──
_agents = {
    "browser_agent": BrowserAgent(),
    "travel_agent": TravelAgent(),
    "os_agent": OSAgent(),
    "coding_agent": CodingAgent(),
    "document_agent": DocumentAgent(),
    "data_agent": DataAgent(),
    "planner": None,  # Uses LLM
    "research_agent": None,  # Uses LLM + browser
}

def get_agent(name: str):
    """Get a specialized agent by name."""
    return _agents.get(name)
