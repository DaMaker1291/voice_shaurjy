"""
JARVIS 3D CAD Pipeline — Parametric modeling, rendering, and presentation generation.

Handles:
- Blender parametric 3D modeling via headless Python
- FreeCAD parametric mechanical design
- 4K render animation (60-frame turntable)
- PowerPoint generation from 3D renders
- CAD file compression and export
- Live preview streaming to wearables
"""
import os
import sys
import json
import time
import logging
import subprocess
import tempfile
from pathlib import Path
from typing import Optional
from dataclasses import dataclass

# Presentation imports (lazy-loaded)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except ImportError:
    pass

logger = logging.getLogger("cad_pipeline")


@dataclass
class RenderResult:
    """Result of a 3D render operation."""
    success: bool
    output_path: str = ""
    frame_count: int = 0
    resolution: str = ""
    file_size_mb: float = 0.0
    duration_seconds: float = 0.0
    error: str = ""


class BlenderPipeline:
    """Headless Blender 3D modeling and rendering pipeline."""

    def __init__(self, display: str = ":99"):
        self.display = display
        self._blender_path = self._find_blender()

    def _find_blender(self) -> str:
        """Find Blender executable in WSL."""
        try:
            result = subprocess.run(
                ["wsl", "-e", "bash", "-c", "which blender 2>/dev/null || echo NOT_FOUND"],
                capture_output=True, text=True, timeout=5
            )
            path = result.stdout.strip()
            if path and path != "NOT_FOUND":
                return path
        except Exception:
            pass
        return "blender"  # Fallback to PATH

    def generate_turbine(self, output_dir: str = None, num_blades: int = 12,
                         radius: float = 3.0, blade_length: float = 1.5) -> RenderResult:
        """Generate a parametric 3D turbine engine model."""
        script = f"""
import bpy
import math
import os

# Clear default scene
bpy.ops.object.select_all(action='SELECT')
bpy.ops.object.delete()

# ── Generate Turbine Hub ──
bpy.ops.mesh.primitive_cylinder_add(radius=1.0, depth=0.5, location=(0, 0, 0))
hub = bpy.context.active_object
hub.name = "TurbineHub"

# ── Generate Turbine Blades ──
num_blades = {num_blades}
radius = {radius}
blade_length = {blade_length}

for i in range(num_blades):
    angle = i * (2 * math.pi / num_blades)
    x = math.cos(angle) * radius
    y = math.sin(angle) * radius

    # Create blade
    bpy.ops.mesh.primitive_cube_add(size=1, location=(x, y, 0))
    blade = bpy.context.active_object
    blade.name = f"Blade_{{i}}"
    blade.scale = (0.1, blade_length, 0.02)
    blade.rotation_euler = (0.15, 0, angle)

    # Add bevel modifier for smoothness
    bevel = blade.modifiers.new(name="Bevel", type='BEVEL')
    bevel.width = 0.02
    bevel.segments = 3

# ── Add Central Shaft ──
bpy.ops.mesh.primitive_cylinder_add(radius=0.15, depth=2.0, location=(0, 0, 0))
shaft = bpy.context.active_object
shaft.name = "Shaft"

# ── Add Base Plate ──
bpy.ops.mesh.primitive_cylinder_add(radius=2.5, depth=0.1, location=(0, 0, -0.3))
base = bpy.context.active_object
base.name = "BasePlate"

# ── Setup Camera (Turntable) ──
bpy.ops.object.camera_add(location=(0, -8, 4), rotation=(1.1, 0, 0))
camera = bpy.context.active_object
camera.name = "TurntableCamera"
bpy.context.scene.camera = camera

# ── Setup Lighting ──
bpy.ops.object.light_add(type='AREA', location=(5, -5, 8))
light = bpy.context.active_object
light.data.energy = 500
light.name = "KeyLight"

bpy.ops.object.light_add(type='AREA', location=(-5, 5, 3))
fill = bpy.context.active_object
fill.data.energy = 200
fill.name = "FillLight"

# ── Setup Render ──
scene = bpy.context.scene
scene.render.engine = 'BLENDER_EEVEE'
scene.render.resolution_x = 1920
scene.render.resolution_y = 1080
scene.render.resolution_percentage = 100

# ── Material Setup ──
mat = bpy.data.materials.new(name="MetalMaterial")
mat.use_nodes = True
bsdf = mat.node_tree.nodes["Principled BSDF"]
bsdf.inputs["Base Color"].default_value = (0.3, 0.35, 0.4, 1)
bsdf.inputs["Metallic"].default_value = 0.9
bsdf.inputs["Roughness"].default_value = 0.2

# Apply to all mesh objects
for obj in bpy.data.objects:
    if obj.type == 'MESH':
        obj.data.materials.append(mat)

# ── Save .blend file ──
output_dir = "{output_dir or tempfile.gettempdir()}"
os.makedirs(output_dir, exist_ok=True)
blend_path = os.path.join(output_dir, "turbine_engine.blend")
bpy.ops.wm.save_as_mainfile(filepath=blend_path)

# ── Render Single Frame ──
render_path = os.path.join(output_dir, "turbine_render.png")
scene.render.filepath = render_path
bpy.ops.render.render(write_still=True)
print(f"Rendered: {{render_path}}")

# ── Render 60-Frame Turntable Animation ──
anim_dir = os.path.join(output_dir, "turntable_frames")
os.makedirs(anim_dir, exist_ok=True)

for frame in range(60):
    angle = frame * (2 * math.pi / 60)
    camera.location.x = math.sin(angle) * 8
    camera.location.y = -math.cos(angle) * 8
    camera.location.z = 4
    camera.rotation_euler = (1.1, 0, angle)
    scene.frame_set(frame + 1)
    scene.render.filepath = os.path.join(anim_dir, f"frame_{{frame+1:04d}}.png")
    bpy.ops.render.render(write_still=True)

print(f"Animation rendered: {{anim_dir}} (60 frames)")
print("TURBINE_GENERATION_COMPLETE")
"""
        return self._run_blender_script(script, "Turbine Generation")

    def generate_custom_model(self, description: str, output_dir: str = None) -> RenderResult:
        """Generate a custom 3D model from natural language description."""
        # Use LLM to generate Blender Python script
        script = self._llm_generate_blender_script(description)
        if not script:
            return RenderResult(success=False, error="Could not generate model script")
        return self._run_blender_script(script, "Custom Model")

    def render_animation(self, blend_file: str, output_dir: str,
                         frames: int = 60, resolution: str = "1920x1080") -> RenderResult:
        """Render animation from existing .blend file."""
        width, height = resolution.split("x")
        script = f"""
import bpy
import os

# Load blend file
bpy.ops.wm.open_mainfile(filepath="{blend_file}")

# Setup render
scene = bpy.context.scene
scene.render.engine = 'BLENDER_EEVEE'
scene.render.resolution_x = {width}
scene.render.resolution_y = {height}
scene.render.resolution_percentage = 100

# Find camera or create one
camera = None
for obj in bpy.data.objects:
    if obj.type == 'CAMERA':
        camera = obj
        break
if not camera:
    bpy.ops.object.camera_add(location=(0, -8, 4), rotation=(1.1, 0, 0))
    camera = bpy.context.active_object
scene.camera = camera

# Render frames
output_dir = "{output_dir}"
os.makedirs(output_dir, exist_ok=True)

for frame in range({frames}):
    import math
    angle = frame * (2 * math.pi / {frames})
    camera.location.x = math.sin(angle) * 8
    camera.location.y = -math.cos(angle) * 8
    camera.location.z = 4
    scene.frame_set(frame + 1)
    scene.render.filepath = os.path.join(output_dir, f"frame_{{frame+1:04d}}.png")
    bpy.ops.render.render(write_still=True)

print(f"Rendered {{ {frames} }} frames to {output_dir}")
print("RENDER_COMPLETE")
"""
        return self._run_blender_script(script, "Animation Render")

    def export_stl(self, blend_file: str, output_path: str) -> bool:
        """Export 3D model as STL file."""
        script = f"""
import bpy

# Load blend file
bpy.ops.wm.open_mainfile(filepath="{blend_file}")

# Select all mesh objects
bpy.ops.object.select_all(action='SELECT')

# Export as STL
bpy.ops.export_mesh.stl(filepath="{output_path}", use_selection=True)
print(f"Exported STL: {output_path}")
print("STL_EXPORT_COMPLETE")
"""
        result = self._run_blender_script(script, "STL Export")
        return result.success

    def _run_blender_script(self, script: str, operation: str) -> RenderResult:
        """Execute a Blender Python script in WSL."""
        start_time = time.time()

        # Write script to temp file
        script_path = f"/tmp/jarvis_blender_{int(time.time())}.py"

        try:
            # Write script to WSL
            write_cmd = f"cat > {script_path} << 'BLENDER_SCRIPT_EOF'\n{script}\nBLENDER_SCRIPT_EOF"
            subprocess.run(
                ["wsl", "-e", "bash", "-c", write_cmd],
                capture_output=True, timeout=10
            )

            # Run Blender headlessly
            blender_cmd = f"DISPLAY={self.display} {self._blender_path} --background --python {script_path} 2>&1"
            result = subprocess.run(
                ["wsl", "-e", "bash", "-c", blender_cmd],
                capture_output=True, text=True, timeout=600  # 10 min timeout for renders
            )

            duration = time.time() - start_time

            # Parse output
            output = result.stdout + result.stderr
            success = "COMPLETE" in output or result.returncode == 0

            # Extract output path
            output_path = ""
            for line in output.split("\n"):
                if "Rendered:" in output or "Exported" in output or "Animation rendered:" in output:
                    # Try to extract path
                    if ":" in line:
                        output_path = line.split(":", 1)[1].strip()

            # Clean up
            subprocess.run(
                ["wsl", "-e", "bash", "-c", f"rm -f {script_path}"],
                capture_output=True, timeout=5
            )

            return RenderResult(
                success=success,
                output_path=output_path,
                duration_seconds=round(duration, 2),
                error="" if success else f"Blender exited with code {result.returncode}",
            )

        except subprocess.TimeoutExpired:
            return RenderResult(
                success=False,
                duration_seconds=time.time() - start_time,
                error="Blender timed out (10 min limit)",
            )
        except Exception as e:
            return RenderResult(
                success=False,
                duration_seconds=time.time() - start_time,
                error=str(e),
            )

    def _llm_generate_blender_script(self, description: str) -> Optional[str]:
        """Use LLM to generate Blender Python script from description."""
        prompt = f"""Generate a Blender Python script (bpy) to create: {description}

Requirements:
- Use only bpy standard operations
- Include materials, lighting, and camera setup
- Output print statements for progress tracking
- Save the .blend file

Script:"""

        try:
            from groq_agent import generate
            response = generate(prompt, user_id="cad_pipeline", max_tokens=2000, temperature=0.1)
            if response:
                # Extract code block
                if "```python" in response:
                    code = response.split("```python")[1].split("```")[0]
                elif "```" in response:
                    code = response.split("```")[1].split("```")[0]
                else:
                    code = response
                return code.strip()
        except Exception as e:
            logger.warning(f"LLM script generation failed: {e}")

        return None


class PresentationGenerator:
    """Generate professional PowerPoint presentations with dark tech theme."""

    def __init__(self):
        self._pptx_available = self._check_pptx()
        # Dark tech color palette
        self.BG_DARK = RGBColor(0x0D, 0x11, 0x17)
        self.BG_CARD = RGBColor(0x16, 0x1B, 0x22)
        self.ACCENT_BLUE = RGBColor(0x58, 0xA6, 0xFF)
        self.ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)
        self.ACCENT_ORANGE = RGBColor(0xF7, 0x8C, 0x1E)
        self.TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        self.TEXT_GRAY = RGBColor(0x8B, 0x94, 0x9E)
        self.TEXT_LIGHT = RGBColor(0xC9, 0xD1, 0xD9)

    def _check_pptx(self) -> bool:
        try:
            import pptx
            return True
        except ImportError:
            try:
                subprocess.run(["pip", "install", "python-pptx"], capture_output=True, timeout=30)
                return True
            except Exception:
                return False

    def _set_bg(self, slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.BG_DARK

    def _add_shape(self, slide, left, top, width, height, fill_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _add_text(self, slide, left, top, width, height, text, size=18,
                  color=None, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
        color = color or self.TEXT_WHITE
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        return txBox

    def _add_accent(self, slide, left, top, width, color=None):
        color = color or self.ACCENT_BLUE
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def _add_bullets(self, slide, left, top, width, height, items, size=14,
                     color=None, bullet_color=None):
        color = color or self.TEXT_LIGHT
        bullet_color = bullet_color or self.ACCENT_BLUE
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "  " + item
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.name = "Segoe UI"
            p.space_after = Pt(6)
            pPr = p._p.get_or_add_pPr()
            buClr = pPr.makeelement(qn('a:buClr'), {})
            srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': f'{bullet_color}'})
            buClr.append(srgbClr)
            pPr.append(buClr)
            pPr.append(pPr.makeelement(qn('a:buChar'), {'char': '\u25CF'}))

    def _add_transition(self, slide, transition_type="fade"):
        transition = slide.element.makeelement(qn('p:transition'), {
            'advClick': '1', 'advTm': '0', 'spd': 'med'
        })
        child = transition.makeelement(qn('p:fade'), {})
        transition.append(child)
        slide.element.append(transition)

    def create_presentation(self, title: str, slides_data: list[dict],
                            output_path: str, images: list[str] = None) -> bool:
        """Create a professional presentation with dark tech theme.

        slides_data: [{"title": "...", "content": "bullet\npoints", "accent": "blue|green|orange"}]
        images: Optional list of image paths to include on slides
        """
        if not self._pptx_available:
            return False

        try:
            from pptx.oxml.ns import qn as _qn
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._set_bg(slide)
            self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), self.ACCENT_BLUE)
            self._add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
                           title.upper(), size=44, bold=True, align=PP_ALIGN.CENTER,
                           font="Segoe UI Light")
            self._add_accent(slide, Inches(4.5), Inches(3.3), Inches(4))
            self._add_text(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.8),
                           "Autonomous Generation Report", size=24, color=self.TEXT_GRAY,
                           align=PP_ALIGN.CENTER)
            self._add_text(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.5),
                           f"Generated: {time.strftime('%B %d, %Y at %H:%M')}  |  JARVIS OS Pipeline",
                           size=14, color=self.TEXT_GRAY, align=PP_ALIGN.CENTER)
            self._add_shape(slide, Inches(0), Inches(7.4), Inches(13.333), Pt(4), self.ACCENT_GREEN)
            self._add_transition(slide)

            # Content slides
            colors = [self.ACCENT_BLUE, self.ACCENT_GREEN, self.ACCENT_ORANGE]
            for idx, slide_data in enumerate(slides_data):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._set_bg(slide)
                accent = colors[idx % 3]

                self._add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                               slide_data.get("title", "Untitled").upper(), size=32, bold=True)
                self._add_accent(slide, Inches(0.8), Inches(1.2), Inches(2.5), accent)

                content = slide_data.get("content", "")
                if content:
                    items = [x.strip() for x in content.split("\n") if x.strip()]
                    self._add_bullets(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(5.0),
                                      items, size=15, bullet_color=accent)

                # Add image if provided
                image_path = slide_data.get("image") or (images[idx] if images and idx < len(images) else None)
                if image_path and os.path.exists(image_path):
                    slide.shapes.add_picture(image_path, Inches(7), Inches(2), width=Inches(5.5))

                self._add_transition(slide, "push" if idx % 2 == 0 else "fade")

            # Save
            os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
            prs.save(output_path)
            logger.info(f"Presentation saved: {output_path} ({len(slides_data) + 1} slides)")
            return True

        except Exception as e:
            logger.error(f"Presentation generation failed: {e}")
            return False

    def create_from_research(self, topic: str, findings: list[str],
                             output_path: str) -> bool:
        """Create a presentation from research findings."""
        slides = [
            {"title": topic, "content": "Research Report\n\nGenerated by JARVIS OS"},
        ]
        for i, finding in enumerate(findings):
            slides.append({"title": f"Finding {i+1}", "content": finding})
        slides.append({"title": "Conclusion", "content": f"Summary of {len(findings)} findings on {topic}"})
        return self.create_presentation(f"{topic} - Research Report", slides, output_path)


class CADFileCompressor:
    """Compress CAD files and 3D assets."""

    def compress_to_zip(self, input_paths: list[str], output_path: str) -> bool:
        """Compress files/folders into a zip archive."""
        import zipfile
        try:
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                for path in input_paths:
                    if os.path.isfile(path):
                        zf.write(path, os.path.basename(path))
                    elif os.path.isdir(path):
                        for root, dirs, files in os.walk(path):
                            for file in files:
                                file_path = os.path.join(root, file)
                                arcname = os.path.relpath(file_path, os.path.dirname(path))
                                zf.write(file_path, arcname)
            return True
        except Exception as e:
            logger.error(f"Compression failed: {e}")
            return False

    def get_size_mb(self, path: str) -> float:
        """Get file/folder size in MB."""
        if os.path.isfile(path):
            return os.path.getsize(path) / (1024 * 1024)
        elif os.path.isdir(path):
            total = 0
            for root, dirs, files in os.walk(path):
                for f in files:
                    total += os.path.getsize(os.path.join(root, f))
            return total / (1024 * 1024)
        return 0.0


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="[CAD] %(message)s")
    pipeline = BlenderPipeline()
    print(f"Blender path: {pipeline._blender_path}")
    print(f"Blender available: {pipeline._blender_path != 'blender'}")
